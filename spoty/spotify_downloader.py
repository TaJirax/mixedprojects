#!/usr/bin/env python3
"""
Spotify Downloader - Modern Edition - Performance Optimized
- V2RayN/SOCKS5 proxy support
- Built-in local HTTP bridge for SOCKS5 proxies
- Optimized download speed matching CLI performance
"""

import os
import re
import sys
import json
import time
import queue
import contextlib
import hashlib
import shutil
import threading
import subprocess
import webbrowser
import zipfile
import ssl
import urllib.request
import urllib.error
import socket
import socketserver
import select
import tarfile
import tempfile
import ipaddress
import io
from html.parser import HTMLParser
from urllib.parse import quote, unquote, urljoin, urlsplit
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor

from blueknight_paths import (
    ARCH, COOKIE_FAILURE, EXE, IS_ANDROID, IS_LINUX, IS_MACOS, IS_WINDOWS, OS_TAG,
    SITE_SESSION_COOKIES, SOURCES, app_data_dir, candidate_label, cookie_candidates,
    cookie_dir, download_dir, download_root, host_matches, jar_line_matches_domain,
    locked_by_browser, read_registry, record_jar, registry_entry, ytdlp_cookie_args)
import pyshell

# The window toolkit. Android supplies its own WebView and never imports this.
if not IS_ANDROID:
    import webview
else:
    webview = None


APP_NAME = "Blue Knight Downloader"
APP_VERSION = "7.0.2"
CREATOR = "Blue Knight"
TELEGRAM_URL = "https://t.me/BlueKnight_Net"

# Every tool publishes one release asset per platform under a name of its own
# choosing. The names are collected here, keyed the same way, so that adding a
# platform is a table entry rather than a branch in the download code.
SPOTDL_VERSION = "4.5.2"
SPOTDL_ASSETS = {"win": "spotdl-{v}-win32.exe", "mac": "spotdl-{v}-darwin",
                 "linux": "spotdl-{v}-linux", "android": None}
YTDLP_VERSION = "2026.07.04"
YTDLP_ASSETS = {"win": "yt-dlp.exe", "mac": "yt-dlp_macos",
                "linux": "yt-dlp_linux_aarch64" if ARCH == "arm64" else "yt-dlp_linux",
                "android": None}
DENO_ASSETS = {
    "win": "deno-x86_64-pc-windows-msvc.zip",
    "mac": f"deno-{'aarch64' if ARCH == 'arm64' else 'x86_64'}-apple-darwin.zip",
    "linux": f"deno-{'aarch64' if ARCH == 'arm64' else 'x86_64'}-unknown-linux-gnu.zip",
    # Deno publishes no Android build. The Android shell answers yt-dlp's
    # JavaScript challenges with the system WebView instead; see js_runtime().
    "android": None,
}


def spotdl_urls(version=SPOTDL_VERSION):
    asset = SPOTDL_ASSETS[OS_TAG]
    if not asset:
        return ()
    asset = asset.format(v=version)
    return (f"https://github.com/spotDL/spotify-downloader/releases/download/"
            f"v{version}/{asset}",
            f"https://sourceforge.net/projects/spotdl.mirror/files/v{version}/"
            f"{asset}/download")


def ytdlp_urls(version=YTDLP_VERSION):
    asset = YTDLP_ASSETS[OS_TAG]
    if not asset:
        return ()
    # Independent hosts, tried in order. SourceForge is the escape hatch for
    # networks that block GitHub release downloads.
    return (f"https://github.com/yt-dlp/yt-dlp/releases/latest/download/{asset}",
            f"https://github.com/yt-dlp/yt-dlp/releases/download/{version}/{asset}",
            f"https://sourceforge.net/projects/yt-dlp.mirror/files/{version}/"
            f"{asset}/download")


def deno_urls():
    asset = DENO_ASSETS[OS_TAG]
    if not asset:
        return ()
    return (f"https://github.com/denoland/deno/releases/latest/download/{asset}",
            f"https://dl.deno.land/release-latest/{asset}")


SPOTDL_DOWNLOAD_URLS = spotdl_urls()
YTDLP_DOWNLOAD_URLS = ytdlp_urls()
DENO_DOWNLOAD_URLS = deno_urls()

# FFmpeg is the one tool with no single project publishing every platform.
# Each entry is (binaries this archive should contain, mirrors to try) — macOS
# needs two archives because evermeet.cx packages ffmpeg and ffprobe apart.
FFMPEG_DOWNLOADS = {
    "win": ((("ffmpeg", "ffprobe"), (
        "https://github.com/BtbN/FFmpeg-Builds/releases/download/latest/ffmpeg-master-latest-win64-gpl.zip",
        "https://www.gyan.dev/ffmpeg/builds/ffmpeg-release-essentials.zip")),),
    "linux": ((("ffmpeg", "ffprobe"), (
        f"https://github.com/BtbN/FFmpeg-Builds/releases/download/latest/"
        f"ffmpeg-master-latest-linux{'arm64' if ARCH == 'arm64' else '64'}-gpl.tar.xz",)),),
    "mac": ((("ffmpeg",), ("https://evermeet.cx/ffmpeg/getrelease/ffmpeg/zip",)),
            (("ffprobe",), ("https://evermeet.cx/ffmpeg/getrelease/ffprobe/zip",))),
    # Android ships FFmpeg inside the APK as a native library; there is nowhere
    # on the device it is allowed to download and mark executable at runtime.
    "android": (),
}

# Every downloadable tool is described once, so finding, downloading and
# updating are one code path instead of one per tool.
TOOLS = {
    "spotdl": {
        "label": "spotDL", "exe": f"spotdl{EXE}", "urls": SPOTDL_DOWNLOAD_URLS,
        "api": "https://api.github.com/repos/spotDL/spotify-downloader/releases/latest",
    },
    "yt-dlp": {
        "label": "yt-dlp", "exe": f"yt-dlp{EXE}", "urls": YTDLP_DOWNLOAD_URLS,
        "api": "https://api.github.com/repos/yt-dlp/yt-dlp/releases/latest",
        # Android has no release binary and no second Python to spawn, so the
        # importable package stands in and pyshell runs it in this process.
        "module": "yt_dlp",
    },
    "deno": {
        "label": "Deno", "exe": f"deno{EXE}", "urls": DENO_DOWNLOAD_URLS,
        "api": "https://api.github.com/repos/denoland/deno/releases/latest",
        "archive_member": f"deno{EXE}",
        "checksum_urls": tuple(url + ".sha256sum" for url in DENO_DOWNLOAD_URLS),
    },
}
# A platform with neither a release asset nor an importable package has no tool
# to find, download or update, and every check for one would only ever fail.
TOOLS = {name: spec for name, spec in TOOLS.items()
         if spec["urls"] or spec.get("module")}

# yt-dlp format specs shared by the YouTube and TikTok pages.
MEDIA_FORMATS = {
    "best": "bv*+ba/b",
    "1080": "bestvideo[height<=1080]+bestaudio/best[height<=1080]",
    "720": "bestvideo[height<=720]+bestaudio/best[height<=720]",
    "480": "bestvideo[height<=480]+bestaudio/best[height<=480]",
}
MEDIA_PATTERNS = {
    "youtube": re.compile(r"https?://(?:[\w-]+\.)*(?:youtube\.com|youtu\.be)/", re.I),
    "ytmusic": re.compile(r"https?://music\.youtube\.com/", re.I),
    "tiktok": re.compile(r"https?://(?:[\w-]+\.)*tiktok\.com/", re.I),
    "instagram": re.compile(r"https?://(?:[\w-]+\.)*instagram\.com/", re.I),
    "soundcloud": re.compile(r"https?://(?:[\w-]+\.)*soundcloud\.com/", re.I),
    "x": re.compile(r"https?://(?:[\w-]+\.)*(?:x\.com|twitter\.com)/", re.I),
}
MEDIA_LABELS = {
    "youtube": "YouTube", "ytmusic": "YouTube Music", "tiktok": "TikTok",
    "instagram": "Instagram", "soundcloud": "SoundCloud", "x": "X",
    "general": "Video",
}
YOUTUBE_DOMAINS = ("youtube.com", "youtu.be", "youtube-nocookie.com")
# Sources served by YouTube's extractor, and so subject to its bot checks and
# its PO-token client dance. YouTube Music is youtube.com wearing a different hat.
YOUTUBE_KINDS = {"youtube", "ytmusic"}
# Sources that will not serve anything useful to a logged-out client.
COOKIE_SOURCES = {"instagram"}
# Sources that usually work logged out, but sometimes demand a session. They get
# cookies only after asking for them, so the normal path stays cookie-free.
COOKIE_ON_DEMAND = {"youtube", "ytmusic", "tiktok", "soundcloud", "x", "general"}
# YouTube is rolling out "PO tokens", and a client that needs one cannot serve a
# download without it. Per yt-dlp's PO Token Guide the exceptions are android_vr,
# web_embedded and tv — so those are the ones worth retrying as. tv_simply,
# web_safari and mweb all need a token, which is why they answered the earlier
# attempts with "Requested format is not available" rather than a video.
# yt-dlp's own client table says android_vr, ios and android need neither a
# proof-of-origin token nor a solved JavaScript player, and that none of the
# three can carry cookies.
#
# The first two rungs are exactly what they were, so nothing that used to
# succeed now takes longer to get there: a signed-in run still reaches the
# cookie-carrying rung on its first escalation, rather than after two clients
# that would ignore the session it is trying to spend. ios and android are
# added underneath, where they cost nothing until everything above has failed.
YT_CLIENT_LADDER = ("android_vr", "web_embedded,tv,default", "ios", "android")
# The rungs that can carry a signed-in session, by yt-dlp's SUPPORTS_COOKIES.
YT_COOKIE_CLIENTS = frozenset({"web_embedded,tv,default"})
YOUTUBE_MEDIA_DENIED = re.compile(
    r"unable to download (?:video|audio) data.*(?:http error 403|forbidden)|"
    r"http error 403:\s*forbidden|fragment.*http error 403", re.I | re.S)
PO_TOKEN_PATTERN = re.compile(r"^mweb\.gvs\+[A-Za-z0-9._~=/+-]{20,4096}$")
SIGNIN_DEMANDED = re.compile(
    r"not a bot|sign ?in to confirm|login required|account.*cookies|"
    r"this video is only available|use --cookies", re.I)
GENERAL_CHALLENGE = re.compile(
    r"captcha|cloudflare|cf-chl|attention required|verify (?:that )?you are human|"
    r"not a bot|login required|sign ?in|use --cookies|http error 40[13]|forbidden", re.I)
# yt-dlp loads the cookie jar before it touches the network, so a bogus URL is
# enough to find out whether a jar is readable — no request, no rate limit.
COOKIE_PROBE_URL = "blueknightprobe://cookies"
# The domain whose session makes a jar worth using, per source.
COOKIE_DOMAINS = {
    "instagram": "instagram.com", "youtube": "youtube.com",
    "ytmusic": "youtube.com", "tiktok": "tiktok.com",
    "soundcloud": "soundcloud.com", "x": "x.com",
}

# Performance optimization: Batch UI updates
UI_UPDATE_INTERVAL = 0.1
LOG_BATCH_SIZE = 10


def asset_path(name):
    """Return a bundled/source asset path, or None if missing."""
    roots = [Path(getattr(sys, "_MEIPASS", "")) / "assets",
             Path(__file__).resolve().parent.parent / "assets",
             Path(__file__).resolve().parent / "assets"]
    for root in roots:
        candidate = root / name
        if candidate.is_file():
            return str(candidate)
    return None


def web_index():
    """Path to the glass UI document."""
    for root in (Path(getattr(sys, "_MEIPASS", "")), Path(__file__).resolve().parent):
        candidate = root / "web" / "index.html"
        if candidate.is_file():
            return str(candidate)
    raise FileNotFoundError("web/index.html is missing")


# Where a sign-in window starts, and which pages to read the session from.
#
# The browser-extension guides say to land on robots.txt, and for an extension
# that is right. Here it is not: WebView2 answers get_cookies() for the document
# it currently holds, and a text file yields an empty list — measured, 0 cookies
# on robots.txt against 5 on the homepage. Real pages only, and more than one,
# because a YouTube session is split across youtube.com and google.com.
SIGNIN_PAGES = {
    "youtube": ("https://accounts.google.com/ServiceLogin?service=youtube",
                ("https://www.youtube.com/", "https://myaccount.google.com/"),
                "www.youtube.com"),
    "instagram": ("https://www.instagram.com/accounts/login/",
                  ("https://www.instagram.com/",),
                  "www.instagram.com"),
    "tiktok": ("https://www.tiktok.com/login",
               ("https://www.tiktok.com/",),
               "www.tiktok.com"),
    "soundcloud": ("https://soundcloud.com/signin",
                   ("https://soundcloud.com/you/library", "https://soundcloud.com/"),
                   "soundcloud.com"),
    "x": ("https://x.com/i/flow/login",
          ("https://x.com/home",),
          "x.com"),
}
# A source whose session belongs to another source. YouTube Music has no login
# of its own — it is youtube.com — so it must read and refresh YouTube's jar
# rather than opening a second window and writing a jar nothing looks for.
COOKIE_SITE_ALIASES = {"ytmusic": "youtube"}


def cookie_site(kind):
    """Whose stored session this source uses."""
    return COOKIE_SITE_ALIASES.get(kind, kind)


_host_matches = host_matches


def normalize_media_url(kind, value):
    """Validate a media source and return the URL yt-dlp should receive."""
    if kind not in MEDIA_LABELS:
        raise ValueError("Unsupported download source.")

    value = (value or "").strip()
    if kind == "tiktok" and value.startswith("@"):
        value = f"https://www.tiktok.com/{value}"
    elif kind == "instagram" and value.startswith("@"):
        value = f"https://www.instagram.com/{value[1:].strip('/')}/"

    parsed = urlsplit(value)
    if parsed.scheme.lower() not in {"http", "https"} or not parsed.hostname:
        raise ValueError("Paste a complete http:// or https:// link.")
    if parsed.username is not None or parsed.password is not None:
        raise ValueError("Links containing a username or password are not supported.")

    host = parsed.hostname.rstrip(".").lower()
    if kind == "general":
        if any(_host_matches(host, domain) for domain in YOUTUBE_DOMAINS):
            raise ValueError("YouTube links must use the dedicated YouTube downloader.")
    elif not MEDIA_PATTERNS[kind].search(value):
        raise ValueError(f"That is not a {MEDIA_LABELS[kind]} link.")
    return value


# What to try when yt-dlp comes back empty-handed, in order. yt-dlp is always
# first and is not listed here. "gallery-dl" resolves photos and mixed posts;
# "direct" treats the link as the file itself, which needs no engine at all.
ENGINES = {
    "x": ("gallery-dl", "html5", "direct"),
    "instagram": ("gallery-dl",),
    "general": ("streamlink", "html5", "gallery-dl", "direct"),
}


def cookie_domain_for(kind, url):
    """The only cookie domain an attempt may borrow for this URL."""
    fixed = COOKIE_DOMAINS.get(kind)
    if fixed:
        return fixed
    if kind != "general":
        return None
    try:
        return (urlsplit(url).hostname or "").lower().rstrip(".") or None
    except ValueError:
        return None

PYTHON_COMPONENTS = {
    "gallery-dl": {"label": "gallery-dl", "module": "gallery_dl", "max_major": 2},
    "streamlink": {"label": "Streamlink", "module": "streamlink", "max_major": 9},
    "img2pdf": {"label": "img2pdf", "module": "img2pdf", "max_major": 1},
    "pillow": {"label": "Pillow", "module": "PIL", "max_major": 13},
    "certifi": {"label": "certifi CA bundle", "module": "certifi", "max_major": 3000},
}
# pikepdf wraps qpdf, a C++ library with no Android build. pypdf does the same
# reading and rewriting in pure Python, so it stands in there; the desktops keep
# pikepdf, which is faster and stricter on damaged files.
# Installed by the Android shell at boot: given a source, re-reads its live
# session from the system CookieManager and rewrites the jar, returning
# (cookies written, whether a session cookie was among them). Left None on the
# desktops, which re-read a real browser profile instead.
SESSION_REFRESHER = None

PDF_LIBRARY = "pypdf" if IS_ANDROID else "pikepdf"
PYTHON_COMPONENTS[PDF_LIBRARY] = {
    "label": PDF_LIBRARY, "module": PDF_LIBRARY,
    "max_major": 7 if IS_ANDROID else 11,
}
if IS_ANDROID:
    # img2pdf requires pikepdf, whose qpdf extension has no Android wheel.
    # Android renders image PDFs with reportlab instead.
    PYTHON_COMPONENTS.pop("img2pdf", None)
    PYTHON_COMPONENTS["reportlab"] = {
        "label": "reportlab", "module": "reportlab", "max_major": 6,
    }
ENGINE_UPDATE_ROOT = (app_data_dir() / "python-engines")
ENGINE_UPDATE_REGISTRY = ENGINE_UPDATE_ROOT / "registry.json"


def activate_engine_updates():
    """Put previously verified component wheel overlays ahead of bundled copies."""
    try:
        registry = json.loads(ENGINE_UPDATE_REGISTRY.read_text("utf-8"))
        root = ENGINE_UPDATE_ROOT.resolve()
    except (OSError, ValueError):
        return
    for name in PYTHON_COMPONENTS:
        try:
            path = Path(registry[name]["path"]).resolve()
            path.relative_to(root)
            if path.is_dir():
                sys.path.insert(0, str(path))
        except (KeyError, OSError, ValueError):
            continue


activate_engine_updates()
# yt-dlp saying "there is no video here" is not the same as yt-dlp failing. The
# first means another engine should look; the second means the job is over.
NO_VIDEO_FOR_YTDLP = re.compile(
    r"no video could be found|there'?s no video|unsupported url|"
    r"no media found|no video formats found|requested format is not available", re.I)

PDF_MAGIC = b"%PDF-"
DOCUMENT_EXTENSIONS = {
    ".pdf", ".epub", ".mobi", ".azw", ".azw3", ".fb2", ".djvu", ".chm",
    ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".ods", ".odp", ".rtf", ".txt", ".md", ".csv", ".cbz",
}
MEDIA_EXTENSIONS = {
    ".mp4", ".m4v", ".webm", ".mkv", ".mov", ".avi", ".flv", ".ts",
    ".mp3", ".m4a", ".aac", ".ogg", ".opus", ".wav", ".flac",
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".avif",
}
MEDIA_CONTENT_TYPES = {
    "application/octet-stream", "application/vnd.apple.mpegurl",
    "application/x-mpegurl", "application/ogg",
}
# A page can link to thousands of files. Without a ceiling one paste turns into
# an unbounded crawl of someone else's server.
PDF_SCAN_LIMIT = 200
MANGA_PAGE_LIMIT = 500
# Real chapter pages are comfortably larger than this. These conservative
# bounds discard spacer GIFs, analytics beacons and favicon-sized placeholders
# without imposing a site-specific page resolution.
MIN_CHAPTER_IMAGE_EDGE = 32
MIN_CHAPTER_IMAGE_AREA = 4096
MEDIA_CONVERT_FORMATS = {
    "mp3", "flac", "wav", "m4a", "ogg", "aac",
    "mp4", "mkv", "webm", "mov", "avi",
}
VIDEO_CONVERT_FORMATS = {"mp4", "mkv", "webm", "mov", "avi"}
VIDEO_CONVERT_CODECS = {"auto", "h264", "h265", "vp9"}
VIDEO_CONVERT_QUALITY = {"master": 14, "high": 18, "balanced": 23, "compact": 28}
VIDEO_UPSCALE_FILTERS = {
    "lanczos": "lanczos+accurate_rnd+full_chroma_int",
    "spline": "spline+accurate_rnd",
    "bicubic": "bicubic",
}
VIDEO_ENHANCEMENT_FILTERS = {
    "none": (),
    "clean": ("hqdn3d=1.5:1.5:6:6", "deband=1thr=0.02:2thr=0.02:3thr=0.02"),
    "detail": ("hqdn3d=0.8:0.8:3:3", "unsharp=5:5:0.65:5:5:0.0"),
    "cinema": ("hqdn3d=1.0:1.0:4:4", "eq=contrast=1.04:saturation=1.06",
               "unsharp=5:5:0.4:5:5:0.0"),
}
AUDIO_ENHANCEMENT_FILTERS = {
    "none": None,
    "normalize": "loudnorm=I=-14:LRA=11:TP=-1.0",
    "music": ("highpass=f=20,lowpass=f=20000,"
              "acompressor=threshold=0.125:ratio=1.6:attack=20:release=250:makeup=1,"
              "loudnorm=I=-14:LRA=11:TP=-1.0"),
    "cleanup": "highpass=f=35,lowpass=f=18000,afftdn=nf=-30,loudnorm=I=-14:LRA=9:TP=-1.0",
}
VIDEO_RESIZE_PRESETS = {
    "source": None, "2160p": (3840, 2160), "1440p": (2560, 1440),
    "1080p": (1920, 1080), "720p": (1280, 720), "480p": (854, 480),
}
IMAGE_DOCUMENT_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".webp", ".avif", ".tif", ".tiff", ".bmp", ".gif",
}
OFFICE_DOCUMENT_EXTENSIONS = {
    ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".ods", ".odp", ".rtf", ".txt", ".csv", ".html", ".htm",
}
EBOOK_CONVERT_EXTENSIONS = {".epub", ".mobi", ".azw", ".azw3", ".fb2"}
# One identity for the whole app: the browser the sign-in window presents, and
# the browser the harvested session is replayed by. A session cookie used from
# a different kind of device than the one that earned it is one of the clearest
# automation signals there is, and answering "not a bot" is what it costs.
#
# On Android the shell hands over the WebView's own agent, because that WebView
# is what the login page actually talks to. The fallbacks below are only for
# when it cannot be read — and the Android one still says Android.
BROWSER_UA = os.environ.get("BLUEKNIGHT_UA", "").strip() or (
    "Mozilla/5.0 (Linux; Android 13; Pixel 7) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/126.0.0.0 Mobile Safari/537.36"
    if IS_ANDROID else
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36")
# Frozen Python builds do not consistently find a usable OpenSSL CA bundle.
# Use certifi explicitly so CDNs whose chain is not present in that embedded
# OpenSSL installation still verify normally.  This retains hostname and
# certificate validation; it is not an unverified-context escape hatch.
def trusted_ssl_context():
    # Imported after activate_engine_updates() so a validated certifi update is
    # used on the next launch instead of the copy frozen into the executable.
    import certifi
    return ssl.create_default_context(cafile=certifi.where())


TRUSTED_SSL_CONTEXT = trusted_ssl_context()


class PdfLinkParser(HTMLParser):
    """Every href and embed/iframe source on a page, in document order.

    Deliberately not a .pdf-suffix filter: plenty of sites serve documents from
    /download?id=… with no extension anywhere in the URL. What a link actually
    is gets decided by fetching it and looking at the bytes.
    """

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.links = []

    def handle_starttag(self, tag, attrs):
        wanted = {"a": "href", "embed": "src", "iframe": "src", "object": "data"}.get(tag)
        if not wanted:
            return
        for name, value in attrs:
            if name == wanted and value:
                self.links.append(value.strip())


class HtmlMediaParser(HTMLParser):
    """Media URLs exposed by ordinary HTML5 markup and social meta tags."""

    META_NAMES = {"og:video", "og:video:url", "og:video:secure_url",
                  "twitter:player:stream", "og:audio", "og:audio:url"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.links = []

    def handle_starttag(self, tag, attrs):
        values = dict(attrs)
        if tag in {"video", "audio", "source"} and values.get("src"):
            self.links.append(values["src"].strip())
        elif tag == "meta":
            key = (values.get("property") or values.get("name") or "").lower()
            if key in self.META_NAMES and values.get("content"):
                self.links.append(values["content"].strip())
        elif tag == "a" and values.get("href"):
            path = urlsplit(values["href"]).path.lower()
            if any(path.endswith(ext) for ext in MEDIA_EXTENSIONS):
                self.links.append(values["href"].strip())


class HtmlImageParser(HTMLParser):
    """Likely page images for sites without a gallery-dl extractor."""

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.links = []

    def handle_starttag(self, tag, attrs):
        values = dict(attrs)
        if tag == "img":
            value = (values.get("data-src") or values.get("data-original")
                     or values.get("data-lazy-src") or values.get("data-url"))
            srcset = (values.get("data-srcset") or values.get("data-lazy-srcset")
                      or values.get("srcset"))
            if not value and srcset:
                value = srcset.split(",")[-1].strip().split()[0]
            value = value or values.get("src")
            if value:
                self.links.append(value.strip())
        elif tag == "source" and (values.get("data-srcset") or values.get("srcset")):
            srcset = values.get("data-srcset") or values["srcset"]
            value = srcset.split(",")[-1].strip().split()[0]
            if value:
                self.links.append(value)
        elif tag == "a" and values.get("href"):
            path = urlsplit(values["href"]).path.lower()
            if Path(path).suffix in {".jpg", ".jpeg", ".png", ".webp", ".avif"}:
                self.links.append(values["href"].strip())


def normalize_page_url(value):
    """Validate any http(s) address a person pasted. Returns the URL."""
    value = (value or "").strip()
    parsed = urlsplit(value)
    if parsed.scheme.lower() not in {"http", "https"} or not parsed.hostname:
        raise ValueError("Paste a complete http:// or https:// link.")
    if parsed.username is not None or parsed.password is not None:
        raise ValueError("Links containing a username or password are not supported.")
    return value


def safe_filename(name, fallback="document", suffix=".pdf"):
    """A remote-supplied name, reduced to something safe to write.

    The name comes from a server, so it is untrusted input: "../../boot.ini" and
    "CON" are both things a hostile Content-Disposition can say. Only the basename
    survives, and only characters Windows actually accepts.
    """
    name = unquote(name or "").replace("\\", "/").split("/")[-1].strip()
    name = re.sub(r'[<>:"|?*\x00-\x1f]', "_", name).strip(". ")
    stem, _, ext = name.rpartition(".")
    if not stem:
        stem, ext = name, ext or suffix.lstrip(".")
    if stem.upper() in {"CON", "PRN", "AUX", "NUL", *(f"COM{i}" for i in range(1, 10)),
                        *(f"LPT{i}" for i in range(1, 10))}:
        stem = f"_{stem}"
    stem = (stem or fallback)[:120]
    ext = (ext if ext.lower() == suffix.lstrip(".") else suffix.lstrip("."))
    return f"{stem}.{ext}"


def detect_document_extension(head, url, content_type=""):
    """Return a validated document suffix, or None for non-document content."""
    path_suffix = Path(urlsplit(url).path).suffix.lower()
    content_type = (content_type or "").split(";", 1)[0].strip().lower()
    if PDF_MAGIC in head[:1024]:
        return ".pdf"
    if head.startswith(b"PK\x03\x04"):
        return ".zip"
    if len(head) >= 68 and head[60:68] == b"BOOKMOBI":
        return path_suffix if path_suffix in {".mobi", ".azw", ".azw3"} else ".mobi"
    if b"<FictionBook" in head[:4096] and path_suffix == ".fb2":
        return ".fb2"
    if head.startswith(b"{\\rtf"):
        return ".rtf"
    if head.startswith(b"AT&TFORM") and b"DJVU" in head[:32]:
        return ".djvu"
    if head.startswith(b"ITSF"):
        return ".chm"
    if head.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        office = {
            "application/msword": ".doc",
            "application/vnd.ms-excel": ".xls",
            "application/vnd.ms-powerpoint": ".ppt",
        }
        return office.get(content_type, path_suffix if path_suffix in {".doc", ".xls", ".ppt"}
                          else ".doc")
    text_types = {".txt": {"text/plain"}, ".md": {"text/plain", "text/markdown"},
                  ".csv": {"text/csv", "text/plain"}}
    if path_suffix in text_types and content_type in text_types[path_suffix]:
        return path_suffix
    return None


def is_valid_pdf_file(path):
    """Check the structural markers required of a complete PDF file."""
    try:
        path = Path(path)
        size = path.stat().st_size
        if size < len(PDF_MAGIC) + len(b"%%EOF"):
            return False
        with path.open("rb") as handle:
            head = handle.read(1024)
            handle.seek(max(0, size - 4096))
            tail = handle.read()
        # ISO 32000 permits the header within the first 1024 bytes and requires
        # an end-of-file marker near the end. This catches HTML error bodies and
        # the common case where a connection closes halfway through a PDF.
        return PDF_MAGIC in head and b"%%EOF" in tail
    except OSError:
        return False


def inspect_zip_document(path):
    """Identify a validated ZIP-based document container."""
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        mimetype = archive.read("mimetype").strip() if "mimetype" in names else b""
        if archive.testzip() is not None:
            return None
        if (mimetype == b"application/epub+zip"
                and "META-INF/container.xml" in names):
            return ".epub"
        odf = {
            b"application/vnd.oasis.opendocument.text": ".odt",
            b"application/vnd.oasis.opendocument.spreadsheet": ".ods",
            b"application/vnd.oasis.opendocument.presentation": ".odp",
        }
        if mimetype in odf:
            return odf[mimetype]
        if "[Content_Types].xml" in names:
            if any(name.startswith("word/") for name in names):
                return ".docx"
            if any(name.startswith("xl/") for name in names):
                return ".xlsx"
            if any(name.startswith("ppt/") for name in names):
                return ".pptx"
        image_exts = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".avif"}
        files = [name for name in names if name and not name.endswith("/")]
        if files and sum(Path(name).suffix.lower() in image_exts for name in files) >= max(1, len(files) // 2):
            return ".cbz"
    return None


def unique_path(folder, filename):
    """A path that does not overwrite an existing download."""
    candidate = folder / filename
    if not candidate.exists():
        return candidate
    stem, suffix = candidate.stem, candidate.suffix
    for n in range(2, 1000):
        candidate = folder / f"{stem} ({n}){suffix}"
        if not candidate.exists():
            return candidate
    return folder / f"{stem} ({int(time.time())}){suffix}"


def natural_sort_key(value):
    """Case-insensitive numeric ordering: page2 precedes page10."""
    return tuple((0, int(part)) if part.isdigit() else (1, part.casefold())
                 for part in re.split(r"(\d+)", str(value)))


def versions_equal(left, right):
    """Compare numeric release strings without caring about zero padding."""
    left = str(left or "").strip().lstrip("vV")
    right = str(right or "").strip().lstrip("vV")
    if re.fullmatch(r"\d+(?:\.\d+)*", left) and re.fullmatch(r"\d+(?:\.\d+)*", right):
        return tuple(map(int, left.split("."))) == tuple(map(int, right.split(".")))
    return left.casefold() == right.casefold()


def unique_chapter_paths(folder, name):
    """Reserve a matching page folder, PDF and CBZ name without overwriting."""
    stem = Path(safe_filename(name, fallback="chapter", suffix=".pdf")).stem
    for number in range(1, 1000):
        candidate = stem if number == 1 else f"{stem} ({number})"
        page_folder = folder / candidate
        pdf = folder / f"{candidate}.pdf"
        cbz = folder / f"{candidate}.cbz"
        if not any(path.exists() for path in (page_folder, pdf, cbz)):
            return page_folder, pdf, cbz
    stamp = int(time.time())
    return (folder / f"{stem} ({stamp})", folder / f"{stem} ({stamp}).pdf",
            folder / f"{stem} ({stamp}).cbz")


def validate_engine_update(paths):
    """Import and exercise a staged Python-engine set in a fresh process."""
    try:
        import importlib

        for path in reversed(paths):
            resolved = str(Path(path).resolve())
            if resolved not in sys.path:
                sys.path.insert(0, resolved)
        importlib.invalidate_caches()

        # certifi is needed while this file initializes TLS, so discard all
        # already-imported engine modules before testing the staged overlays.
        roots = {spec["module"].split(".", 1)[0] for spec in PYTHON_COMPONENTS.values()}
        for module_name in list(sys.modules):
            if module_name.split(".", 1)[0] in roots:
                del sys.modules[module_name]

        certifi_module = importlib.import_module("certifi")
        gallery_module = importlib.import_module("gallery_dl")
        streamlink_module = importlib.import_module("streamlink")
        img2pdf_module = importlib.import_module("img2pdf")
        pdf_module = importlib.import_module(PDF_LIBRARY)
        image_module = importlib.import_module("PIL.Image")

        if not Path(certifi_module.where()).is_file():
            raise RuntimeError("certifi did not provide a CA bundle")
        if not getattr(gallery_module, "__version__", None):
            raise RuntimeError("gallery-dl did not report a version")
        if not getattr(streamlink_module, "__version__", None):
            raise RuntimeError("Streamlink did not report a version")

        # Exercise the document stack together; importing alone does not catch
        # incompatible Pillow/img2pdf/PDF-library binary combinations.
        with tempfile.TemporaryDirectory() as temp:
            image_path = Path(temp) / "probe.png"
            image_module.new("RGB", (32, 32), "white").save(image_path)
            payload = img2pdf_module.convert([str(image_path)])
            if not payload.startswith(PDF_MAGIC):
                raise RuntimeError("img2pdf returned invalid output")
            if pdf_pages(io.BytesIO(payload)) != 1:
                raise RuntimeError(f"{PDF_LIBRARY} could not read the img2pdf probe")

        # This resolves a built-in URL scheme and does not touch the network.
        stream_name, _, _ = streamlink_module.Streamlink().resolve_url(
            "hls://example.com/probe.m3u8")
        if stream_name != "hls":
            raise RuntimeError("Streamlink plugin loading failed")
        return True
    except Exception:
        return False


def write_netscape_jar(cookies, path):
    """Save browser cookies as the Netscape file yt-dlp reads. Returns the count."""
    lines = ["# Netscape HTTP Cookie File",
             "# Written by Blue Knight Downloader. Editing this is not needed.", ""]
    written = 0
    for cookie in cookies:
        for name, morsel in cookie.items():
            domain = morsel["domain"] or ""
            if not domain:
                continue
            expires = morsel["expires"]
            if expires:
                with contextlib.suppress(Exception):
                    expires = int(time.mktime(
                        time.strptime(expires, "%a, %d %b %Y %H:%M:%S %Z")))
            lines.append("\t".join([
                domain,
                "TRUE" if domain.startswith(".") else "FALSE",
                morsel["path"] or "/",
                "TRUE" if morsel["secure"] else "FALSE",
                str(expires or int(time.time()) + 60 * 60 * 24 * 180),
                name,
                morsel.value,
            ]))
            written += 1
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return written


# Reading the clipboard is the one thing every desktop does through a different
# program. Android has no command for it at all and answers in the shell instead.
CLIPBOARD_READERS = {
    "win": (["powershell", "-NoProfile", "-Command", "Get-Clipboard -Raw"],),
    "mac": (["pbpaste"],),
    # Wayland first: on a Wayland session xclip either fails or reads a stale
    # X11 selection, while wl-paste is absent on a pure X11 one and falls through.
    "linux": (["wl-paste", "--no-newline"], ["xclip", "-selection", "clipboard", "-o"],
              ["xsel", "--clipboard", "--output"]),
    "android": (),
}


def read_clipboard():
    """Clipboard text, empty when there is nothing to paste."""
    for command in CLIPBOARD_READERS[OS_TAG]:
        try:
            out = pyshell.run(command, stdout=subprocess.PIPE, stderr=subprocess.DEVNULL,
                              timeout=6)
            if out.returncode == 0 and out.stdout.strip():
                return out.stdout.strip()
        except Exception:
            continue
    return ""


# Writing takes a different program from reading on every desktop. Android has
# neither and answers in the shell, which owns the ClipboardManager.
CLIPBOARD_WRITERS = {
    "win": (["clip"],),
    "mac": (["pbcopy"],),
    "linux": (["wl-copy"], ["xclip", "-selection", "clipboard"],
              ["xsel", "--clipboard", "--input"]),
    "android": (),
}


def write_clipboard(text):
    """Put text on the clipboard. True when a program accepted it."""
    for command in CLIPBOARD_WRITERS[OS_TAG]:
        try:
            child = pyshell.popen(command, stdin=subprocess.PIPE,
                                  stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            child.communicate(text, timeout=6)
            if child.returncode == 0:
                return True
        except Exception:
            continue
    return False


class Var:
    """Minimal observable value, standing in for the old Tk variables."""

    def __init__(self, value=""):
        self._value = value

    def get(self):
        return self._value

    def set(self, value):
        self._value = value


def android_shim(target, name):
    """Give a packaged Android executable back its real name.

    Two rules collide on Android. The system only unpacks a bundled binary, and
    only marks it executable, when it sits in the native library directory
    under a lib*.so name. yt-dlp, meanwhile, identifies a program by its file
    name — it will never accept "libffmpeg.so" as ffmpeg, so merging a video
    with its audio track and writing an MP3 would both fail with the tool
    apparently missing.

    A symlink satisfies both: the link carries the name the tool is looked up
    by, and the target stays the file the kernel is willing to execute. The
    link cannot be a copy, because app storage is mounted no-exec from API 29.
    """
    link = app_data_dir() / "bin" / name
    try:
        if link.is_symlink() and link.resolve() == Path(target).resolve():
            return link
        link.parent.mkdir(parents=True, exist_ok=True)
        link.unlink(missing_ok=True)
        link.symlink_to(target)
        return link
    except OSError:
        # No symlink support: the raw path still runs, it just cannot be found
        # by name — which the caller reports rather than failing silently.
        return Path(target)


def bundled_tool(name):
    """Find a packaged tool in both PyInstaller one-file and one-folder layouts."""
    # Android unpacks a bundled executable only from the native library
    # directory, and only when it is named like a library. The shell passes
    # that directory in; the lib*.so name is a packaging rule, not a format.
    tools_override = os.environ.get("BLUEKNIGHT_TOOLS")
    if tools_override:
        stem = Path(name).stem
        packaged = Path(tools_override) / f"lib{stem}.so"
        if packaged.is_file():
            return str(android_shim(packaged, stem))
        direct = Path(tools_override) / name
        if direct.is_file():
            return str(direct)

    roots = []
    bundle_root = getattr(sys, "_MEIPASS", None)
    if bundle_root:
        roots.append(Path(bundle_root))
    if getattr(sys, "frozen", False):
        executable_dir = Path(sys.executable).resolve().parent
        roots.extend((executable_dir, executable_dir / "_internal"))
    roots.append(Path(__file__).resolve().parent)

    seen = set()
    for root in roots:
        for candidate in (root / "tools" / name, root / "vendor" / name, root / name):
            key = str(candidate)
            if key not in seen and candidate.is_file():
                return str(candidate)
            seen.add(key)
    return None


def pdf_pages(source):
    """How many pages a PDF holds. Raises if it cannot be read at all.

    Two libraries answer this, and which one is installed depends on the
    platform rather than on anything the caller cares about — see PDF_LIBRARY.
    """
    if PDF_LIBRARY == "pikepdf":
        import pikepdf
        with pikepdf.Pdf.open(source) as document:
            return len(document.pages)
    from pypdf import PdfReader
    return len(PdfReader(source).pages)


def pdf_rewrite(source, target):
    """Read a PDF and write it back out, which is what proves it is intact.

    A file that opens, parses and re-serialises is a real PDF; one that was
    renamed or truncated fails here rather than on the reader's device.
    """
    if PDF_LIBRARY == "pikepdf":
        import pikepdf
        with pikepdf.Pdf.open(source) as document:
            document.save(target)
        return

    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(source)
    if not reader.pages:
        raise RuntimeError("The PDF holds no pages")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(target, "wb") as handle:
        writer.write(handle)


def make_executable(path):
    """Give a freshly downloaded tool the permission bit it arrives without.

    A release asset is a plain file: Windows infers runnability from the name,
    every other platform from the mode, and an archive member's recorded mode
    is not carried through the copy that verifies its checksum.
    """
    if IS_WINDOWS:
        return
    with contextlib.suppress(OSError):
        path = Path(path)
        path.chmod(path.stat().st_mode | 0o755)


def extract_binaries(archive_path, wanted):
    """Pull named programs out of a release archive.

    wanted maps a bare program name to where it should land. Members are matched
    on their base name rather than a full path, because the same binary sits at
    `ffmpeg-.../bin/ffmpeg` in one project's archive and at the root of another's.
    Each member is streamed rather than fully extracted, which verifies that
    member's checksum without inflating the rest of a large distribution.
    """
    archive_path = Path(archive_path)
    targets = {name.lower(): Path(target) for name, target in wanted.items()}
    found = set()

    def base(member):
        return Path(member.replace("\\", "/")).name.lower()

    def claim(member, opener):
        stem, suffix = os.path.splitext(base(member))
        # An archive also ships documentation, and `doc/ffmpeg.html` has the
        # same stem as the program. Only a program-shaped name qualifies.
        if suffix not in ("", ".exe"):
            return
        target = targets.get(stem)
        if target is None or stem in found:
            return
        with opener() as source, target.open("wb") as output:
            shutil.copyfileobj(source, output, length=1024 * 1024)
        found.add(stem)

    if zipfile.is_zipfile(archive_path):
        with zipfile.ZipFile(archive_path) as archive:
            for member in archive.namelist():
                claim(member, lambda m=member: archive.open(m))
    else:
        # tarfile picks the compression itself, which is what makes the same
        # call work for the .tar.xz Linux builds and any .tar.gz mirror.
        with tarfile.open(archive_path) as archive:
            for member in archive.getmembers():
                if member.isfile():
                    claim(member.name,
                          lambda m=member: contextlib.closing(archive.extractfile(m)))

    missing = sorted(set(targets) - found)
    if missing:
        raise RuntimeError(
            f"{archive_path.name} is missing {', '.join(missing)}")


def normalize_proxy_url(value, proxy_type="http"):
    """Validate a proxy URL and add the selected scheme when it is omitted."""
    value = value.strip()
    if not value:
        raise ValueError("Proxy URL is empty")
    if "://" not in value:
        value = f"{proxy_type}://{value}"

    parsed = urlsplit(value)
    scheme = parsed.scheme.lower()
    if scheme not in {"http", "https", "socks5", "socks5h"}:
        raise ValueError("Proxy scheme must be http, https, socks5, or socks5h")
    if not parsed.hostname:
        raise ValueError("Proxy URL must include a host")
    try:
        port = parsed.port
    except ValueError as exc:
        raise ValueError("Proxy URL has an invalid port") from exc
    if port is None:
        port = 1080 if scheme.startswith("socks5") else 8080

    host = parsed.hostname
    if ":" in host:
        host = f"[{host}]"
    auth = ""
    if parsed.username is not None:
        auth = quote(unquote(parsed.username), safe="")
        if parsed.password is not None:
            auth += ":" + quote(unquote(parsed.password), safe="")
        auth += "@"
    return f"{scheme}://{auth}{host}:{port}"


def display_proxy_url(value):
    """Hide proxy credentials before displaying the URL in the activity log."""
    parsed = urlsplit(value)
    host = parsed.hostname or ""
    if ":" in host:
        host = f"[{host}]"
    auth = "***@" if parsed.username is not None else ""
    return f"{parsed.scheme}://{auth}{host}:{parsed.port}"


def _recv_exact(sock, size):
    data = bytearray()
    while len(data) < size:
        chunk = sock.recv(size - len(data))
        if not chunk:
            raise ConnectionError("Proxy closed the connection")
        data.extend(chunk)
    return bytes(data)


class Socks5HttpBridge:
    """Expose a local HTTP CONNECT proxy backed by a SOCKS5 proxy."""

    def __init__(self, proxy_url):
        parsed = urlsplit(normalize_proxy_url(proxy_url, "socks5"))
        if not parsed.scheme.startswith("socks5"):
            raise ValueError("The local bridge requires a SOCKS5 proxy")
        self.host = parsed.hostname
        self.port = parsed.port
        self.username = unquote(parsed.username) if parsed.username else None
        self.password = unquote(parsed.password) if parsed.password else ""
        self.server = None
        self.thread = None

    def _open_target(self, target_host, target_port):
        sock = socket.create_connection((self.host, self.port), timeout=15)
        methods = b"\x00\x02" if self.username is not None else b"\x00"
        sock.sendall(b"\x05" + bytes([len(methods)]) + methods)
        version, method = _recv_exact(sock, 2)
        if version != 5 or method == 0xFF:
            sock.close()
            raise ConnectionError("SOCKS5 proxy rejected authentication methods")
        if method == 2:
            if self.username is None:
                sock.close()
                raise ConnectionError("SOCKS5 proxy requires a username and password")
            user = self.username.encode("utf-8")
            password = self.password.encode("utf-8")
            if len(user) > 255 or len(password) > 255:
                sock.close()
                raise ValueError("SOCKS5 credentials are too long")
            sock.sendall(b"\x01" + bytes([len(user)]) + user + bytes([len(password)]) + password)
            if _recv_exact(sock, 2)[1] != 0:
                sock.close()
                raise ConnectionError("SOCKS5 authentication failed")
        elif method != 0:
            sock.close()
            raise ConnectionError("SOCKS5 proxy selected an unsupported authentication method")

        try:
            address = ipaddress.ip_address(target_host)
            atyp = 1 if address.version == 4 else 4
            encoded_host = address.packed
        except ValueError:
            encoded_host = target_host.encode("idna")
            if len(encoded_host) > 255:
                sock.close()
                raise ValueError("Target hostname is too long")
            atyp = 3
            encoded_host = bytes([len(encoded_host)]) + encoded_host

        sock.sendall(b"\x05\x01\x00" + bytes([atyp]) + encoded_host + target_port.to_bytes(2, "big"))
        version, result, _, bound_type = _recv_exact(sock, 4)
        if version != 5 or result != 0:
            sock.close()
            raise ConnectionError(f"SOCKS5 connection failed (code {result})")
        if bound_type == 1:
            _recv_exact(sock, 4)
        elif bound_type == 4:
            _recv_exact(sock, 16)
        elif bound_type == 3:
            _recv_exact(sock, _recv_exact(sock, 1)[0])
        else:
            sock.close()
            raise ConnectionError("SOCKS5 proxy returned an invalid address")
        _recv_exact(sock, 2)
        sock.settimeout(None)
        return sock

    @staticmethod
    def _split_authority(authority, default_port):
        parsed = urlsplit("//" + authority)
        if not parsed.hostname:
            raise ValueError("Invalid proxy request target")
        return parsed.hostname, parsed.port or default_port

    @staticmethod
    def _relay(client, upstream):
        sockets = (client, upstream)
        while True:
            readable, _, _ = select.select(sockets, [], [], 60)
            if not readable:
                return
            for source in readable:
                data = source.recv(65536)
                if not data:
                    return
                (upstream if source is client else client).sendall(data)

    def start(self):
        bridge = self

        class Handler(socketserver.BaseRequestHandler):
            def handle(self):
                self.request.settimeout(15)
                header = bytearray()
                while b"\r\n\r\n" not in header:
                    chunk = self.request.recv(4096)
                    if not chunk:
                        return
                    header.extend(chunk)
                    if len(header) > 65536:
                        raise ValueError("HTTP proxy request headers are too large")

                head, body = bytes(header).split(b"\r\n\r\n", 1)
                lines = head.split(b"\r\n")
                method, target, version = lines[0].decode("latin-1").split(" ", 2)
                if method.upper() == "CONNECT":
                    host, port = bridge._split_authority(target, 443)
                    upstream = bridge._open_target(host, port)
                    self.request.sendall(b"HTTP/1.1 200 Connection Established\r\n\r\n")
                else:
                    parsed = urlsplit(target)
                    if not parsed.hostname:
                        raise ValueError("HTTP proxy requests must use an absolute URL")
                    host = parsed.hostname
                    port = parsed.port or (443 if parsed.scheme == "https" else 80)
                    upstream = bridge._open_target(host, port)
                    path = parsed.path or "/"
                    if parsed.query:
                        path += "?" + parsed.query
                    clean_headers = [line for line in lines[1:] if not line.lower().startswith(b"proxy-connection:")]
                    request_head = f"{method} {path} {version}\r\n".encode("latin-1")
                    upstream.sendall(request_head + b"\r\n".join(clean_headers) + b"\r\n\r\n" + body)
                try:
                    bridge._relay(self.request, upstream)
                finally:
                    upstream.close()

        class Server(socketserver.ThreadingTCPServer):
            allow_reuse_address = True
            daemon_threads = True

            def handle_error(self, request, client_address):
                pass

        self.server = Server(("127.0.0.1", 0), Handler)
        self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
        self.thread.start()
        return f"http://127.0.0.1:{self.server.server_address[1]}"

    def close(self):
        if self.server:
            self.server.shutdown()
            self.server.server_close()
            self.server = None


class SpotifyDownloader:
    """Download engine. The window talks to it through the js_api bridge."""

    def __init__(self):
        self.window = None

        # Performance optimization: Batching
        self.ui_queue = queue.Queue()
        self.log_batch = []
        self.last_ui_update = 0
        self.process = None
        self.is_downloading = False
        self.is_converting = False
        self.closing = False
        self.setup_running = False
        self.executor = ThreadPoolExecutor(max_workers=3)

        # Store FFmpeg in the per-user application data directory
        self.app_data = app_data_dir()
        self.ffmpeg_dir = self.app_data / "ffmpeg"
        self.ffmpeg_exe = self.ffmpeg_dir / "bin" / f"ffmpeg{EXE}"
        self.tools_dir = self.app_data / "tools"
        self.spotdl_exe = self.tools_dir / f"spotdl{EXE}"

        # Proxy settings with proper defaults for v2rayN
        self.download_format = Var("mp3")
        self.quality = Var("320")
        self.use_proxy = Var(False)
        self.proxy_type = Var("http")  # http or socks5
        self.proxy_url = Var("http://127.0.0.1:10809")  # v2rayN HTTP default
        self.youtube_po_token = Var("")
        self.save_folder = Var(str(download_root()))
        self.status = Var("Preparing")

        self.spotdl_cmd = None
        self.ffmpeg_cmd = None
        self.ytdlp_cmd = None
        self.deno_cmd = None
        self.gallerydl_version = None
        self.streamlink_version = None
        self.updating = False

        # Pre-compile regex patterns for performance
        self._re_url = re.compile(r"https?://open\.spotify\.com/(track|album|playlist)/", re.I)
        self._re_found_songs = re.compile(r"Found\s+(\d+)\s+songs", re.I)
        self._re_downloading = re.compile(r"Downloading", re.I)
        self._re_downloaded = re.compile(r"Downloaded", re.I)
        self._re_brackets = re.compile(r"\[.*?\]")
        self._re_error = re.compile(r"(ERROR|Error|error)")
        self._re_warning = re.compile(r"(WARNING|Warning|warning)")
        self._re_request_error = re.compile(r"(RequestError|Failed to complete request)")
        self._re_ytdlp_item = re.compile(r"Downloading item (\d+) of (\d+)", re.I)

        self.hello_sent = False
        self._session_cookies = set()
        self._login_window = None
        self._login_kind = None
        self._media_done = 0
        self._media_total = None
        self._stream_index = 0

    # ------------------------------------------------------------------
    # State the page reads through the bridge
    # ------------------------------------------------------------------
    def snapshot(self):
        return {
            "format": self.download_format.get(),
            "quality": self.quality.get(),
            "proxy_enabled": bool(self.use_proxy.get()),
            "proxy_type": self.proxy_type.get(),
            "proxy_url": self.proxy_url.get(),
            "youtube_po_token": self.youtube_po_token.get(),
            "folder": self.save_folder.get(),
        }

    def begin(self):
        """First contact from the page: greet, then start the dependency check."""
        if self.hello_sent:
            return
        self.hello_sent = True
        self.log("Welcome. Paste a Spotify link to begin.", "info")
        self.log("Behind a filter? Set a proxy in Connection settings.", "info")
        threading.Thread(target=self.start_automatic_setup, daemon=True).start()

    def drain(self):
        """Drain queued events for the page. Called on a short interval."""
        events = []
        while True:
            try:
                action, args = self.ui_queue.get_nowait()
            except queue.Empty:
                break
            if action == "batch_log":
                # log() queues one batch_log per line, so a chatty download used
                # to arrive as hundreds of separate events, each costing the page
                # its own console write. Runs of them fold into one.
                lines = [list(pair) for pair in args]
                if events and events[-1][0] == "log":
                    events[-1][1].extend(lines)
                else:
                    events.append(["log", lines])
            elif action == "finished":
                self.download_finished()
                events.append(["finished", []])
            else:
                events.append([action, list(args)])
        return events

    def set_option(self, name, value):
        options = {
            "format": self.download_format,
            "quality": self.quality,
            "proxy_type": self.proxy_type,
            "proxy_url": self.proxy_url,
            "proxy_enabled": self.use_proxy,
            "youtube_po_token": self.youtube_po_token,
            "folder": self.save_folder,
        }
        var = options.get(name)
        if var is not None:
            var.set(value)

    def notify(self, message, kind="info"):
        self.ui("toast", message, kind)

    def open_folder(self):
        folder = Path(self.save_folder.get()).expanduser()
        folder.mkdir(parents=True, exist_ok=True)
        try:
            if sys.platform == "win32":
                os.startfile(str(folder))
            elif sys.platform == "darwin":
                subprocess.Popen(["open", str(folder)])
            else:
                subprocess.Popen(["xdg-open", str(folder)])
        except Exception as exc:
            self.notify(f"Could not open the folder: {exc}", "err")

    # ------------------------------------------------------------------
    # Proxy helpers
    # ------------------------------------------------------------------
    def set_proxy_preset(self, proxy_type, url):
        """Set proxy to a preset configuration"""
        self.proxy_type.set(proxy_type)
        self.proxy_url.set(url)
        self.use_proxy.set(True)
        self.log(f"Proxy preset set: {proxy_type.upper()} - {url}", "info")
        
        if proxy_type == "socks5":
            self.log("SOCKS5 will use the built-in local bridge.", "info")
        else:
            self.log("HTTP proxy ready for spotDL", "success")

    def test_proxy(self):
        """Make an HTTPS request through the configured proxy."""
        try:
            proxy = normalize_proxy_url(self.proxy_url.get(), self.proxy_type.get())
        except ValueError as exc:
            self.notify(str(exc), "err")
            return

        self.log(f"Testing proxy: {display_proxy_url(proxy)}...", "info")

        def do_test():
            bridge = None
            try:
                effective_proxy = proxy
                if urlsplit(proxy).scheme.startswith("socks5"):
                    bridge = Socks5HttpBridge(proxy)
                    effective_proxy = bridge.start()
                handler = urllib.request.ProxyHandler({"http": effective_proxy, "https": effective_proxy})
                opener = urllib.request.build_opener(handler)
                request = urllib.request.Request(
                    "https://api.ipify.org?format=json",
                    headers={"User-Agent": f"{APP_NAME}/{APP_VERSION}"},
                )
                with opener.open(request, timeout=15) as response:
                    result = response.read().decode("utf-8", "replace")
                self.log(f"Proxy test succeeded: {result[:120]}", "success")
                self.notify("Proxy connection works.", "ok")
            except Exception as exc:
                detail = str(exc)[:200]
                self.log(f"Proxy test failed: {detail}", "error")
                self.notify(
                    f"Proxy test failed: {detail}\n\n"
                    "Check that the proxy app is running and that the type and port match its local settings.",
                    "err",
                )
            finally:
                if bridge:
                    bridge.close()

        threading.Thread(target=do_test, daemon=True).start()

    # ------------------------------------------------------------------
    # Setup (unchanged)
    # ------------------------------------------------------------------
    def start_automatic_setup(self):
        if self.setup_running:
            return
        self.setup_running = True
        threading.Thread(target=self.setup_dependencies, daemon=True).start()

    def setup_dependencies(self):
        try:
            self.ui("setup", "Checking FFmpeg...", "This is required for audio conversion")
            self.ffmpeg_cmd = self.get_or_download_ffmpeg()
            
            if not self.ffmpeg_cmd:
                raise RuntimeError("Failed to setup FFmpeg")

            if IS_ANDROID:
                # spotDL depends on native Python extensions which have no
                # Android wheels. The Android path below resolves Spotify
                # metadata and downloads the matching YouTube Music audio.
                self.spotdl_cmd = "Android Spotify engine"
            else:
                self.ui("setup", "Checking spotDL...", "This downloads Spotify tracks")
                self.spotdl_cmd = self.get_or_download_tool("spotdl")
                if not self.spotdl_cmd:
                    raise RuntimeError("spotDL was installed, but its command could not be located.")

            self.ui("setup", "Checking yt-dlp...", "This downloads YouTube, TikTok, Instagram and video-site media")
            self.ytdlp_cmd = self.get_or_download_tool("yt-dlp")
            if not self.ytdlp_cmd:
                raise RuntimeError("yt-dlp was installed, but its command could not be located.")

            if not IS_ANDROID:
                self.ui("setup", "Checking YouTube challenge runtime...",
                        "Deno runs yt-dlp's bundled EJS challenge solver")
                self.deno_cmd = self.get_or_download_tool("deno")
                if not self.deno_cmd:
                    raise RuntimeError("Deno was installed, but its command could not be located.")

            self.ui("setup", "Checking fallback engines...",
                    "This adds live streams, HTML galleries and mixed media")
            try:
                import gallery_dl
                import streamlink
                self.gallerydl_version = gallery_dl.__version__
                self.streamlink_version = streamlink.__version__
            except ImportError as exc:
                raise RuntimeError(f"A bundled fallback engine is missing: {exc.name}") from exc

            self.ui("setup_done", "Ready to download", "Everything is installed and up to date.")
            self.ui("components", self.component_summary())
            self.log("✓ FFmpeg: " + str(self.ffmpeg_cmd), "success")
            self.log("✓ spotDL: " + str(self.spotdl_cmd), "success")
            self.log("✓ yt-dlp: " + str(self.ytdlp_cmd), "success")
            if self.deno_cmd:
                self.log(f"✓ Deno: {self.deno_cmd}", "success")
            self.log(f"✓ gallery-dl: {self.gallerydl_version}", "success")
            self.log(f"✓ Streamlink: {self.streamlink_version}", "success")
            self.log("💡 Use 'V2RayN HTTP' preset for proxy support", "info")

        except Exception as exc:
            self.ui("setup_error", "Automatic setup needs attention", str(exc))
        finally:
            self.setup_running = False

    @contextlib.contextmanager
    def _net_opener(self):
        """A urllib opener that honours the configured proxy.

        Fetching tools is exactly the traffic a restricted network blocks, so
        it goes through the same proxy the downloads themselves use.
        """
        bridge = None
        try:
            handlers = [urllib.request.HTTPSHandler(context=TRUSTED_SSL_CONTEXT)]
            if self.use_proxy.get() and self.proxy_url.get().strip():
                proxy = normalize_proxy_url(self.proxy_url.get(), self.proxy_type.get())
                if urlsplit(proxy).scheme.startswith("socks5"):
                    bridge = Socks5HttpBridge(proxy)
                    proxy = bridge.start()
                handlers.insert(
                    0, urllib.request.ProxyHandler({"http": proxy, "https": proxy}))
            yield urllib.request.build_opener(*handlers)
        finally:
            if bridge:
                bridge.close()

    def _download_with_fallbacks(self, urls, destination, label):
        """Download atomically, trying independent hosts in order."""
        destination.parent.mkdir(parents=True, exist_ok=True)
        partial = destination.with_suffix(destination.suffix + ".part")
        errors = []
        with self._net_opener() as opener:
            for attempt, url in enumerate(urls, 1):
                try:
                    partial.unlink(missing_ok=True)
                    host = urlsplit(url).hostname or "download server"
                    self.ui("setup", f"Downloading {label}...", f"Trying {host} ({attempt}/{len(urls)})")
                    self.log(f"Downloading {label} from {host} ({attempt}/{len(urls)})…", "info")
                    request = urllib.request.Request(url, headers={"User-Agent": f"{APP_NAME}/{APP_VERSION}"})
                    with opener.open(request, timeout=45) as response, partial.open("wb") as output:
                        shutil.copyfileobj(response, output, length=1024 * 1024)
                    if partial.stat().st_size == 0:
                        raise RuntimeError("server returned an empty file")
                    os.replace(partial, destination)
                    return destination
                except Exception as exc:
                    errors.append(f"{urlsplit(url).hostname}: {exc}")
                    self.log(f"{label} source {attempt} failed; trying the next source.", "warning")
        partial.unlink(missing_ok=True)
        raise RuntimeError(f"All {label} download sources failed: " + " | ".join(errors))

    def latest_version(self, name):
        """Published version for a tool, or None when the check cannot run."""
        api = TOOLS[name].get("api")
        if not api:
            return None
        try:
            with self._net_opener() as opener:
                request = urllib.request.Request(api, headers={
                    "User-Agent": f"{APP_NAME}/{APP_VERSION}",
                    "Accept": "application/vnd.github+json",
                })
                with opener.open(request, timeout=20) as response:
                    tag = json.loads(response.read().decode("utf-8", "replace")).get("tag_name", "")
            return tag.strip().lstrip("vV") or None
        except Exception:
            return None

    @staticmethod
    def python_engine_version(module_name):
        try:
            module = __import__(module_name)
            return str(module.__version__)
        except (ImportError, AttributeError):
            return None

    def latest_python_engine(self, name):
        """Latest supported wheel metadata from PyPI's HTTPS index."""
        spec = PYTHON_COMPONENTS[name]
        url = f"https://pypi.org/pypi/{name}/json"
        with self._net_opener() as opener:
            with opener.open(urllib.request.Request(url, headers={
                    "User-Agent": f"{APP_NAME}/{APP_VERSION}",
                    "Accept": "application/json"}), timeout=30) as response:
                payload = json.loads(response.read().decode("utf-8"))
        cpython = f"cp{sys.version_info.major}{sys.version_info.minor}-"
        versions = []
        for raw_version, files in payload.get("releases", {}).items():
            if not re.fullmatch(r"\d+(?:\.\d+)*", raw_version):
                continue
            version_key = tuple(int(part) for part in raw_version.split("."))
            if version_key[0] >= spec["max_major"]:
                continue
            versions.append((version_key, raw_version, files))
        for _, version, files in sorted(versions, reverse=True):
            wheels = [item for item in files
                      if item.get("packagetype") == "bdist_wheel"
                      and item.get("filename", "").endswith(".whl")
                      and not item.get("yanked", False)
                      and ("py3-none-any" in item["filename"]
                           or "py3-none-win_amd64" in item["filename"]
                           or (cpython in item["filename"]
                               and "win_amd64" in item["filename"]))]
            if wheels:
                wheels.sort(key=lambda item: "win_amd64" not in item["filename"])
                wheel = wheels[0]
                return (version, wheel["url"], wheel["digests"]["sha256"],
                        wheel["filename"])
        raise RuntimeError(f"PyPI has no supported Windows wheel for {spec['label']}.")

    def install_python_engine(self, name, version, url, digest, filename):
        """Download, verify and stage one wheel without activating it."""
        ENGINE_UPDATE_ROOT.mkdir(parents=True, exist_ok=True)
        wheel_path = ENGINE_UPDATE_ROOT / (filename + ".part")
        wheel_path.unlink(missing_ok=True)
        try:
            with self._net_opener() as opener:
                request = urllib.request.Request(
                    url, headers={"User-Agent": f"{APP_NAME}/{APP_VERSION}"})
                with opener.open(request, timeout=60) as response, wheel_path.open("wb") as target:
                    shutil.copyfileobj(response, target, length=1024 * 1024)
            actual = hashlib.sha256(wheel_path.read_bytes()).hexdigest()
            if actual.lower() != digest.lower():
                raise RuntimeError("downloaded wheel failed its PyPI SHA-256 check")

            safe_version = re.sub(r"[^0-9A-Za-z_.-]", "_", version)
            target_dir = ENGINE_UPDATE_ROOT / f"{name}-{safe_version}-{digest[:12].lower()}"
            if not target_dir.is_dir():
                with tempfile.TemporaryDirectory(dir=ENGINE_UPDATE_ROOT) as temp:
                    staging = Path(temp) / "package"
                    staging.mkdir()
                    root = staging.resolve()
                    with zipfile.ZipFile(wheel_path) as wheel:
                        for member in wheel.infolist():
                            destination = (staging / member.filename).resolve()
                            destination.relative_to(root)
                        wheel.extractall(staging)
                    os.replace(staging, target_dir)
            return {"version": version, "path": str(target_dir), "sha256": digest.lower()}
        finally:
            wheel_path.unlink(missing_ok=True)

    def validate_python_engine_updates(self, records):
        """Smoke-test staged wheels together in a clean child process."""
        paths = [record["path"] for record in records.values()]
        command = [sys.executable]
        if not getattr(sys, "frozen", False):
            command.append(str(Path(__file__).resolve()))
        command += ["--validate-engine-update", json.dumps(paths)]
        result = subprocess.run(
            command, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            text=True, encoding="utf-8", errors="replace", timeout=90,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
        )
        if result.returncode != 0:
            detail = (result.stderr or result.stdout or "compatibility probe failed").strip()
            raise RuntimeError(detail[-180:])

    @staticmethod
    def commit_python_engine_updates(records):
        """Atomically activate a validated group on the next app launch."""
        try:
            registry = json.loads(ENGINE_UPDATE_REGISTRY.read_text("utf-8"))
            if not isinstance(registry, dict):
                registry = {}
        except (OSError, ValueError):
            registry = {}
        registry.update(records)
        pending = ENGINE_UPDATE_REGISTRY.with_suffix(".json.new")
        pending.write_text(json.dumps(registry, indent=2), encoding="utf-8")
        os.replace(pending, ENGINE_UPDATE_REGISTRY)

    @staticmethod
    def _is_native_executable(path, minimum_size=1024 * 1024):
        """True for a real program image for this machine, not an HTML error page.

        A mirror that answers a download with a login page or a 404 body still
        writes a file of plausible size, so the first bytes are checked against
        the executable formats: PE on Windows, ELF on Linux and Android, and
        Mach-O — thin or universal, either byte order — on macOS.
        """
        magics = (b"MZ",) if IS_WINDOWS else (
            (b"\xcf\xfa\xed\xfe", b"\xce\xfa\xed\xfe", b"\xca\xfe\xba\xbe",
             b"\xbe\xba\xfe\xca") if IS_MACOS else (b"\x7fELF",))
        try:
            path = Path(path)
            with path.open("rb") as handle:
                head = handle.read(4)
            return (path.stat().st_size >= minimum_size
                    and any(head.startswith(magic) for magic in magics))
        except OSError:
            return False

    def find_tool(self, name):
        """Locate a tool, newest copy first.

        A downloaded copy in app data only exists after an update, so it always
        wins over the copy shipped inside the bundle.
        """
        spec = TOOLS[name]
        downloaded = self.tools_dir / spec["exe"]
        if self._is_native_executable(downloaded):
            return str(downloaded)
        packaged = bundled_tool(spec["exe"])
        if packaged:
            return packaged
        found = shutil.which(name)
        if found:
            return found
        # Nothing on disk. An importable package is still a usable tool, and
        # pyshell knows how to run one; this is the whole Android path.
        if spec.get("module"):
            if IS_ANDROID and pyshell.module_available(name):
                return name
            if not IS_ANDROID and pyshell.module_version(name) != "unknown":
                return name
        return None

    @staticmethod
    def tool_download_urls(name, version=None):
        """Release URLs for a specific version, with independent mirrors."""
        if name == "spotdl" and version:
            return spotdl_urls(version.lstrip("vV"))
        if name == "yt-dlp" and version:
            return ytdlp_urls(version)
        return TOOLS[name]["urls"]

    def get_or_download_tool(self, name, force=False, version=None):
        """Return a usable tool path, downloading from the mirrors when needed."""
        spec = TOOLS[name]
        if IS_ANDROID and spec.get("module"):
            # Python command packages are compiled into the APK by Chaquopy.
            # Android has no compatible native release asset to fetch, so never
            # turn a missing/broken bundle into a doomed network download.
            existing = self.find_tool(name)
            if existing:
                return existing
            raise RuntimeError(
                f"The bundled {spec['label']} package is missing or cannot start. "
                "Install a complete BlueKnight Downloader APK.")
        if not force:
            existing = self.find_tool(name)
            if existing:
                return existing

        target = self.tools_dir / spec["exe"]
        candidate = target.with_name(spec["exe"] + ".new")
        archive_path = None
        if spec.get("archive_member"):
            archive_path = target.with_name(spec["exe"] + ".zip.new")
            self._download_with_fallbacks(
                self.tool_download_urls(name, version), archive_path, spec["label"])
            try:
                if spec.get("checksum_urls"):
                    checksum_path = target.with_name(spec["exe"] + ".sha256.new")
                    try:
                        self._download_with_fallbacks(
                            spec["checksum_urls"], checksum_path, f"{spec['label']} checksum")
                        manifest = checksum_path.read_text("ascii", "replace")
                        match = re.search(r"\b([0-9A-Fa-f]{64})\b", manifest)
                        if not match:
                            raise RuntimeError(f"{spec['label']} checksum manifest is invalid.")
                        actual = hashlib.sha256(archive_path.read_bytes()).hexdigest()
                        if actual.lower() != match.group(1).lower():
                            raise RuntimeError(f"{spec['label']} archive failed its SHA-256 check.")
                    finally:
                        checksum_path.unlink(missing_ok=True)
                with zipfile.ZipFile(archive_path) as archive:
                    member = spec["archive_member"]
                    if member not in archive.namelist():
                        raise RuntimeError(f"{spec['label']} archive does not contain {member}.")
                    with archive.open(member) as source, candidate.open("wb") as output:
                        shutil.copyfileobj(source, output, length=1024 * 1024)
            finally:
                archive_path.unlink(missing_ok=True)
        else:
            self._download_with_fallbacks(
                self.tool_download_urls(name, version), candidate, spec["label"])
        if not self._is_native_executable(candidate):
            candidate.unlink(missing_ok=True)
            raise RuntimeError(
                f"The downloaded {spec['label']} file is not a valid executable.")
        make_executable(candidate)
        candidate_version = self.tool_version(candidate)
        if candidate_version == "unknown":
            candidate.unlink(missing_ok=True)
            raise RuntimeError(f"The downloaded {spec['label']} executable did not start.")
        if version and candidate_version.lstrip("vV") != version.lstrip("vV"):
            candidate.unlink(missing_ok=True)
            raise RuntimeError(
                f"The {spec['label']} server returned {candidate_version}, expected {version}.")
        target.parent.mkdir(parents=True, exist_ok=True)
        os.replace(candidate, target)
        return str(target)

    def tool_version(self, path):
        try:
            stem = Path(path).stem.lower()
            if stem.startswith("lib") and stem[3:] in {"ffmpeg", "ffprobe"}:
                stem = stem[3:]
            version_flag = "-version" if stem in {"ffmpeg", "ffprobe"} else "--version"
            result = pyshell.run([path, version_flag], timeout=45)
            lines = result.stdout.strip().splitlines()
            if result.returncode != 0 or not lines:
                return "unknown"
            if stem == "deno":
                match = re.match(r"deno\s+([0-9][^\s]*)", lines[0])
                return match.group(1) if match else lines[0].strip()
            if stem in {"ffmpeg", "ffprobe"}:
                match = re.match(r"ff(?:mpeg|probe) version\s+([^\s]+)", lines[0], re.I)
                return match.group(1) if match else lines[0].strip()
            return lines[-1].strip()
        except Exception:
            return "unknown"

    def update_tools(self):
        """Refresh spotDL and yt-dlp from their mirrors, then report per tool."""
        if IS_ANDROID:
            self.notify(
                "Android components are verified and updated with each app release.", "ok")
            return
        if self.updating or self.is_downloading or self.is_converting or self.setup_running:
            self.notify("Wait for the current job to finish.", "err")
            return
        self.updating = True
        threading.Thread(target=self._update_tools, daemon=True).start()

    def probe_cookie(self, candidate, domain=None):
        """Read this jar and see what is in it. Returns (ok, detail). Offline.

        Two separate things can be wrong, and only one of them looks like an
        error. The jar may be unreadable — a Chromium browser holds a lock while
        it runs, and Chrome 127+ encrypts app-bound so DPAPI never opens it. Or
        the jar reads perfectly and simply holds no session for the site, which
        yt-dlp reports much later as an empty response or a bot check. Asking
        yt-dlp to export the jar answers both, with a bogus URL and no request.
        """
        if candidate[0] == "file":
            # Already a Netscape jar: read it. Handing it to yt-dlp only to have
            # it written back out would mean two --cookies flags, and the second
            # would win — which is how an exported jar came back empty.
            return self._probe_cookie_file(Path(candidate[1]), domain)

        if not self.ytdlp_cmd:
            return False, "yt-dlp is not installed yet"

        jar = Path(tempfile.gettempdir()) / f"blueknight_cookies_{os.getpid()}.txt"
        jar.unlink(missing_ok=True)
        try:
            out = pyshell.run(
                [self.ytdlp_cmd, "--simulate", "--no-warnings",
                 *ytdlp_cookie_args(candidate), "--cookies", str(jar), COOKIE_PROBE_URL],
                timeout=40)
            text = (out.stdout or "") + (out.stderr or "")
            failure = COOKIE_FAILURE.search(text)
            if failure:
                return False, failure.group(0)[:120]

            lines = []
            if jar.is_file():
                lines = [ln for ln in jar.read_text("utf-8", "replace").splitlines()
                         if ln.strip() and not ln.startswith("#")]
            if not lines:
                return False, "no cookies read"
            if domain and not any(jar_line_matches_domain(ln, domain) for ln in lines):
                return False, f"{len(lines)} cookies, none for {domain}"
            return True, (f"{len(lines)} cookies, {domain} session" if domain
                          else f"{len(lines)} cookies")
        except Exception as exc:
            return False, str(exc)[:120]
        finally:
            jar.unlink(missing_ok=True)

    @staticmethod
    def _probe_cookie_file(path, domain=None):
        """What is inside an exported cookies.txt."""
        try:
            lines = [ln for ln in path.read_text("utf-8", "replace").splitlines()
                     if ln.strip() and not ln.lstrip().startswith("#")]
        except Exception as exc:
            return False, str(exc)[:120]
        if not lines:
            return False, "the file is empty"
        if domain and not any(jar_line_matches_domain(ln, domain) for ln in lines):
            return False, f"{len(lines)} cookies, none for {domain}"
        return True, (f"{len(lines)} cookies, {domain} session" if domain
                      else f"{len(lines)} cookies")

    def usable_cookies(self, domain=None, require_session=False):
        """Cookie candidates that pass the probe, ones holding a session first.

        require_session drops jars that read fine but carry nothing for the
        site — retrying a signed-out jar against a sign-in demand is a loop.
        """
        good, readable = [], []
        self._session_cookies = set()
        for candidate in cookie_candidates(domain):
            ok, detail = self.probe_cookie(candidate, domain)
            if not ok and locked_by_browser(candidate):
                # yt-dlp calls every Chromium browser "Chrome" here, which reads
                # as nonsense when the browser in question is Opera GX.
                detail = "open right now — close it to use its cookies"
            self._batch_log(f"Cookies · {candidate_label(candidate)}: {detail}", "success" if ok else "warning")
            if ok:
                good.append(candidate)
                self._session_cookies.add(candidate)
            elif domain and detail.endswith(f"none for {domain}"):
                # Readable but signed out. Worth a try if nothing better exists.
                readable.append(candidate)
        self._flush_logs()
        return good if require_session else good + readable

    def _update_tools(self):
        self.ui("status", "Updating")
        self.ui("updating", True)
        changed, failed = [], []
        try:
            for name, spec in TOOLS.items():
                label = spec["label"]
                try:
                    path = self.find_tool(name)
                    current = self.tool_version(path) if path else None

                    if path:
                        # Never re-download what is already current: ask first,
                        # and only spend the bandwidth when there is a newer one.
                        self.log(f"Checking {label} (have {current})…", "info")
                        latest = self.latest_version(name)
                        if latest is None:
                            self.log(f"{label}: could not reach the update server. "
                                     f"Keeping {current}.", "warning")
                            self._remember_tool(name, path)
                            continue
                        if versions_equal(latest, current):
                            self.log(f"{label} is already current ({current}).", "success")
                            self._remember_tool(name, path)
                            continue
                        self.log(f"{label} {current} → {latest} available.", "info")
                    else:
                        self.log(f"{label} is missing. Installing…", "warning")

                    path = self.get_or_download_tool(
                        name, force=True, version=latest if path else None)
                    after = self.tool_version(path)
                    self._remember_tool(name, path)
                    changed.append(f"{label} {current or 'missing'} → {after}")
                    self.log(f"✓ {label} is now {after}", "success")
                except Exception as exc:
                    failed.append(f"{label}: {str(exc)[:160]}")
                    self.log(f"{label} update failed: {str(exc)[:200]}", "error")

            staged_engines = {}
            staged_labels = []
            for name, spec in PYTHON_COMPONENTS.items():
                label = spec["label"]
                try:
                    current = self.python_engine_version(spec["module"])
                    self.log(f"Checking {label} (have {current or 'missing'})…", "info")
                    latest, url, digest, filename = self.latest_python_engine(name)
                    if versions_equal(current, latest):
                        self.log(f"{label} is already current ({current}).", "success")
                        continue
                    staged_engines[name] = self.install_python_engine(
                        name, latest, url, digest, filename)
                    staged_labels.append((label, current or "missing", latest))
                    self.log(f"Staged {label} {latest}; validating the update set…", "info")
                except Exception as exc:
                    failed.append(f"{label}: {str(exc)[:160]}")
                    self.log(f"{label} update failed: {str(exc)[:200]}", "error")

            if staged_engines:
                try:
                    self.validate_python_engine_updates(staged_engines)
                    self.commit_python_engine_updates(staged_engines)
                    for label, current, latest in staged_labels:
                        changed.append(f"{label} {current} → {latest} (restart required)")
                        self.log(
                            f"✓ {label} {latest} passed compatibility checks; "
                            "restart to activate it.", "success")
                except Exception as exc:
                    failed.append(f"Python engine set: {str(exc)[:160]}")
                    self.log(
                        "Python engine compatibility check failed; keeping every current version: "
                        f"{str(exc)[:180]}", "error")

            try:
                current_ffmpeg = self.tool_version(self.ffmpeg_cmd) if self.ffmpeg_cmd else "missing"
                self.log(f"Checking FFmpeg (have {current_ffmpeg})…", "info")
                updated_ffmpeg = self.get_or_download_ffmpeg(force=True)
                updated_version = self.tool_version(updated_ffmpeg)
                self.ffmpeg_cmd = updated_ffmpeg
                if updated_version != current_ffmpeg:
                    changed.append(f"FFmpeg {current_ffmpeg} → {updated_version}")
                    self.log(f"✓ FFmpeg is now {updated_version}", "success")
                else:
                    self.log(f"FFmpeg is already current ({current_ffmpeg}).", "success")
            except Exception as exc:
                failed.append(f"FFmpeg: {str(exc)[:160]}")
                self.log(
                    f"FFmpeg update failed; keeping the working version: {str(exc)[:180]}",
                    "error")

            # Instagram needs a login, so a readable cookie jar is a component
            # like any other. Checked here because yt-dlp must exist first.
            self.log("Checking cookies for Instagram…", "info")
            working = self.usable_cookies(COOKIE_DOMAINS["instagram"])
            if working and working[0] in self._session_cookies:
                self.log(f"Instagram session found in {candidate_label(working[0])}.", "success")
            elif working:
                self.log(f"{candidate_label(working[0])} cookies are readable but hold no Instagram session — "
                         "log into Instagram there, or export a cookies.txt.", "warning")
            elif cookie_candidates(COOKIE_DOMAINS['instagram']):
                self.log("No cookie jar could be read. Close Edge and Chrome, use Firefox, or "
                         "drop a cookies.txt beside the app.", "warning")
            else:
                self.log("No cookie source found. Instagram needs a cookies.txt beside the app, "
                         "in your Downloads, or in the download folder.", "warning")

            self.log("Checking the saved SoundCloud session…", "info")
            if registry_entry("soundcloud"):
                soundcloud_session = self.usable_cookies(
                    COOKIE_DOMAINS["soundcloud"], require_session=True)
                if soundcloud_session:
                    self.log(
                        f"SoundCloud Go/Go+ session found in "
                        f"{candidate_label(soundcloud_session[0])}.", "success")
                else:
                    self.log(
                        "The saved SoundCloud session is expired or unreadable; "
                        "sign in again on the SoundCloud page.", "warning")
            else:
                self.log(
                    "SoundCloud is in public mode. Go/Go+ members can optionally sign in "
                    "on the SoundCloud page.", "info")

            self.ui("components", self.component_summary())
            if failed:
                self.notify(
                    "Some updates failed:\n" + "\n".join(failed)
                    + "\n\nThe versions you already have still work. On a restricted network, "
                      "turn on a proxy in Connection settings and try again.",
                    "err",
                )
            elif changed:
                self.notify("Updated:\n" + "\n".join(changed), "ok")
            else:
                self.notify("Everything is already up to date.", "ok")
        finally:
            self.updating = False
            self.ui("updating", False)
            self.status.set("Ready")
            self.ui("status", "Ready")

    def _remember_tool(self, name, path):
        if name == "spotdl":
            self.spotdl_cmd = path
        elif name == "yt-dlp":
            self.ytdlp_cmd = path
        elif name == "deno":
            self.deno_cmd = path

    def component_summary(self):
        python_versions = {
            name: self.python_engine_version(spec["module"]) or "bundled"
            for name, spec in PYTHON_COMPONENTS.items()
        }
        pdf_label = PYTHON_COMPONENTS[PDF_LIBRARY]["label"]
        desktop_converters = (
            f"\nOffice converter: {self._converter_program('soffice') or 'LibreOffice not found'}"
            f"\nEbook converter: {self._converter_program('ebook-convert') or 'Calibre not found'}"
            if not IS_ANDROID else "")
        return (f"Spotify engine: {self.spotdl_cmd}\n"
                f"yt-dlp: {self.ytdlp_cmd}\n"
                f"Deno/EJS: {self.deno_cmd or 'not used on Android'}\n"
                f"gallery-dl: {python_versions['gallery-dl']}\n"
                f"Streamlink: {python_versions['streamlink']}\n"
                f"{'reportlab' if IS_ANDROID else 'img2pdf'}: "
                f"{python_versions['reportlab' if IS_ANDROID else 'img2pdf']}\n"
                f"Pillow: {python_versions['pillow']}\n"
                f"{pdf_label}: {python_versions[PDF_LIBRARY]}\n"
                f"certifi: {python_versions['certifi']}\n"
                f"FFmpeg: {self.tool_version(self.ffmpeg_cmd) if self.ffmpeg_cmd else 'missing'}"
                f"{desktop_converters}")

    def get_or_download_ffmpeg(self, force=False):
        """Return FFmpeg, staging and launching updates before activation."""
        if not force:
            # A successfully updated copy must win over the bundled baseline.
            if (self._is_native_executable(self.ffmpeg_exe)
                    and self.tool_version(self.ffmpeg_exe) != "unknown"):
                return str(self.ffmpeg_exe)
            packaged_ffmpeg = bundled_tool(f"ffmpeg{EXE}")
            if packaged_ffmpeg:
                return packaged_ffmpeg
            system_ffmpeg = shutil.which("ffmpeg")
            if system_ffmpeg:
                return system_ffmpeg
            try:
                import imageio_ffmpeg
                binary = imageio_ffmpeg.get_ffmpeg_exe()
                if binary and Path(binary).exists():
                    return binary
            except Exception:
                pass

        archives = FFMPEG_DOWNLOADS[OS_TAG]
        if not archives:
            raise RuntimeError(
                "FFmpeg is not downloadable on this platform and none is bundled.")

        archive_path = self.app_data / "ffmpeg.archive.new"
        try:
            self.app_data.mkdir(parents=True, exist_ok=True)
            self.ffmpeg_exe.parent.mkdir(parents=True, exist_ok=True)
            with tempfile.TemporaryDirectory(dir=self.ffmpeg_dir) as temp:
                stage = Path(temp)
                # Both binaries are staged and launch-checked before either is
                # activated, so a half-downloaded pair can never replace a
                # working one.
                for wanted, urls in archives:
                    archive_path.unlink(missing_ok=True)
                    self._download_with_fallbacks(urls, archive_path, "FFmpeg")
                    self.ui("setup", "Extracting FFmpeg...", "Setting up audio converter")
                    extract_binaries(archive_path,
                                     {name: stage / f"{name}{EXE}" for name in wanted})
                for name in ("ffmpeg", "ffprobe"):
                    candidate = stage / f"{name}{EXE}"
                    make_executable(candidate)
                    if (not self._is_native_executable(candidate)
                            or self.tool_version(candidate) == "unknown"):
                        raise RuntimeError(
                            f"The downloaded {candidate.name} did not pass its launch check")

                # ffmpeg is the activation point and is replaced last. A failure
                # before this line leaves every running download and the prior
                # executable untouched.
                os.replace(stage / f"ffprobe{EXE}", self.ffmpeg_dir / "bin" / f"ffprobe{EXE}")
                os.replace(stage / f"ffmpeg{EXE}", self.ffmpeg_exe)
            return str(self.ffmpeg_exe)
        except Exception as e:
            raise RuntimeError(f"FFmpeg setup failed: {e}")
        finally:
            archive_path.unlink(missing_ok=True)

    def ffmpeg_location(self):
        """What --ffmpeg-location can actually be handed.

        yt-dlp looks a program up by file name inside the folder it is given.
        On Android the packaged binaries are named libffmpeg.so and
        libffprobe.so — a packaging rule, not a format — so the folder only
        answers once the shims exist there under the real names. Passing the
        native library folder regardless, which is what this did, meant yt-dlp
        found no ffmpeg at all: every merge of a separate video and audio
        stream, and every MP3 extraction, would have failed with the tool
        reported missing.

        When the shims cannot be made, the binary's own path goes instead.
        yt-dlp accepts a file there and picks the program out of its name, and
        "libffmpeg.so" contains "ffmpeg".
        """
        ffmpeg = Path(self.ffmpeg_cmd)
        if IS_ANDROID:
            # Resolving ffprobe is what puts its shim beside ffmpeg's, and it
            # has to happen here rather than at start-up: a folder holding one
            # of the pair is a folder yt-dlp cannot probe formats with, and
            # nothing else guarantees the order. android_shim is idempotent, so
            # asking again costs a stat.
            with contextlib.suppress(Exception):
                bundled_tool(f"ffprobe{EXE}")
        folder = ffmpeg.parent
        if (folder / f"ffmpeg{EXE}").is_file():
            return str(folder)
        return str(ffmpeg)

    def js_runtime(self):
        """The JavaScript engine yt-dlp gets, as (name, path).

        This is not optional for YouTube. Its media URLs are signed by a
        JavaScript challenge, and without an engine to run it yt-dlp falls back
        to its own interpreter, which current player versions defeat — the
        unsigned URL then comes back as HTTP 403 Forbidden. That is exactly the
        Android-only failure, because Deno publishes no Android build and the
        flag was simply never passed there.

        QuickJS is the answer on the phone: one small C interpreter that
        cross-compiles with the NDK, and one yt-dlp supports by name.
        """
        if self.deno_cmd and Path(self.deno_cmd).is_file():
            return "deno", self.deno_cmd
        node = shutil.which("node")
        if node:
            return "node", node
        quickjs = (bundled_tool(f"quickjs{EXE}") or bundled_tool(f"qjs{EXE}")
                   or shutil.which("qjs"))
        if quickjs:
            return "quickjs", quickjs
        return None, None

    def find_spotdl(self):
        return self.find_tool("spotdl")

    # ------------------------------------------------------------------
    # OPTIMIZED DOWNLOAD - With HTTP and SOCKS5 proxy support
    # ------------------------------------------------------------------
    def start_download(self, url, fmt=None, quality=None):
        if fmt:
            self.download_format.set(fmt)
        if quality:
            self.quality.set(quality)

        if self.updating or self.setup_running or not self.spotdl_cmd:
            self.notify("Still finishing setup. Try again in a moment.", "err")
            self.ui("finished")
            return
        if self.is_downloading:
            return

        url = (url or "").strip()
        if not self._re_url.search(url):
            self.notify("That is not a Spotify track, album, or playlist link.", "err")
            self.ui("finished")
            return

        try:
            output = download_dir("spotify", self.save_folder.get().strip())
        except Exception as exc:
            self.notify(f"Could not create the download folder.\n\n{exc}", "err")
            self.ui("finished")
            return

        self.is_downloading = True
        self.status.set("Downloading")
        options = {
            "format": self.download_format.get(),
            "quality": self.quality.get(),
            "proxy_enabled": self.use_proxy.get(),
            "proxy_url": self.proxy_url.get().strip(),
            "proxy_type": self.proxy_type.get(),
        }
        threading.Thread(
            target=(self.download_with_android_spotify
                    if IS_ANDROID else self.download_with_spotdl),
            args=(url, output, options),
            daemon=True,
        ).start()

    def _spotify_json(self, url, token=None):
        headers = {"User-Agent": BROWSER_UA, "Accept": "application/json"}
        if token:
            headers["Authorization"] = f"Bearer {token}"
        with self._net_opener() as opener:
            request = urllib.request.Request(url, headers=headers)
            with opener.open(request, timeout=30) as response:
                return json.loads(response.read().decode("utf-8", "replace"))

    @staticmethod
    def _spotify_track_record(item, album=None):
        return {
            "title": str(item.get("name") or "Unknown track"),
            "artists": ", ".join(
                str(artist.get("name")) for artist in item.get("artists", [])
                if artist.get("name")) or "Unknown artist",
            "album": str((item.get("album") or {}).get("name") or album or ""),
            "track_number": int(item.get("track_number") or 0),
        }

    def android_spotify_tracks(self, url):
        """Resolve a public Spotify track, album or playlist without spotDL.

        Spotify's web player exposes a short-lived anonymous access token. The
        token is used only for public metadata; audio still comes from the
        matching YouTube Music result, exactly as it does in spotDL.
        """
        match = re.search(
            r"open\.spotify\.com/(track|album|playlist)/([A-Za-z0-9]+)", url, re.I)
        if not match:
            raise ValueError("That is not a Spotify track, album, or playlist link.")
        kind, item_id = match.group(1).lower(), match.group(2)
        token_payload = self._spotify_json(
            "https://open.spotify.com/get_access_token?reason=transport&productType=web_player")
        token = token_payload.get("accessToken")
        if not token:
            raise RuntimeError("Spotify did not provide a public metadata session.")

        base = "https://api.spotify.com/v1"
        if kind == "track":
            return [self._spotify_track_record(
                self._spotify_json(f"{base}/tracks/{item_id}", token))]

        payload = self._spotify_json(f"{base}/{kind}s/{item_id}", token)
        album = payload.get("name") if kind == "album" else None
        page = payload.get("tracks") or {}
        tracks = []
        while page:
            for wrapped in page.get("items", []):
                item = wrapped.get("track") if kind == "playlist" else wrapped
                if item and not item.get("is_local"):
                    tracks.append(self._spotify_track_record(item, album))
            next_url = page.get("next")
            page = self._spotify_json(next_url, token) if next_url else None
        if not tracks:
            raise RuntimeError(f"Spotify returned no downloadable tracks for this {kind}.")
        return tracks

    def download_with_android_spotify(self, url, output, options):
        """Android's pure-Python Spotify-to-YouTube-Music download path."""
        started = time.time()
        failures = []
        completed = 0
        try:
            from ytmusicapi import YTMusic
            import mutagen
            import yt_dlp

            tracks = self.android_spotify_tracks(url)
            total = len(tracks)
            self.ui("plan", total, 1, str(output))
            music = YTMusic()

            for index, track in enumerate(tracks, 1):
                if not self.is_downloading:
                    break
                label = f"{track['artists']} - {track['title']}"
                self.ui("track", index, total, label[:90])
                self._batch_log(f"[{index}/{total}] Finding {label}", "download")
                try:
                    results = music.search(label, filter="songs", limit=5)
                    match = next((item for item in results if item.get("videoId")), None)
                    if not match:
                        raise RuntimeError("no matching YouTube Music result")
                    video_url = "https://music.youtube.com/watch?v=" + match["videoId"]
                    target = unique_path(output, safe_filename(
                        label, fallback=f"track-{index}", suffix=f".{options['format']}"))
                    outtmpl = str(target.parent / f"{target.stem}.%(ext)s")

                    def progress(status):
                        if not self.is_downloading:
                            raise RuntimeError("Download stopped")
                        if status.get("status") == "downloading":
                            match_percent = re.search(
                                r"([0-9]+(?:\.[0-9]+)?)", status.get("_percent_str", ""))
                            percent = float(match_percent.group(1)) if match_percent else 0.0
                            self.ui("progress", percent, status.get("_speed_str", ""),
                                    status.get("_eta_str", ""))

                    ydl_options = {
                        "format": "bestaudio/best",
                        "outtmpl": outtmpl,
                        "noplaylist": True,
                        "quiet": True,
                        "no_warnings": True,
                        # Same rule as the command-line path: the folder when
                        # both shims are there, the binary itself when not.
                        "ffmpeg_location": self.ffmpeg_location(),
                        "progress_hooks": [progress],
                        "postprocessors": [{
                            "key": "FFmpegExtractAudio",
                            "preferredcodec": options["format"],
                            "preferredquality": options["quality"],
                        }],
                    }
                    if options["proxy_enabled"] and options["proxy_url"]:
                        ydl_options["proxy"] = normalize_proxy_url(
                            options["proxy_url"], options["proxy_type"])
                    with yt_dlp.YoutubeDL(ydl_options) as downloader:
                        downloader.download([video_url])

                    if not target.is_file():
                        candidates = sorted(target.parent.glob(target.stem + ".*"))
                        target = candidates[-1] if candidates else target
                    if not target.is_file() or target.stat().st_size == 0:
                        raise RuntimeError("the audio converter returned no file")
                    with contextlib.suppress(Exception):
                        tags = mutagen.File(str(target), easy=True)
                        if tags is not None:
                            tags["title"] = [track["title"]]
                            tags["artist"] = [track["artists"]]
                            if track["album"]:
                                tags["album"] = [track["album"]]
                            if track["track_number"]:
                                tags["tracknumber"] = [str(track["track_number"])]
                            tags.save()
                    completed += 1
                    self._batch_log(f"Downloaded {target.name}", "success")
                except Exception as exc:
                    failures.append(f"{label}: {str(exc)[:140]}")
                    self._batch_log(f"Skipped {label}: {str(exc)[:140]}", "error")

            self._flush_logs()
            if not self.is_downloading:
                self.log("Download stopped.", "warning")
            elif completed:
                self.ui("success", max(time.time() - started, 0.1))
                if failures:
                    self.notify(
                        f"Downloaded {completed} track(s); {len(failures)} could not be matched.",
                        "err")
            else:
                raise RuntimeError(failures[0] if failures else "No tracks were downloaded.")
        except Exception as exc:
            self._flush_logs()
            detail = str(exc)[:300]
            self.log(f"Spotify download failed: {detail}", "error")
            self.ui("failure", detail)
        finally:
            self.ui("finished")

    def _prepare_proxy(self, proxy_url, proxy_type):
        """Return an HTTP proxy URL accepted by spotDL and an optional bridge."""
        proxy = normalize_proxy_url(proxy_url, proxy_type)
        if urlsplit(proxy).scheme.startswith("socks5"):
            bridge = Socks5HttpBridge(proxy)
            local_proxy = bridge.start()
            self._batch_log(f"SOCKS5 bridge listening at {local_proxy}", "info")
            return local_proxy, bridge
        return proxy, None

    def download_with_spotdl(self, url, output, options):
        """Optimized download with proxy support"""
        fmt = options["format"]
        output_template = output / "{artists} - {title}.{output-ext}"
        cmd = [self.spotdl_cmd, "download", url, "--output", str(output_template)]

        # Format settings
        if fmt == "mp3":
            cmd += ["--format", "mp3", "--bitrate", f"{options['quality']}k"]
        elif fmt == "flac":
            cmd += ["--format", "flac"]

        proxy_enabled = options["proxy_enabled"]
        proxy_url = options["proxy_url"]
        proxy_type = options["proxy_type"]
        bridge = None
        env = os.environ.copy()

        if proxy_enabled and proxy_url:
            try:
                effective_proxy, bridge = self._prepare_proxy(proxy_url, proxy_type)
            except Exception as exc:
                self.log(f"Proxy setup failed: {exc}", "error")
                self.ui("failure", f"Proxy setup failed:\n\n{exc}")
                self.ui("finished")
                return
            cmd += ["--proxy", effective_proxy]
            for name in ("HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy"):
                env[name] = effective_proxy
            self._batch_log(
                f"Using proxy: {display_proxy_url(normalize_proxy_url(proxy_url, proxy_type))}",
                "info",
            )
        else:
            self._batch_log("No proxy configured - using direct connection", "info")

        # FFmpeg setup (minimal)
        if self.ffmpeg_cmd:
            ffmpeg_dir = str(Path(self.ffmpeg_cmd).parent)
            if ffmpeg_dir not in env.get("PATH", ""):
                env["PATH"] = ffmpeg_dir + os.pathsep + env.get("PATH", "")
            env["FFMPEG_PATH"] = self.ffmpeg_cmd
            cmd += ["--ffmpeg", self.ffmpeg_cmd]

        self._batch_log(f"Starting {fmt.upper()} download...", "download")
        # Until spotDL reports how many songs the link holds, say one.
        self.ui("plan", 1, 1, str(output))

        try:
            # Use larger buffer for better I/O performance
            self.process = pyshell.popen(cmd, bufsize=8192, env=env)
            
            total = None
            count = 0
            started = time.time()
            last_progress_update = 0
            error_lines = []
            
            for raw in self.process.stdout:
                if not self.is_downloading:
                    break
                    
                line = raw.strip()
                if not line:
                    continue

                # Check for important lines only
                if self._re_found_songs.search(line):
                    total = int(self._re_found_songs.search(line).group(1))
                    # The real count, replacing whatever the link implied.
                    self.ui("plan", total, 1, str(output))
                    self._batch_log(f"Found {total} song(s)", "success")
                    continue

                if self._re_downloading.search(line):
                    count += 1
                    track = line.split(":", 1)[-1].strip() if ":" in line else line
                    track = self._re_brackets.sub("", track).strip()
                    
                    current_time = time.time()
                    if current_time - last_progress_update > 0.5:
                        self.ui("track", count, total, track[:90])
                        last_progress_update = current_time
                    
                    if count % 5 == 0:
                        self._batch_log(f"[{count}/{total or '?'}] {track[:80]}", "download")
                    continue

                if self._re_downloaded.search(line):
                    self._batch_log("✓ Downloaded", "success")
                    continue

                if self._re_error.search(line) and len(line) < 200:
                    error_lines.append(line)
                    self._batch_log(line[:200], "error")
                elif self._re_request_error.search(line):
                    error_lines.append(line)
                    self._batch_log(line[:200], "error")
                elif "Invalid proxy" in line or "proxy server" in line.lower():
                    error_lines.append(line)
                    self._batch_log(f"PROXY ERROR: {line[:200]}", "error")

            self.process.wait()
            elapsed = max(time.time() - started, 0.1)
            code = self.process.returncode

            if not self.is_downloading:
                self._flush_logs()
                self.log("Download stopped.", "warning")
                return

            if code == 0:
                self._flush_logs()
                self.ui("success", elapsed)
            else:
                error_text = "\n".join(error_lines[-10:]) if error_lines else "Unknown error"
                self._flush_logs()
                
                # Enhanced proxy troubleshooting
                if "Invalid proxy" in error_text or "proxy server" in error_text.lower():
                    troubleshooting = (
                        "\n\n🔧 PROXY ISSUE DETECTED:"
                        "\n• Confirm that the selected HTTP/SOCKS5 type matches the local port"
                        "\n• v2rayN commonly uses HTTP 10809 and SOCKS5 10808"
                        "\n• Test the proxy from Advanced settings"
                        "\n• Make sure the local proxy application is running"
                    )
                elif proxy_enabled and self._re_request_error.search(error_text):
                    troubleshooting = (
                        "\n\n🔧 PROXY CONNECTION ISSUE:"
                        "\n• Verify v2rayN is running"
                        "\n• Check the proxy port (HTTP: 10809, SOCKS5: 10808)"
                        "\n• Try 'Test Proxy' button"
                        "\n• Ensure your VPN server is working"
                        "\n• Try a different proxy server"
                    )
                else:
                    troubleshooting = self._get_troubleshooting(error_text, proxy_enabled)

                self.log(f"spotDL exited with code {code}", "error")
                self.log(f"Error: {error_text[:300]}", "error")
                self.log(troubleshooting, "warning")
                
                self.ui("failure", f"Exit code {code}\n\n{error_text[:300]}{troubleshooting}")
        except FileNotFoundError:
            self._flush_logs()
            self.log("spotDL not found. Try 'Update spotDL'.", "error")
            self.ui("failure", "spotDL could not be launched.")
        except Exception as exc:
            self._flush_logs()
            self.log(f"Download error: {str(exc)[:300]}", "error")
            self.ui("failure", f"Download failed: {str(exc)[:300]}")
        finally:
            if bridge:
                bridge.close()
            self.ui("finished")

    # ------------------------------------------------------------------
    # Video and audio sites, through yt-dlp
    # ------------------------------------------------------------------
    def start_media_download(self, kind, url, quality="best", audio_only=False, limit=None):
        if self.is_downloading:
            return
        if self.updating or not self.ytdlp_cmd or not self.ffmpeg_cmd:
            self.notify("Still finishing setup. Try again in a moment.", "err")
            self.ui("finished")
            return

        try:
            url = normalize_media_url(kind, url)
            token = self.youtube_po_token.get().strip()
            if kind in YOUTUBE_KINDS and token and not PO_TOKEN_PATTERN.fullmatch(token):
                raise ValueError("The YouTube PO token must use mweb.gvs+TOKEN format.")
            if limit is not None:
                limit = int(limit)
                if not 1 <= limit <= 10000:
                    raise ValueError("Item limit must be between 1 and 10000.")
        except (TypeError, ValueError) as exc:
            self.notify(str(exc), "err")
            self.ui("finished")
            return

        try:
            output = download_dir(kind, self.save_folder.get().strip())
        except Exception as exc:
            self.notify(f"Could not create the download folder.\n\n{exc}", "err")
            self.ui("finished")
            return

        self.is_downloading = True
        self.status.set("Downloading")
        threading.Thread(
            target=self.download_with_ytdlp,
            args=(kind, url, output, quality, audio_only, limit),
            daemon=True,
        ).start()

    # ------------------------------------------------------------------
    # Documents. No yt-dlp here: PDFs and ebooks are files, not streams.
    # ------------------------------------------------------------------
    def start_pdf_download(self, url, limit=None):
        if self.is_downloading:
            return
        if self.updating or self.setup_running:
            self.notify("Components are being updated. Try again when the update finishes.", "err")
            self.ui("finished")
            return
        try:
            url = normalize_page_url(url)
            if limit is not None:
                limit = int(limit)
                if not 1 <= limit <= PDF_SCAN_LIMIT:
                    raise ValueError(f"Document limit must be between 1 and {PDF_SCAN_LIMIT}.")
        except (TypeError, ValueError) as exc:
            self.notify(str(exc), "err")
            self.ui("finished")
            return

        try:
            output = download_dir("pdf", self.save_folder.get().strip())
        except Exception as exc:
            self.notify(f"Could not create the download folder.\n\n{exc}", "err")
            self.ui("finished")
            return

        self.is_downloading = True
        self.status.set("Downloading")
        threading.Thread(target=self.download_pdfs, args=(url, output, limit),
                         daemon=True).start()

    def start_manga_download(self, url, limit=None):
        if self.is_downloading:
            return
        if self.updating or self.setup_running:
            self.notify("Components are being updated. Try again when the update finishes.", "err")
            self.ui("finished")
            return
        try:
            url = normalize_page_url(url)
            if limit is not None:
                limit = int(limit)
                if not 1 <= limit <= MANGA_PAGE_LIMIT:
                    raise ValueError(f"Page limit must be between 1 and {MANGA_PAGE_LIMIT}.")
        except (TypeError, ValueError) as exc:
            self.notify(str(exc), "err")
            self.ui("finished")
            return
        try:
            output = download_dir("pdf", self.save_folder.get().strip())
        except Exception as exc:
            self.notify(f"Could not create the download folder.\n\n{exc}", "err")
            self.ui("finished")
            return
        self.is_downloading = True
        self.status.set("Downloading")
        threading.Thread(target=self.download_manga, args=(url, output, limit),
                         daemon=True).start()

    def _open_url(self, opener, url, timeout=45, headers=None):
        request_headers = {"User-Agent": BROWSER_UA, "Accept": "*/*"}
        request_headers.update(headers or {})
        return opener.open(urllib.request.Request(url, headers=request_headers), timeout=timeout)

    def _save_stream(self, response, output, url, index, total,
                     expect=None, suffix=".pdf", name_hint=None):
        """Stream one response to disk. Returns the path, or None if it is not wanted.

        `expect` is the magic number the body must start with. For PDFs that check
        is the whole point — servers mislabel downloads constantly, and a
        200-with-an-HTML-error-page saved as .pdf is the classic way a
        "successful" scrape produces 40 broken files. Media has no single magic
        number, so there the guard is only that it must not be a web page.
        """
        head = response.read(max(len(expect), 512) if expect else 512)
        if expect:
            if not head.startswith(expect):
                return None
        else:
            content_type = (response.headers.get_content_type() or "").lower()
            hinted = Path(name_hint or urlsplit(url).path).suffix.lower()
            looks_like_media = (content_type.startswith(("video/", "audio/", "image/"))
                                or content_type in MEDIA_CONTENT_TYPES
                                or hinted in MEDIA_EXTENSIONS)
            looks_like_text = (content_type.startswith("text/")
                               or content_type in {"application/json", "application/xml",
                                                   "application/xhtml+xml"}
                               or head.lstrip().lower().startswith((b"<html", b"<!doctype")))
            if not looks_like_media or looks_like_text:
                return None

        disposition = response.headers.get("Content-Disposition", "")
        match = re.search(r'filename\*?=(?:UTF-8\'\')?"?([^";]+)', disposition, re.I)
        name = (match.group(1) if match
                else name_hint or urlsplit(url).path or "download")
        if suffix == ".bin":
            suffix = {
                "video/mp4": ".mp4", "video/webm": ".webm", "video/mp2t": ".ts",
                "audio/mpeg": ".mp3", "audio/mp4": ".m4a", "audio/ogg": ".ogg",
                "image/jpeg": ".jpg", "image/png": ".png", "image/gif": ".gif",
                "image/webp": ".webp",
            }.get(response.headers.get_content_type(), suffix)
        path = unique_path(output, safe_filename(name, suffix=suffix))
        partial = path.with_suffix(path.suffix + ".part")

        try:
            size = int(response.headers.get("Content-Length") or 0)
        except ValueError:
            size = 0
        done = len(head)
        self.ui("file", path.name, 0)
        with partial.open("wb") as out:
            out.write(head)
            while self.is_downloading:
                chunk = response.read(256 * 1024)
                if not chunk:
                    break
                out.write(chunk)
                done += len(chunk)
                now = time.time()
                if now - self.last_ui_update >= UI_UPDATE_INTERVAL:
                    self.last_ui_update = now
                    pct = (done / size * 100) if size else 0
                    self.ui("bytes", pct, path.name, f"{done / 1048576:.1f}MiB", "", 0)
        if not self.is_downloading:
            partial.unlink(missing_ok=True)
            return None
        os.replace(partial, path)
        self.ui("track", index, total, path.name[:90])
        return path

    def _save_document_response(self, response, output, url, index, total):
        """Save one validated PDF or ebook response. Returns (path, probe bytes)."""
        head = response.read(4096)
        disposition = response.headers.get("Content-Disposition", "")
        match = re.search(r'filename\*?=(?:UTF-8\'\')?"?([^";]+)', disposition, re.I)
        name = match.group(1) if match else urlsplit(url).path or "document"
        suffix = detect_document_extension(
            head, name or url, response.headers.get("Content-Type", ""))
        if not suffix:
            return None, head

        path = unique_path(output, safe_filename(name, suffix=suffix))
        partial = path.with_suffix(path.suffix + ".part")
        try:
            size = int(response.headers.get("Content-Length") or 0)
        except ValueError:
            size = 0
        done = len(head)
        self.ui("file", path.name, 0)
        try:
            with partial.open("wb") as out:
                out.write(head)
                while self.is_downloading:
                    chunk = response.read(256 * 1024)
                    if not chunk:
                        break
                    out.write(chunk)
                    done += len(chunk)
                    now = time.time()
                    if now - self.last_ui_update >= UI_UPDATE_INTERVAL:
                        self.last_ui_update = now
                        pct = (done / size * 100) if size else 0
                        self.ui("bytes", pct, path.name, f"{done / 1048576:.1f}MiB", "", 0)
            if not self.is_downloading:
                partial.unlink(missing_ok=True)
                return None, head
            if size and done != size:
                raise RuntimeError(
                    f"The server closed the document early ({done} of {size} bytes).")
            if suffix == ".zip":
                actual_suffix = inspect_zip_document(partial)
                if not actual_suffix:
                    partial.unlink(missing_ok=True)
                    return None, head
                path = unique_path(output, safe_filename(name, suffix=actual_suffix))
            elif suffix == ".pdf" and not is_valid_pdf_file(partial):
                raise RuntimeError("The server returned an incomplete or invalid PDF.")
            os.replace(partial, path)
            self.ui("track", index, total, path.name[:90])
            return path, head
        except Exception:
            partial.unlink(missing_ok=True)
            raise

    def _resolve_html5_media(self, url):
        """Resolve media exposed directly in HTML markup or social meta tags."""
        with self._net_opener() as opener:
            with self._open_url(opener, url, headers={"Accept": "text/html,*/*;q=0.8"}) as response:
                content_type = response.headers.get_content_type()
                if content_type and content_type != "text/html":
                    return []
                final_url = response.geturl()
                body = response.read(8 * 1024 * 1024).decode(
                    response.headers.get_content_charset() or "utf-8", "replace")
        parser = HtmlMediaParser()
        parser.feed(body)
        found, seen = [], set()
        for value in parser.links:
            media_url = urljoin(final_url, value)
            parsed = urlsplit(media_url)
            if parsed.scheme not in {"http", "https"} or media_url in seen:
                continue
            seen.add(media_url)
            ext = Path(parsed.path).suffix.lower()
            found.append((media_url, Path(parsed.path).name, ext.lstrip(".") or "bin"))
        return found

    def _resolve_html_images(self, url):
        """Resolve image URLs from a conventional chapter page."""
        with self._net_opener() as opener:
            with self._open_url(opener, url, headers={"Accept": "text/html,*/*;q=0.8"}) as response:
                final_url = response.geturl()
                body = response.read(12 * 1024 * 1024).decode(
                    response.headers.get_content_charset() or "utf-8", "replace")
        parser = HtmlImageParser()
        parser.feed(body)
        found, seen = [], set()
        for value in parser.links:
            image_url = urljoin(final_url, value)
            parsed = urlsplit(image_url)
            if parsed.scheme not in {"http", "https"} or image_url in seen:
                continue
            seen.add(image_url)
            ext = Path(parsed.path).suffix.lower().lstrip(".") or "bin"
            found.append((image_url, Path(parsed.path).name, ext))
        return found

    @staticmethod
    def _pick_streamlink_stream(streams, quality):
        if not streams:
            return None, None
        if quality == "best":
            return "best", streams.get("best") or next(iter(streams.values()))
        cap = int(quality) if str(quality).isdigit() else None
        choices = []
        if cap:
            for name, stream in streams.items():
                match = re.search(r"(\d{3,4})p", name)
                if match and int(match.group(1)) <= cap:
                    choices.append((int(match.group(1)), name, stream))
        if choices:
            _, name, stream = max(choices, key=lambda item: item[0])
            return name, stream
        return "best", streams.get("best") or next(iter(streams.values()))

    def _download_with_streamlink(self, url, output, quality, audio_only):
        """Record a Streamlink-supported live/VOD stream and return its path."""
        try:
            from streamlink import Streamlink
        except ImportError as exc:
            raise RuntimeError("Streamlink is not installed in this build.") from exc

        session = Streamlink()
        session.set_option("http-timeout", 30.0)
        session.set_option("stream-timeout", 20.0)
        session.set_option("http-headers", {"User-Agent": BROWSER_UA})
        if self.use_proxy.get() and self.proxy_url.get().strip():
            session.set_option(
                "http-proxy", normalize_proxy_url(self.proxy_url.get(), self.proxy_type.get()))
        streams = session.streams(url)
        stream_name, stream = self._pick_streamlink_stream(streams, quality)
        if stream is None:
            return None

        stamp = time.strftime("%Y%m%d_%H%M%S")
        host = (urlsplit(url).hostname or "stream").replace(".", "_")
        raw = unique_path(output, safe_filename(f"{host}_{stamp}", suffix=".ts"))
        partial = raw.with_suffix(raw.suffix + ".part")
        self.ui("plan", 1, 1, str(output))
        self.ui("file", raw.name, 0)
        self._batch_log(f"Streamlink opened {stream_name}. Stop saves the recording.", "info")
        self._flush_logs()
        done = 0
        handle = stream.open()
        try:
            with partial.open("wb") as target:
                while self.is_downloading:
                    chunk = handle.read(256 * 1024)
                    if not chunk:
                        break
                    target.write(chunk)
                    done += len(chunk)
                    now = time.time()
                    if now - self.last_ui_update >= UI_UPDATE_INTERVAL:
                        self.last_ui_update = now
                        self.ui("bytes", 0, raw.name, f"{done / 1048576:.1f}MiB", "", 0)
        finally:
            handle.close()
        if not done:
            partial.unlink(missing_ok=True)
            return None
        os.replace(partial, raw)

        codec = self.download_format.get() if self.download_format.get() in {"mp3", "flac"} else "mp3"
        target = raw.with_suffix(f".{codec}" if audio_only else ".mp4")
        command = ([self.ffmpeg_cmd, "-y", "-i", str(raw), "-vn", "-c:a",
                    "libmp3lame" if codec == "mp3" else "flac", str(target)] if audio_only else
                   [self.ffmpeg_cmd, "-y", "-i", str(raw), "-c", "copy",
                    "-movflags", "+faststart", str(target)])
        result = subprocess.run(
            command, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
        if result.returncode == 0 and target.is_file() and target.stat().st_size:
            raw.unlink(missing_ok=True)
            return target
        return raw

    # ------------------------------------------------------------------
    # Engines other than yt-dlp
    # ------------------------------------------------------------------
    def _resolve_with_gallery_dl(self, url, kind=None):
        """Media URLs behind a page, via gallery-dl. Returns [(url, name, ext)].

        gallery-dl is used only to *find* the media, not to fetch it. Its own
        downloader would bring a second set of proxy settings, a second cookie
        story and a second progress format; resolving here and streaming through
        the app's own fetcher keeps one of each.
        """
        try:
            import gallery_dl.config
            import gallery_dl.job
        except ImportError:
            raise RuntimeError("gallery-dl is not installed in this build.")

        gallery_dl.config.clear()
        gallery_dl.config.set((), "verbosity", 0)
        # X serves guests through the syndication endpoint; without this, photo
        # posts need a login that public posts should not require.
        gallery_dl.config.set(("extractor", "twitter"), "syndication", True)
        # The same session yt-dlp would have used, in gallery-dl's spelling: a
        # path for an exported jar, a [browser, profile] list for a live profile.
        cookie_domain = COOKIE_DOMAINS.get(kind or "")
        if cookie_domain:
            for candidate in self.usable_cookies(cookie_domain):
                if candidate[0] == "file":
                    gallery_dl.config.set(("extractor",), "cookies", candidate[1])
                elif len(candidate) > 2:
                    gallery_dl.config.set(("extractor",), "cookies",
                                          [candidate[1], candidate[2]])
                else:
                    continue
                break

        # DataJob dumps its findings as JSON to a file handle it binds from
        # sys.stdout at import time — so redirect_stdout cannot catch it, and in
        # a windowed build sys.stdout is None and writing to it would crash.
        # Handing it a buffer of its own is the only thing that holds in both.
        job = gallery_dl.job.DataJob(url, file=io.StringIO())
        job.run()
        if job.exception:
            raise job.exception

        found = []
        for entry in job.data:
            # (3, url, metadata) is gallery-dl's "here is a file" record.
            if len(entry) >= 3 and entry[0] == 3 and isinstance(entry[1], str):
                meta = entry[2] if isinstance(entry[2], dict) else {}
                ext = str(meta.get("extension") or "").lstrip(".") or "bin"
                name = str(meta.get("filename") or "") or urlsplit(entry[1]).path
                found.append((entry[1], name, ext))
        return found

    def _download_media_urls(self, items, output, label, ordered_names=False,
                             chapter_images=False):
        """Stream already-resolved media URLs. Returns saved paths in source order."""
        saved = []
        self.ui("plan", len(items), 1, str(output))
        with self._net_opener() as opener:
            for index, (media_url, name, ext) in enumerate(items, 1):
                if not self.is_downloading:
                    break
                self.ui("track", index - 1, len(items), "")
                try:
                    name_hint = f"{index:05d}.{ext}" if ordered_names else name
                    with self._open_url(opener, media_url) as response:
                        path = self._save_stream(response, output, media_url, index,
                                                 len(items), suffix=f".{ext}",
                                                 name_hint=name_hint)
                    if path and chapter_images:
                        from PIL import Image
                        try:
                            with Image.open(path) as image:
                                width, height = image.size
                                image.verify()
                            # Tracking pixels and spacer GIFs are common among
                            # ordinary <img> tags.  They are not chapter pages,
                            # and accepting them can add blank pages or violate
                            # a PDF writer's minimum page dimension.
                            if (width < MIN_CHAPTER_IMAGE_EDGE
                                    or height < MIN_CHAPTER_IMAGE_EDGE
                                    or width * height < MIN_CHAPTER_IMAGE_AREA):
                                path.unlink(missing_ok=True)
                                path = None
                        except Exception:
                            path.unlink(missing_ok=True)
                            path = None
                    if path:
                        saved.append(path)
                        self._batch_log(f"Saved {path.name}", "success")
                    else:
                        reason = ("not a readable chapter page" if chapter_images
                                  else "not media")
                        self._batch_log(f"Skipped {media_url[:70]}: {reason}", "warning")
                except Exception as exc:
                    self._batch_log(f"Skipped {media_url[:70]}: {str(exc)[:80]}", "warning")
        self._flush_logs()
        return saved

    def _no_engine_message(self, kind):
        """What to say when every engine has looked and found nothing."""
        tried = " → ".join(("yt-dlp",) + ENGINES.get(kind, ()))
        detail = getattr(self, "_last_ytdlp_error", "") or ""
        hint = ("\n\n🔧 Nothing downloadable was found at that link."
                f"\n• Engines tried: {tried}"
                "\n• Text-only posts and quote-tweets carry no media to save"
                "\n• Protected or deleted posts need the account that can see them")
        if kind == "x":
            hint += "\n• Use Sign in on the X page if the post is followers-only"
        return (detail[:220] + hint) if detail else hint.lstrip("\n")

    def _try_other_engines(self, kind, url, output, started, quality="best", audio_only=False):
        """Everything that is not yt-dlp, in order. True when one of them worked.

        yt-dlp only knows video. A post that is a photo, or a page it has no
        extractor for, is not a failure to report — it is a job for a different
        engine.
        """
        for engine in ENGINES.get(kind, ()):
            if not self.is_downloading:
                return False
            try:
                if engine == "streamlink":
                    self._batch_log("yt-dlp found no downloadable video — checking for a "
                                    "live or HLS stream with Streamlink.", "info")
                    self._flush_logs()
                    path = self._download_with_streamlink(url, output, quality, audio_only)
                    if path:
                        self._batch_log(f"Streamlink saved {path.name}.", "success")
                        self._flush_logs()
                        self.ui("success", time.time() - started)
                        return True
                    continue
                if audio_only:
                    continue
                if engine == "gallery-dl":
                    self._batch_log("No video there — looking for photos and other "
                                    "media with gallery-dl.", "info")
                    self._flush_logs()
                    items = self._resolve_with_gallery_dl(url, kind)
                    if not items:
                        continue
                    self._batch_log(f"gallery-dl found {len(items)} file(s).", "info")
                elif engine == "html5":
                    self._batch_log("Inspecting the page's HTML5 video and audio sources.", "info")
                    self._flush_logs()
                    items = self._resolve_html5_media(url)
                    if not items:
                        continue
                    self._batch_log(f"HTML5 resolver found {len(items)} media file(s).", "info")
                elif engine == "direct":
                    self._batch_log("Trying the link as a direct media file.", "info")
                    self._flush_logs()
                    items = [(url, "", (Path(urlsplit(url).path).suffix or ".bin").lstrip("."))]
                else:
                    continue

                if self._download_media_urls(items, output, MEDIA_LABELS.get(kind, kind)):
                    self.ui("success", time.time() - started)
                    return True
            except Exception as exc:
                self._batch_log(f"{engine} could not handle this link: {str(exc)[:110]}",
                                "warning")
                self._flush_logs()
        return False

    def download_pdfs(self, url, output, limit):
        """Fetch a PDF/ebook, or every validated document a page links to."""
        started = time.time()
        saved, failed = 0, 0
        self.ui("plan", 1, 1, str(output))
        self._batch_log(f"Documents: {url}", "download")
        try:
            with self._net_opener() as opener:
                # One fetch answers both questions: a direct PDF is saved here and
                # now, and anything else is a page whose links are worth reading.
                try:
                    with self._open_url(opener, url) as response:
                        page_url = response.geturl()
                        path, probe = self._save_document_response(
                            response, output, page_url, 1, 1)
                        if path:
                            self._batch_log(f"Saved {path.name}", "success")
                            self._flush_logs()
                            self.ui("success", time.time() - started)
                            return
                        body = (probe + response.read(8 * 1024 * 1024)).decode(
                            response.headers.get_content_charset() or "utf-8", "replace")
                except urllib.error.HTTPError as exc:
                    raise RuntimeError(f"The site answered HTTP {exc.code} ({exc.reason}).")

                parser = PdfLinkParser()
                parser.feed(body)
                targets, seen = [], set()
                for href in parser.links:
                    if href.startswith(("javascript:", "mailto:", "#")):
                        continue
                    full = urljoin(page_url, href)
                    if urlsplit(full).scheme not in ("http", "https") or full in seen:
                        continue
                    seen.add(full)
                    targets.append(full)
                # Links that name a document are the likely ones, so they go first and
                # the limit spends itself on them rather than on navigation chrome.
                targets.sort(key=lambda u: Path(urlsplit(u).path).suffix.lower()
                             not in DOCUMENT_EXTENSIONS)
                targets = targets[:min(limit or PDF_SCAN_LIMIT, PDF_SCAN_LIMIT)]

                if not targets:
                    raise RuntimeError("That page is not a PDF and links to no documents.")
                self._batch_log(f"Checking {len(targets)} links for documents…", "info")
                self.ui("plan", len(targets), 1, str(output))

                for index, target in enumerate(targets, 1):
                    if not self.is_downloading:
                        break
                    self.ui("track", index - 1, len(targets), "")
                    try:
                        with self._open_url(opener, target, timeout=30) as response:
                            path, _ = self._save_document_response(
                                response, output, response.geturl(), index, len(targets))
                        if path:
                            saved += 1
                            self._batch_log(f"Saved {path.name}", "success")
                    except Exception as exc:
                        failed += 1
                        self._batch_log(f"Skipped {target[:70]}: {str(exc)[:80]}", "warning")

            self._flush_logs()
            if not self.is_downloading:
                self.log("Download stopped.", "warning")
            elif saved:
                self.log(f"{saved} document(s) saved to {output}."
                         + (f" {failed} link(s) could not be read." if failed else ""), "success")
                self.ui("success", time.time() - started)
            else:
                self.ui("failure",
                        "No PDF or ebook files were found behind those links.\n\n"
                        "🔧 Some sites build the document viewer with JavaScript, so the "
                        "link to the file is not in the page source. Open the PDF itself "
                        "in your browser and paste that address instead.")
        except Exception as exc:
            self._flush_logs()
            detail = str(exc)[:300]
            self.log(f"Document download failed: {detail}", "error")
            self.ui("failure", detail)
        finally:
            self.ui("finished")

    def download_manga(self, url, output, limit):
        """Resolve a chapter into an ordered page folder, PDF and CBZ."""
        started = time.time()
        self._batch_log(f"Manga/manhwa chapter: {url}", "download")
        try:
            items = []
            try:
                items = self._resolve_with_gallery_dl(url, "manga")
                items = [item for item in items
                         if f".{item[2].lower()}" in {".jpg", ".jpeg", ".png", ".webp", ".avif"}]
                if items:
                    self._batch_log(f"gallery-dl found {len(items)} chapter page(s).", "info")
            except Exception as exc:
                self._batch_log(f"gallery-dl did not match this chapter: {str(exc)[:100]}",
                                "warning")
            if not items:
                items = self._resolve_html_images(url)
                if items:
                    self._batch_log(f"HTML image resolver found {len(items)} candidate page(s).",
                                    "info")
            items = items[:min(limit or MANGA_PAGE_LIMIT, MANGA_PAGE_LIMIT)]
            if not items:
                raise RuntimeError("No chapter images were found at that link.")

            slug = unquote(Path(urlsplit(url).path.rstrip("/")).name)
            slug = slug or (urlsplit(url).hostname or "chapter")
            page_folder, pdf_path, cbz_path = unique_chapter_paths(output, slug)
            page_folder.mkdir(parents=True)
            pages = self._download_media_urls(
                items, page_folder, "chapter page", ordered_names=True,
                chapter_images=True)
            if not self.is_downloading:
                self.log("Download stopped.", "warning")
                return
            if not pages:
                raise RuntimeError("The chapter links did not return readable images.")

            # Failed page requests can leave gaps. Rename through temporary names
            # so the permanent folder is always a contiguous reading sequence.
            staged = []
            for index, page in enumerate(pages, 1):
                temporary = page_folder / f".ordered-{index:05d}{page.suffix.lower()}"
                os.replace(page, temporary)
                staged.append(temporary)
            pages = []
            for index, page in enumerate(staged, 1):
                ordered = page_folder / f"{index:05d}{page.suffix.lower()}"
                os.replace(page, ordered)
                pages.append(ordered)

            pdf_partial = pdf_path.with_suffix(".pdf.part")
            cbz_partial = cbz_path.with_suffix(".cbz.part")
            self.ui("stage", "building ordered folder, PDF and CBZ")
            try:
                with zipfile.ZipFile(cbz_partial, "w", compression=zipfile.ZIP_STORED) as archive:
                    for page in pages:
                        archive.write(page, page.name)
                os.replace(cbz_partial, cbz_path)
            finally:
                cbz_partial.unlink(missing_ok=True)

            try:
                with tempfile.TemporaryDirectory() as temp:
                    self._write_images_pdf(pages, pdf_partial, Path(temp))
                if not is_valid_pdf_file(pdf_partial):
                    raise RuntimeError("The PDF builder returned an incomplete file.")
                os.replace(pdf_partial, pdf_path)
            finally:
                pdf_partial.unlink(missing_ok=True)

            self._batch_log(f"Saved {len(pages)} ordered pages in {page_folder.name}, "
                            f"plus {pdf_path.name} and {cbz_path.name}.",
                            "success")
            self._flush_logs()
            self.ui("success", time.time() - started)
        except Exception as exc:
            self._flush_logs()
            detail = str(exc)[:300]
            self.log(f"Manga/manhwa download failed: {detail}", "error")
            self.ui("failure", detail +
                    "\n\nOnly public or account-authorized chapter pages are supported. "
                    "DRM and paywalls are not bypassed.")
        finally:
            self.ui("finished")

    def _ytdlp_cmd_for(self, kind, url, output, quality, audio_only, limit, candidate=None,
                       clients=None, impersonate=False, fresh=False):
        cmd = [
            self.ytdlp_cmd, url,
            # Do not inherit a machine-wide yt-dlp.conf. The UI is the source of
            # truth, and a stale format/client override can break a new build.
            "--ignore-config",
            "--output", str(output / "%(title).150B.%(ext)s"),
            "--newline", "--no-warnings", "--ignore-errors",
            # Partials stay in .part files and are renamed on success. Writing
            # straight to the final name leaves a truncated file behind when a
            # run fails, and the next run tries to resume it — which Instagram
            # answers with HTTP 416.
            "--retries", "10", "--fragment-retries", "10",
            "--extractor-retries", "5", "--socket-timeout", "30",
            "--ffmpeg-location", self.ffmpeg_location(),
            # A machine-readable progress line beats scraping the pretty bar.
            # Pipe-separated because titles and speeds both contain spaces.
            "--progress-template",
            "PROGRESS|%(progress._percent_str)s|%(progress._speed_str)s"
            "|%(progress._eta_str)s|%(info.title).70s",
        ]
        if fresh:
            # A .part file can refer to an expired GoogleVideo delivery URL.
            # Keep normal resume support, but restart a delivery-403 retry.
            cmd += ["--no-continue"]
        if audio_only:
            codec = self.download_format.get() if self.download_format.get() in ("mp3", "flac") else "mp3"
            cmd += ["--extract-audio", "--audio-format", codec, "--audio-quality", f"{self.quality.get()}K"]
            if kind == "ytmusic":
                # A music library is only as good as its tags. YouTube Music
                # carries real artist/album/track metadata and cover art, and a
                # file without them is unsortable in any player.
                cmd += ["--embed-metadata", "--embed-thumbnail",
                        "--parse-metadata", "%(release_year,upload_date>%Y)s:%(meta_date)s"]
        elif kind == "instagram":
            # Posts are as often stills as video: a height-capped video selector
            # matches nothing on a photo, and merging into mp4 makes no sense.
            cmd += ["--format", "best"]
        else:
            selector = MEDIA_FORMATS.get(quality, MEDIA_FORMATS["best"])
            if clients:
                # The token-free clients carry a thinner format list, so a strict
                # height cap can match nothing. Take the cap, or whatever exists.
                selector += "/best"
            cmd += ["--format", selector, "--merge-output-format", "mp4"]
        if limit:
            cmd += ["--playlist-end", str(limit)]
        if candidate:
            cmd += ytdlp_cookie_args(candidate)
            # Replay the session as the browser that earned it — but only where
            # that is actually known. On Android every jar comes from this app's
            # own WebView, so the agent is certain. A desktop jar may have been
            # lifted out of Firefox or Edge by --cookies-from-browser, and
            # announcing Chrome over those cookies would manufacture exactly the
            # mismatch this is here to avoid.
            #
            # YouTube is left out on both: its extractor impersonates a
            # different client per player_client, and pinning one agent across
            # all of them contradicts the client it claims to be.
            if IS_ANDROID and kind not in YOUTUBE_KINDS:
                cmd += ["--user-agent", BROWSER_UA]

        # Some extractors and generic pages also ship JavaScript challenges.
        # Supplying Deno is harmless when unused and lets yt-dlp invoke it when
        # an extractor supports EJS or another external challenge script.
        runtime_name, runtime = self.js_runtime()
        if runtime:
            # yt-dlp enables only Deno by default, so a lower-priority engine is
            # ignored unless the defaults are cleared first.
            if runtime_name != "deno":
                cmd += ["--no-js-runtimes"]
            cmd += ["--js-runtimes", f"{runtime_name}:{runtime}"]

        # Impersonation is deliberately a challenge-only General fallback.
        # yt-dlp warns that forcing it globally can reduce speed and stability.
        if kind == "general" and impersonate:
            # A generic family lets yt-dlp select the newest Chrome profile its
            # bundled curl_cffi actually supports. An empty "any" target is
            # fragile on Windows because some launchers discard empty arguments.
            cmd += ["--impersonate", "chrome"]

        if kind in YOUTUBE_KINDS:
            # YouTube documents rate limits separately from bot challenges.
            # Modest spacing keeps playlists from burning through a session.
            cmd += ["--sleep-requests", "1", "--sleep-interval", "3",
                    "--max-sleep-interval", "6"]

            token = self.youtube_po_token.get().strip()
            selected_clients = clients
            extractor_args = []
            if token and PO_TOKEN_PATTERN.fullmatch(token):
                selected_clients = selected_clients or "mweb,default"
                if "mweb" not in selected_clients.split(","):
                    selected_clients = "mweb," + selected_clients
                extractor_args.append(f"po_token={token}")
            if selected_clients:
                extractor_args.insert(0, f"player_client={selected_clients}")
            if extractor_args:
                cmd += ["--extractor-args", "youtube:" + ";".join(extractor_args)]
        if self.use_proxy.get() and self.proxy_url.get().strip():
            # yt-dlp speaks SOCKS5 directly, so the HTTP bridge is not needed here.
            cmd += ["--proxy", normalize_proxy_url(self.proxy_url.get(), self.proxy_type.get())]
        return cmd

    def download_with_ytdlp(self, kind, url, output, quality, audio_only, limit):
        self._media_done = 0
        self._media_total = None
        self._stream_index = 0
        # A merged download writes video and audio separately, so per-file
        # percentage would otherwise reach 100% halfway through the job.
        streams = 2 if (not audio_only and kind != "instagram"
                        and "+" in MEDIA_FORMATS.get(quality, "")) else 1
        self.ui("plan", 1, streams, str(output))
        started = time.time()
        self._batch_log(f"{MEDIA_LABELS[kind]}: {url}", "download")

        # A cookie source that exists is not a cookie source that reads: a running
        # Chromium browser locks its DB, and Chrome 127+ encrypts it app-bound.
        # Probing costs a second and tells us before the download instead of after.
        cookies = [None]
        stored_site = cookie_site(kind)
        use_stored_session = stored_site in SIGNIN_PAGES and registry_entry(stored_site)
        if kind in COOKIE_SOURCES or use_stored_session:
            cookies = self.usable_cookies(
                COOKIE_DOMAINS.get(kind), require_session=(kind == "soundcloud")) or [None]
            if cookies == [None] and kind == "instagram":
                self._batch_log(
                    "No readable cookies — log into Instagram in Firefox, or drop a cookies.txt "
                    "beside the app. Trying anyway.", "warning")

        if cookies == [None] and kind == "soundcloud" and use_stored_session:
            self._batch_log(
                "The saved SoundCloud session is unavailable or expired. "
                "Trying the public stream; sign in again for Go/Go+ access.", "warning")

        # Each attempt is a cookie jar, a set of YouTube player clients, and an
        # optional browser-impersonation pass. The ladder is climbed when a run says why it
        # failed, and every rung is used at most once — otherwise a jar that can
        # never satisfy the site gets retried forever.
        # A phone starts on the rung the desktop escalates to. yt-dlp's default
        # YouTube clients need a proof-of-origin token and a solved JavaScript
        # challenge; without both, the run is answered with "sign in to confirm
        # you're not a bot" and then 403 on the media itself. android_vr asks
        # for neither, so on Android it is the opening move rather than the
        # consolation — and the ladder above still has its remaining rungs.
        # ...but only when there is no session to spend. The token-free clients
        # cannot carry cookies at all, so opening with one where the user has
        # deliberately signed in would quietly throw their login away and fail
        # anything members-only or age-gated. A jar keeps the cookie-capable
        # default, and the ladder is still there underneath it.
        lead = (YT_CLIENT_LADDER[0]
                if IS_ANDROID and kind in YOUTUBE_KINDS else None)
        plan = [{"cookies": c, "clients": lead if (lead and not c) else None,
                 "impersonate": False, "fresh": False}
                for c in cookies]
        spent = set()
        if lead and all(not step["cookies"] for step in plan):
            spent.add("clients-0")      # opened with it; do not offer it again

        try:
            while plan:
                attempt = plan.pop(0)
                candidate = attempt["cookies"]
                if candidate:
                    self._batch_log(f"Cookies: {candidate_label(candidate)}", "info")
                verdict = self._run_ytdlp(
                    kind, url, output, quality, audio_only, limit, candidate, started,
                    # Instagram belongs here too: it probes cookies up front, but a
                    # stored session that expired between runs only shows up as a
                    # sign-in demand mid-download, and that is what the ladder fixes.
                    may_borrow_cookies=(kind in COOKIE_ON_DEMAND or kind in COOKIE_SOURCES),
                    clients=attempt["clients"], impersonate=attempt["impersonate"],
                    fresh=attempt["fresh"])

                if verdict is None:
                    return

                if verdict == "other-engine":
                    if self._try_other_engines(kind, url, output, started, quality, audio_only):
                        return
                    self._flush_logs()
                    self.ui("failure", self._no_engine_message(kind))
                    return

                if verdict in ("signin", "challenge"):
                    nxt = self._next_signin_step(kind, url, attempt, spent,
                                                 challenge=(verdict == "challenge"))
                    if nxt:
                        plan.insert(0, nxt)
                        continue
                    self._flush_logs()
                    if verdict == "challenge" and kind in YOUTUBE_KINDS:
                        detail = (
                            "YouTube refused the media URL after fresh-link and "
                            "alternate-client retries.")
                        self.log(detail, "error")
                        self.ui(
                            "failure",
                            detail + "\n\nTry the app's proxy, a different connection, or "
                            "a signed-in YouTube session. A PO token may be required when "
                            "YouTube binds delivery URLs to proof-of-origin.")
                        return
                    self._report_signin_failure(kind, url)
                    return

                # verdict == "cookies": that jar is unreadable.
                if plan:
                    self._batch_log(
                        f"Could not read cookies from {candidate_label(candidate)} — trying the next source.",
                        "warning")
                    self._flush_logs()
                    continue
                self._flush_logs()
                self._report_ytdlp_failure(kind, f"Could not read cookies from {candidate_label(candidate)}.")
                return
        finally:
            self.ui("finished")

    def _next_signin_step(self, kind, url, attempt, spent, challenge=False):
        """The next thing worth trying against a sign-in demand, or None.

        Cheapest first: another player client costs nothing, a cookie jar needs
        an account. Each rung is spent once, so a jar that can never satisfy the
        site is not retried forever.
        """
        if kind == "general" and challenge and "impersonate" not in spent:
            spent.add("impersonate")
            self._batch_log(
                "The site returned a browser challenge - retrying with browser impersonation.",
                "warning")
            self._flush_logs()
            return {**attempt, "impersonate": True}

        if kind in YOUTUBE_KINDS:
            for index, clients in enumerate(YT_CLIENT_LADDER):
                key = f"clients-{index}"
                if key in spent:
                    continue
                spent.add(key)
                self._batch_log(
                    f"YouTube delivery failed — requesting a fresh media URL as {clients}.",
                    "warning")
                self._flush_logs()
                return {**attempt, "clients": clients, "fresh": True}

        # The stored login outlives its cookies, because the site rotates them.
        # If this app has ever signed in here, re-read the session before
        # deciding there is none — it costs seconds and needs no one's attention.
        site = cookie_site(kind)
        domain = cookie_domain_for(kind, url)
        if site in SIGNIN_PAGES and "refresh" not in spent and read_registry().get(site):
            spent.add("refresh")
            self._batch_log(f"Refreshing the stored {MEDIA_LABELS.get(site, site)} session…",
                            "warning")
            self._flush_logs()
            if self.refresh_session(site, announce=False):
                session = self.usable_cookies(domain, require_session=True)
                if session:
                    self._batch_log("Retrying with the refreshed session.", "success")
                    self._flush_logs()
                    return {**attempt, "cookies": session[0]}

        if "cookies" not in spent:
            spent.add("cookies")
            # Only a jar that actually holds this site's session can help here.
            session = self.usable_cookies(domain, require_session=True)
            if session:
                suffix = f" for {domain}" if kind == "general" else ""
                self._batch_log(
                    f"Retrying with the {candidate_label(session[0])} session{suffix}.",
                    "warning")
                self._flush_logs()
                return {**attempt, "cookies": session[0]}
            self._batch_log("No signed-in session for this site — skipping cookies.",
                            "warning")

        return None

    # ------------------------------------------------------------------
    # Signing in, without an extension or a text file
    # ------------------------------------------------------------------
    def sign_in(self, kind):
        """Open a real login window. Its cookies become the jar we download with.

        This is the same thing the browser extension does, minus the extension:
        the app owns this window, so nothing locks the database and nothing has
        to be exported by hand.
        """
        page = SIGNIN_PAGES.get(kind)
        if not page:
            self.notify(f"{kind.title()} does not need a sign-in.", "err")
            return
        if self._login_window is not None:
            self.notify("A sign-in window is already open.", "err")
            return

        login_url = page[0]
        self._login_kind = kind
        self._login_window = webview.create_window(
            f"Sign in to {kind.title()} · Blue Knight",
            url=login_url, width=980, height=800, min_size=(600, 560),
            background_color="#05080d",
        )
        self._login_window.events.closed += self._login_closed
        self.ui("signin", kind)
        self.log(f"Sign in to {kind.title()} in the window that opened, then press "
                 f"“Finish sign-in”.", "info")

    def _login_closed(self):
        """The window went away. If cookies were never taken, say so."""
        if self._login_window is not None:
            self._login_window = None
            self._login_kind = None
            self.ui("signin", None)
            self.log("Sign-in window closed before the cookies were saved.", "warning")

    def _harvest_session(self, window, kind, patience=12, first_page_only=False):
        """Read the session out of a window and write it as this site's jar.

        Waits for a session cookie rather than a fixed pause: a login redirect
        can still be in flight when the page looks finished, and a jar without
        the session cookie is indistinguishable from being signed out. The wait
        is short when this runs behind a download, long when a person is at the
        keyboard finishing a login.
        """
        _, harvest_urls, site = SIGNIN_PAGES[kind]
        if first_page_only:
            harvest_urls = harvest_urls[:1]
        wanted = SITE_SESSION_COOKIES.get(kind, ())
        harvested, seen, found_session = [], set(), False

        for url in harvest_urls:
            try:
                window.load_url(url)
            except Exception as exc:
                self.log(f"Sign-in: could not open {url} ({str(exc)[:90]})", "warning")
                continue
            for _ in range(patience):                 # one second per turn
                time.sleep(1.0)
                try:
                    cookies = window.get_cookies() or []
                except Exception:
                    continue
                for cookie in cookies:
                    for name, morsel in cookie.items():
                        key = (morsel["domain"], name)
                        if key not in seen:
                            seen.add(key)
                            harvested.append(cookie)
                        if name in wanted:
                            found_session = True
                if found_session:
                    break

        jar = cookie_dir(kind) / f"{site}_cookies.txt"
        count = write_netscape_jar(harvested, jar)
        return jar, count, found_session

    def finish_sign_in(self):
        """Harvest the session from the login window and write it as a jar."""
        window, kind = self._login_window, self._login_kind
        if not window or not kind:
            self.notify("No sign-in window is open.", "err")
            return

        try:
            jar, count, signed_in = self._harvest_session(window, kind)
        except Exception as exc:
            self.log(f"Could not read the sign-in cookies: {str(exc)[:200]}", "error")
            self.notify(f"Could not read the cookies.\n\n{str(exc)[:200]}", "err")
            return
        finally:
            self._login_window = None
            self._login_kind = None
            self.ui("signin", None)
            with contextlib.suppress(Exception):
                window.destroy()

        if not count:
            self.log(f"No cookies came back from {kind.title()}. Was the sign-in finished?",
                     "warning")
            self.notify("No cookies were found. Finish signing in, then try again.", "err")
            return
        if not signed_in:
            self.log(f"{jar.name}: {count} cookies, but no {kind.title()} session cookie.",
                     "warning")
            self.notify(f"Saved {count} cookies, but you are not signed in yet. "
                        f"Complete the login, then press Finish sign-in again.", "err")
            return

        entry = record_jar(kind, jar, "sign-in")
        self.log(f"Signed in to {kind.title()}: {count} cookies saved to "
                 f"{jar.parent.name}/{jar.name}.", "success")
        self.ui("cookie_status", self.cookie_status())
        self.notify(f"Signed in to {kind.title()}. {entry['cookies']} cookies kept.", "ok")

    def refresh_session(self, kind, announce=True):
        """Take a fresh jar from the browser profile, without asking anyone.

        The window the app owns stays signed in between runs, and YouTube rotates
        cookies on any open tab — so the jar goes stale while the login does not.
        Re-reading it silently is what keeps downloads working over time.
        """
        if kind not in SIGNIN_PAGES or self._login_window is not None:
            return False

        if IS_ANDROID:
            # A phone has no browser profile to re-read. The live session is in
            # the system CookieManager — the same place the in-app sign-in put
            # it — and only the shell can reach that, so it installs the hook.
            # Without this the stored login is never refreshed, and a rotated
            # YouTube cookie looks exactly like never having signed in at all.
            if SESSION_REFRESHER is None:
                return False
            try:
                count, signed_in = SESSION_REFRESHER(kind)
            except Exception as exc:
                self.log(f"Could not refresh the {kind.title()} session: {str(exc)[:150]}",
                         "warning")
                return False
            if not signed_in:
                return False
            if announce:
                self.log(f"Refreshed the {kind.title()} session ({count} cookies).",
                         "success")
            self.ui("cookie_status", self.cookie_status())
            return True

        window = None
        try:
            window = webview.create_window(
                f"{kind.title()} session", url=SIGNIN_PAGES[kind][1][0],
                width=900, height=700, hidden=True)
            # A download is waiting on this, so give it seconds, not a minute.
            jar, count, signed_in = self._harvest_session(
                window, kind, patience=5, first_page_only=True)
        except Exception as exc:
            self.log(f"Could not refresh the {kind.title()} session: {str(exc)[:150]}",
                     "warning")
            return False
        finally:
            if window is not None:
                with contextlib.suppress(Exception):
                    window.destroy()

        if not signed_in:
            if announce:
                self.log(f"No stored {kind.title()} session — use Sign in on the "
                         f"{kind.title()} page.", "warning")
            return False
        record_jar(kind, jar, "refresh")
        self.ui("cookie_status", self.cookie_status())
        if announce:
            self.log(f"Refreshed the {kind.title()} session: {count} cookies.", "success")
        return True

    def cookie_status(self):
        """What the page shows next to each Sign in button."""
        status = {}
        for kind in SIGNIN_PAGES:
            entry = registry_entry(kind)
            if entry and entry.get("signed_in"):
                age = max(0, int(time.time()) - entry.get("saved_at", 0))
                hours = age // 3600
                status[kind] = {
                    "signed_in": True,
                    "detail": f"{entry['cookies']} cookies · "
                              + ("just now" if hours < 1 else f"{hours}h old"),
                }
            else:
                status[kind] = {"signed_in": False, "detail": ""}
        return status

    def cookie_export_recipe(self, kind):
        """The export that actually survives, spelled out.

        A jar copied from a normal window dies when the site rotates the
        session; one taken in a private window and never used again does not.
        """
        site = {"youtube": "www.youtube.com", "instagram": "www.instagram.com",
                "tiktok": "www.tiktok.com", "soundcloud": "soundcloud.com",
                "x": "x.com"}.get(kind, "the site")
        return (f"\n\n🍪 Export a cookies.txt — this lasts:"
                "\n1. Install the “Get cookies.txt LOCALLY” extension"
                "\n2. Open a private/incognito window and sign in"
                f"\n3. In that same tab, go to https://{site}/robots.txt"
                "\n4. Export as Netscape format"
                f"\n5. Save it in the app's cookies folder, or your Downloads"
                "\n\nThe app picks up any *cookies*.txt in those folders and uses the one"
                " holding this site's session. Do not open that private window again —"
                " closing it is what keeps the session alive.")

    def _report_signin_failure(self, kind, url=None):
        domain = cookie_domain_for(kind, url or "")
        label = domain if kind == "general" and domain else kind.title()
        detail = f"{label} wants a signed-in session and no usable cookies were found."
        self.log(f"{kind.title()} download failed: {detail}", "error")

        # A browser that is merely running is the most common — and most easily
        # fixed — reason we have no session, so name it instead of listing options.
        locked = [candidate_label(c) for c in cookie_candidates(domain)
                  if locked_by_browser(c)]
        opening = ""
        if locked:
            names = " and ".join(dict.fromkeys(locked))
            opening = (f"\n\n🔧 {names} is open, and a running browser keeps its cookies locked."
                       f"\n• Close {names} completely — check the tray — then press Download again")
        token_hint = ""
        if kind in YOUTUBE_KINDS:
            token_hint = (
                "\n\nYouTube's JavaScript challenge solver was tried first. If the IP still "
                "requires proof-of-origin, add an mweb.gvs+TOKEN in Connection settings.")
        if kind == "general":
            export_hint = (
                f"\n\nExport a Netscape cookies.txt while signed in to {domain or 'the site'}, "
                "then place it beside the app or in Downloads. Only cookies matching the "
                "requested hostname are eligible for this retry.")
        else:
            export_hint = self.cookie_export_recipe(kind)
        self.ui("failure", detail + opening + token_hint + export_hint)

    def _report_ytdlp_failure(self, kind, detail):
        self.log(f"{kind.title()} download failed: {detail}", "error")
        self.ui("failure", detail +
                "\n\n🔧 Cookies are the sticking point:"
                "\n• A running Edge/Chrome/Brave locks its cookie database — close it"
                "\n• Chrome 127 and newer encrypt it app-bound; yt-dlp cannot open that at all"
                "\n• Log in with Firefox, which stays readable"
                "\n• Or export a cookies.txt into this folder, your Downloads, or the download"
                " folder — an exported file wins over every browser")

    def _run_ytdlp(self, kind, url, output, quality, audio_only, limit, candidate, started,
                   may_borrow_cookies=False, clients=None, impersonate=False, fresh=False):
        """One yt-dlp run.

        Returns None when the job is over (success, stop, or a failure already
        reported), or a verdict the caller can act on: "cookies" for an unreadable
        jar, "signin" for a site that wants a session, or "challenge" for a
        browser/media-delivery challenge that has another safe fallback.
        """
        error_lines = []
        try:
            cmd = self._ytdlp_cmd_for(kind, url, output, quality, audio_only, limit, candidate,
                                      clients, impersonate, fresh)
            self.process = pyshell.popen(cmd, bufsize=8192)
            for raw in self.process.stdout:
                if not self.is_downloading:
                    break
                line = raw.strip()
                if not line:
                    continue

                if line.startswith("PROGRESS|"):
                    now = time.time()
                    if now - self.last_ui_update < UI_UPDATE_INTERVAL:
                        continue
                    self.last_ui_update = now
                    parts = (line.split("|") + ["", "", "", ""])[1:5]
                    percent, speed, eta, title = (p.strip() for p in parts)
                    try:
                        pct = float(percent.rstrip("%"))
                    except ValueError:
                        continue
                    # The stream index is what stops a two-file download from
                    # showing 100% while the audio has not started.
                    self.ui("bytes", pct, title, speed, eta, self._stream_index)
                    continue

                item = self._re_ytdlp_item.search(line)
                if item:
                    self._media_done, self._media_total = int(item.group(1)), int(item.group(2))
                    self.ui("total", self._media_total)
                    self.ui("track", self._media_done, self._media_total, "")
                    continue

                if "[download] Destination:" in line or "has already been downloaded" in line:
                    # A new destination means the previous stream finished.
                    self._stream_index += 1
                    name = line.split("Destination:")[-1].strip() or line.strip()
                    self.ui("file", Path(name).name, self._stream_index)
                    self._batch_log(line[:200], "success")
                elif "[Merger]" in line or "Merging formats" in line:
                    self.ui("stage", "merging video and audio")
                elif "[ExtractAudio]" in line:
                    self.ui("stage", "converting audio")
                elif self._re_error.search(line):
                    error_lines.append(line)
                    self._batch_log(line[:200], "error")

            self.process.wait()
            code = self.process.returncode

            if not self.is_downloading:
                self._flush_logs()
                self.log("Download stopped.", "warning")
                return None
            self._flush_logs()
            if code == 0:
                self.ui("success", time.time() - started)
                return None
            joined = "\n".join(error_lines[-6:])
            # Say nothing yet when the caller still has a move to make.
            if candidate and COOKIE_FAILURE.search(joined):
                return "cookies"
            if kind in YOUTUBE_KINDS and YOUTUBE_MEDIA_DENIED.search(joined):
                return "challenge"
            if kind == "general" and GENERAL_CHALLENGE.search(joined):
                return "challenge"
            if may_borrow_cookies and SIGNIN_DEMANDED.search(joined):
                return "signin"
            # "There is no video here" is a routing answer, not a failure: a photo
            # post and a site yt-dlp has never heard of both land here, and both
            # are somebody else's job.
            if ENGINES.get(kind) and NO_VIDEO_FOR_YTDLP.search(joined):
                self._last_ytdlp_error = joined
                return "other-engine"
            raise RuntimeError(joined or f"yt-dlp exited with code {code}")
        except FileNotFoundError:
            self._flush_logs()
            self.log("yt-dlp could not be launched. Run Check for updates.", "error")
            self.ui("failure", "yt-dlp could not be launched.")
        except Exception as exc:
            detail = str(exc)[:300]
            self._flush_logs()
            self.log(f"{kind.title()} download failed: {detail}", "error")
            hint = ""
            if SIGNIN_DEMANDED.search(detail):
                hint = ("\n\n🔧 The site asked for a signed-in session."
                        "\n• Log in with Firefox, or drop a cookies.txt beside the app"
                        "\n• A proxy in Connection settings also helps when it is region bait")
            elif "drm" in detail.lower():
                # Naming the wall beats a retry loop against it. Nothing in this
                # app removes DRM, and no setting is going to change that.
                hint = ("\n\n🔧 This title is DRM-protected, so it cannot be downloaded."
                        "\n• Netflix and Crunchyroll encrypt their whole catalogue"
                        "\n• SoundCloud Go+ tracks are protected the same way"
                        "\n• Unprotected titles on the same site still work normally")
            elif kind == "tiktok":
                hint = "\n\n🔧 TikTok rate-limits hard. Enable a proxy and retry in a minute."
            elif kind == "instagram":
                hint = ("\n\n🔧 Instagram serves almost nothing logged out."
                        "\n• Log into Instagram in Firefox, Edge, or Chrome, then retry"
                        "\n• A running Chromium browser locks its cookies — close Edge/Chrome,"
                        " or use Firefox"
                        "\n• Or export a cookies.txt and drop it beside the app"
                        "\n• Private accounts need the logged-in user to follow them")
            self.ui("failure", detail + hint)
        finally:
            # "finished" belongs to the caller: a cookie retry is still the same job.
            self.process = None
        return None

    def _batch_log(self, message, level="info"):
        """Batch log messages"""
        self.log_batch.append((message, level))
        if len(self.log_batch) >= LOG_BATCH_SIZE:
            self._flush_logs()

    def _flush_logs(self):
        """Write batched logs to UI"""
        if not self.log_batch:
            return
        self.ui_queue.put(("batch_log", self.log_batch.copy()))
        self.log_batch.clear()

    def _get_troubleshooting(self, error_text, proxy_enabled):
        """Generate troubleshooting"""
        if proxy_enabled and "RequestError" in error_text:
            return (
                "\n🔧 PROXY TROUBLESHOOTING:"
                "\n• Verify proxy is running"
                "\n• Verify that HTTP or SOCKS5 matches the configured port"
                "\n• Try a matching v2rayN preset"
                "\n• Test proxy with 'Test Proxy' button"
            )
        elif "RequestError" in error_text:
            return (
                "\n🔧 TROUBLESHOOTING:"
                "\n• Check internet connection"
                "\n• Try VPN/proxy if region-blocked"
                "\n• Use 'V2RayN HTTP' preset for proxy"
                "\n• Update spotDL in Advanced settings"
            )
        elif "blocked by YouTube" in error_text.lower():
            return (
                "\n🔧 TROUBLESHOOTING:"
                "\n• YouTube Music blocked connection"
                "\n• Enable proxy (V2RayN HTTP preset)"
                "\n• Try 'best' format instead"
            )
        else:
            return (
                "\n🔧 TROUBLESHOOTING:"
                "\n• Check internet connection"
                "\n• Update spotDL"
                "\n• Verify Spotify link"
            )

    def stop_download(self):
        # spotDL runs as a child process; yt-dlp runs in-process and stops via its hook.
        if self.is_downloading:
            self.is_downloading = False
            if self.process:
                try:
                    self.process.terminate()
                except Exception:
                    pass
            self._flush_logs()
            self.log("Download stopped.", "warning")
            self.status.set("Ready")
            self.ui("status", "Stopped")

    # ------------------------------------------------------------------
    # Event queue drained by the page
    # ------------------------------------------------------------------
    def ui(self, action, *args):
        self.ui_queue.put((action, args))

    def log(self, message, level="info"):
        self.ui_queue.put(("batch_log", [(message, level)]))

    def download_finished(self):
        self.is_downloading = False
        self.process = None
        if self.status.get() != "Error":
            self.status.set("Ready")

    # ------------------------------------------------------------------
    # Utilities
    # ------------------------------------------------------------------
    def check_setup(self):
        if self.updating:
            self.notify("The component update is still running.", "err")
            return
        self.log("Checking components...", "info")

        if (self._is_windows_executable(self.ffmpeg_exe)
                and self.tool_version(self.ffmpeg_exe) != "unknown"):
            self.ffmpeg_cmd = str(self.ffmpeg_exe)
        else:
            self.ffmpeg_cmd = bundled_tool(f"ffmpeg{EXE}") or shutil.which("ffmpeg")

        self.spotdl_cmd = ("Android Spotify engine" if IS_ANDROID
                           else self.find_tool("spotdl"))
        self.ytdlp_cmd = self.find_tool("yt-dlp")
        self.deno_cmd = None if IS_ANDROID else self.find_tool("deno")

        required = [("Spotify engine", self.spotdl_cmd), ("yt-dlp", self.ytdlp_cmd),
                    ("FFmpeg", self.ffmpeg_cmd)]
        if not IS_ANDROID:
            required.append(("Deno", self.deno_cmd))
        missing = [label for label, found in required if not found]
        self.ui("components", self.component_summary())
        if not missing:
            self.status.set("Ready")
            self.log("All components ready.", "success")
            self.notify("Everything is installed and ready.", "ok")
        else:
            self.notify("Missing: " + ", ".join(missing) + ". Run Check for updates, or restart the app.", "err")

    # ------------------------------------------------------------------
    # Conversion
    # ------------------------------------------------------------------
    @staticmethod
    def normalize_media_conversion_options(options):
        """Validate the converter payload received from the web UI."""
        if isinstance(options, str):
            options = {"format": options}
        if not isinstance(options, dict):
            raise ValueError("Invalid converter settings.")
        output_format = str(options.get("format", "mp4")).lower()
        if output_format not in MEDIA_CONVERT_FORMATS:
            raise ValueError("Choose a supported media output format.")
        codec = str(options.get("codec", "auto")).lower()
        if codec not in VIDEO_CONVERT_CODECS:
            raise ValueError("Choose a supported video codec.")
        quality = str(options.get("quality", "balanced")).lower()
        if quality not in VIDEO_CONVERT_QUALITY:
            raise ValueError("Choose a supported video quality.")
        resize = str(options.get("resize", "source")).lower()
        if resize not in {*VIDEO_RESIZE_PRESETS, "custom"}:
            raise ValueError("Choose a supported resize preset.")
        fps = str(options.get("fps", "source")).lower()
        if fps not in {"source", "24", "30", "60"}:
            raise ValueError("Choose a supported frame rate.")
        frame_mode = str(options.get("frame_mode", "standard")).lower()
        if frame_mode not in {"standard", "motion"}:
            raise ValueError("Choose a supported frame processing mode.")
        upscale = str(options.get("upscale", "lanczos")).lower()
        if upscale not in VIDEO_UPSCALE_FILTERS:
            raise ValueError("Choose a supported upscale engine.")
        video_enhance = str(options.get("video_enhance", "none")).lower()
        if video_enhance not in VIDEO_ENHANCEMENT_FILTERS:
            raise ValueError("Choose a supported video enhancement preset.")
        audio_enhance = str(options.get("audio_enhance", "none")).lower()
        if audio_enhance not in AUDIO_ENHANCEMENT_FILTERS:
            raise ValueError("Choose a supported audio enhancement preset.")
        bitrate = str(options.get("audio_bitrate", "192"))
        if bitrate not in {"128", "192", "256", "320"}:
            raise ValueError("Choose a supported audio bitrate.")
        width = int(options.get("width") or 0)
        height = int(options.get("height") or 0)
        if resize == "custom":
            if not width and not height:
                raise ValueError("Enter a custom width, height, or both.")
            if width and not 16 <= width <= 7680:
                raise ValueError("Custom width must be between 16 and 7680 pixels.")
            if height and not 16 <= height <= 4320:
                raise ValueError("Custom height must be between 16 and 4320 pixels.")
        if output_format == "webm":
            codec = "vp9"
        elif output_format == "avi":
            codec = "h264"
        elif output_format == "mov" and codec == "vp9":
            codec = "h264"
        elif codec == "auto":
            codec = "h264"
        return {
            "format": output_format, "codec": codec, "quality": quality,
            "resize": resize, "fps": fps, "frame_mode": frame_mode,
            "audio_bitrate": bitrate,
            "width": width, "height": height, "upscale": upscale,
            "video_enhance": video_enhance, "audio_enhance": audio_enhance,
        }

    @staticmethod
    def media_conversion_command(ffmpeg, source, target, options):
        """Build a complete resize, codec, quality and audio transcode command."""
        output_format = options["format"]
        bitrate = options["audio_bitrate"]
        base = [
            ffmpeg, "-hide_banner", "-y", "-i", str(source),
            "-map", "0:v:0?", "-map", "0:a:0?", "-map_metadata", "0",
        ]
        audio_filter = AUDIO_ENHANCEMENT_FILTERS[options["audio_enhance"]]
        # loudnorm processes internally at 192 kHz. Bring enhanced output back
        # to a widely supported professional rate so FLAC/AAC muxers do not
        # reject the filter's temporary sample rate.
        audio_args = ["-af", audio_filter, "-ar", "48000"] if audio_filter else []
        if output_format == "mp3":
            return base + audio_args + [
                "-vn", "-c:a", "libmp3lame", "-b:a", f"{bitrate}k", str(target)]
        if output_format == "flac":
            return base + audio_args + [
                "-vn", "-c:a", "flac", "-compression_level", "8", str(target)]
        if output_format == "wav":
            return base + audio_args + ["-vn", "-c:a", "pcm_s24le", str(target)]
        if output_format in {"m4a", "aac"}:
            return base + audio_args + [
                "-vn", "-c:a", "aac", "-b:a", f"{bitrate}k", str(target)]
        if output_format == "ogg":
            audio_codec = "opus" if IS_ANDROID else "libopus"
            experimental = ["-strict", "-2"] if IS_ANDROID else []
            return base + audio_args + [
                "-vn", "-c:a", audio_codec, *experimental,
                "-b:a", f"{bitrate}k", str(target)]

        filters = list(VIDEO_ENHANCEMENT_FILTERS[options["video_enhance"]][:1])
        resize = options["resize"]
        scale_flags = VIDEO_UPSCALE_FILTERS[options["upscale"]]
        if resize == "custom":
            width = options["width"] or -2
            height = options["height"] or -2
            filters.append(
                f"scale={width}:{height}:force_original_aspect_ratio=decrease:"
                f"force_divisible_by=2:flags={scale_flags}")
        elif VIDEO_RESIZE_PRESETS[resize]:
            width, height = VIDEO_RESIZE_PRESETS[resize]
            filters.append(
                f"scale={width}:{height}:force_original_aspect_ratio=decrease:"
                f"force_divisible_by=2:flags={scale_flags}")
        filters.extend(VIDEO_ENHANCEMENT_FILTERS[options["video_enhance"]][1:])
        if options["fps"] != "source":
            if options["frame_mode"] == "motion":
                filters.append(
                    f"minterpolate=fps={options['fps']}:mi_mode=mci:mc_mode=aobmc:"
                    "me_mode=bidir:vsbmc=1")
            else:
                filters.append(f"fps={options['fps']}")
        if filters:
            base += ["-vf", ",".join(filters)]

        crf = str(VIDEO_CONVERT_QUALITY[options["quality"]])
        codec = options["codec"]
        preset = "slow" if options["quality"] in {"master", "high"} else (
            "fast" if options["quality"] == "compact" else "medium")
        if IS_ANDROID:
            # The APK's FFmpeg uses Android's hardware codecs instead of the
            # desktop-only libx264/libx265/libvpx libraries. Bitrate is used
            # because MediaCodec encoders do not accept CRF or x264 presets.
            android_codec = {
                "h264": "h264_mediacodec",
                "h265": "hevc_mediacodec",
                "vp9": "vp9_mediacodec",
            }[codec]
            video_bitrate = {
                "master": "16M", "high": "10M", "balanced": "6M", "compact": "3M",
            }[options["quality"]]
            command = base + ["-c:v", android_codec, "-b:v", video_bitrate,
                              "-pix_fmt", "yuv420p"]
        elif codec == "h265":
            command = base + ["-c:v", "libx265", "-preset", preset, "-crf", crf]
            if output_format in {"mp4", "mov"}:
                command += ["-tag:v", "hvc1"]
        elif codec == "vp9":
            cpu_used = "2" if options["quality"] in {"master", "high"} else "4"
            command = base + ["-c:v", "libvpx-vp9", "-crf", crf, "-b:v", "0",
                              "-row-mt", "1", "-cpu-used", cpu_used]
        else:
            command = base + ["-c:v", "libx264", "-preset", preset, "-crf", crf,
                              "-pix_fmt", "yuv420p"]
        command += audio_args
        if output_format == "webm":
            command += ["-c:a", "opus" if IS_ANDROID else "libopus"]
            if IS_ANDROID:
                command += ["-strict", "-2"]
            command += ["-b:a", f"{bitrate}k"]
        elif output_format == "avi":
            command += ["-c:a", "libmp3lame", "-b:a", f"{bitrate}k"]
        else:
            command += ["-c:a", "aac", "-b:a", f"{bitrate}k"]
        if output_format in {"mp4", "mov"}:
            command += ["-movflags", "+faststart"]
        return command + [str(target)]

    def start_media_conversion(self, path, options):
        try:
            options = self.normalize_media_conversion_options(options)
        except (TypeError, ValueError) as exc:
            self.notify(str(exc), "err")
            return
        if self.updating or self.setup_running:
            self.notify("Components are being updated. Try again when the update finishes.", "err")
            return
        if self.is_converting:
            self.notify("A conversion is already running.", "err")
            return
        self.is_converting = True
        threading.Thread(
            target=self.convert_media_file, args=(path, options), daemon=True).start()

    def convert_media_file(self, path, options):
        source = Path(path)
        partial = None
        output_format = options["format"]
        self.ui("converting", True, "media")
        try:
            if not self.ffmpeg_cmd or not Path(self.ffmpeg_cmd).is_file():
                raise RuntimeError("FFmpeg is not ready. Run Check setup first.")
            if not source.is_file():
                raise RuntimeError("The selected media file no longer exists.")
            target_parent = (download_dir("converted", self.save_folder.get().strip())
                             if IS_ANDROID else source.parent)
            target = unique_path(target_parent, safe_filename(
                source.stem, fallback="converted", suffix=f".{output_format}"))
            partial = target.with_name(target.stem + ".working" + target.suffix)
            partial.unlink(missing_ok=True)
            self.log(f"Converting {source.name} to {output_format.upper()}...", "info")
            command = self.media_conversion_command(
                self.ffmpeg_cmd, source, partial, options)
            result = subprocess.run(
                command, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                text=True, encoding="utf-8", errors="replace",
                creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
            )
            if result.returncode != 0 or not partial.is_file() or partial.stat().st_size == 0:
                detail = (result.stderr or result.stdout or "FFmpeg returned no output").strip()
                raise RuntimeError(detail.splitlines()[-1][:220])
            os.replace(partial, target)
            self.log(f"Converted: {target.name}", "success")
            self.notify(f"Created {target.name}", "ok")
        except Exception as exc:
            self.log(f"Media conversion failed: {str(exc)[:250]}", "error")
            self.notify(str(exc)[:220], "err")
        finally:
            if partial:
                partial.unlink(missing_ok=True)
            self.is_converting = False
            self.ui("converting", False, "media")

    @staticmethod
    def _converter_program(name):
        """LibreOffice or Calibre, wherever this platform installs it.

        Both ship a launcher on PATH on Linux, and neither reliably does on
        Windows or macOS, so the standard install location is checked too.
        """
        found = shutil.which(name)
        if found:
            return found
        soffice = name == "soffice"
        if IS_WINDOWS:
            roots = [os.environ.get("PROGRAMFILES"), os.environ.get("PROGRAMFILES(X86)")]
            relative = ("LibreOffice/program/soffice.exe" if soffice
                        else "Calibre2/ebook-convert.exe")
        elif IS_MACOS:
            roots = ["/Applications", str(Path.home() / "Applications")]
            relative = ("LibreOffice.app/Contents/MacOS/soffice" if soffice
                        else "calibre.app/Contents/MacOS/ebook-convert")
        else:
            roots = ["/usr/lib", "/usr/local/lib", "/opt", "/snap/bin",
                     str(Path.home() / ".local/bin")]
            relative = ("libreoffice/program/soffice" if soffice else "calibre/ebook-convert")
        for root in filter(None, roots):
            for candidate in (Path(root) / Path(relative), Path(root) / name):
                if candidate.is_file():
                    return str(candidate)
        return None

    @staticmethod
    def _write_images_pdf(image_sources, target, workspace):
        from PIL import Image, ImageOps

        normalized = []
        for source in image_sources:
            try:
                if isinstance(source, tuple):
                    _, payload = source
                    image = Image.open(io.BytesIO(payload))
                else:
                    image = Image.open(source)
                with image:
                    image.seek(0)
                    oriented = ImageOps.exif_transpose(image)
                    if oriented.width < 4 or oriented.height < 4:
                        continue
                    rgba = oriented.convert("RGBA")
                    canvas = Image.new("RGB", rgba.size, "white")
                    canvas.paste(rgba, mask=rgba.getchannel("A"))
                    # Number only successfully decoded inputs. This keeps the
                    # normalized sequence contiguous when a folder or CBZ has
                    # a corrupt image or tracking pixel in the middle.
                    destination = workspace / f"{len(normalized) + 1:05d}.png"
                    canvas.save(destination, "PNG", optimize=True)
                normalized.append(str(destination))
            except Exception:
                # One damaged image should not destroy an otherwise valid
                # chapter. The caller reports the number actually merged.
                continue
        if not normalized:
            raise RuntimeError("No readable images were found.")

        # A fixed A4 page prevents tiny or unusually tall webtoon strips from
        # creating invalid page sizes. Android uses reportlab because img2pdf
        # depends on pikepdf/qpdf, which has no Android wheel.
        if IS_ANDROID:
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.lib.utils import ImageReader
            from reportlab.pdfgen import canvas as pdf_canvas

            document = pdf_canvas.Canvas(str(target), pagesize=A4)
            for image_path in normalized:
                with Image.open(image_path) as image:
                    page = landscape(A4) if image.width > image.height else A4
                    page_width, page_height = page
                    document.setPageSize(page)
                    scale = min(page_width / image.width, page_height / image.height)
                    width, height = image.width * scale, image.height * scale
                    document.drawImage(
                        ImageReader(image_path), (page_width - width) / 2,
                        (page_height - height) / 2, width, height,
                        preserveAspectRatio=True)
                    document.showPage()
            document.save()
        else:
            import img2pdf
            a4 = (img2pdf.mm_to_pt(210), img2pdf.mm_to_pt(297))
            layout = img2pdf.get_layout_fun(
                a4, fit=img2pdf.FitMode.into, auto_orient=True)
            target.write_bytes(img2pdf.convert(normalized, layout_fun=layout))
        return len(normalized)

    def start_document_conversion(self, path):
        if self.updating or self.setup_running:
            self.notify("Components are being updated. Try again when the update finishes.", "err")
            return
        if self.is_converting:
            self.notify("A conversion is already running.", "err")
            return
        self.is_converting = True
        threading.Thread(target=self.convert_document_file, args=(path,), daemon=True).start()

    @staticmethod
    def _android_document_to_pdf(source, target):
        """Render Android-supported text/office/ebook formats to a readable PDF."""
        from xml.sax.saxutils import escape
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        suffix = source.suffix.lower()
        blocks = []
        if suffix == ".docx":
            from docx import Document
            document = Document(str(source))
            blocks.extend(paragraph.text for paragraph in document.paragraphs)
            for table in document.tables:
                blocks.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            workbook = load_workbook(source, read_only=True, data_only=True)
            for sheet in workbook.worksheets:
                blocks.append(sheet.title)
                blocks.extend(" | ".join("" if value is None else str(value) for value in row)
                              for row in sheet.iter_rows(values_only=True))
            workbook.close()
        elif suffix == ".pptx":
            from pptx import Presentation
            presentation = Presentation(str(source))
            for index, slide in enumerate(presentation.slides, 1):
                blocks.append(f"Slide {index}")
                blocks.extend(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        elif suffix in {".odt", ".ods", ".odp"}:
            from odf import teletype
            from odf.opendocument import load
            document = load(str(source))
            blocks.append(teletype.extractText(document.text))
        elif suffix == ".epub":
            from ebooklib import ITEM_DOCUMENT, epub
            book = epub.read_epub(str(source))
            for item in book.get_items_of_type(ITEM_DOCUMENT):
                text = re.sub(r"<[^>]+>", " ", item.get_content().decode("utf-8", "replace"))
                blocks.append(re.sub(r"\s+", " ", text))
        elif suffix == ".fb2":
            import xml.etree.ElementTree as element_tree
            blocks.append(" ".join(element_tree.parse(source).getroot().itertext()))
        elif suffix in {".txt", ".csv", ".html", ".htm", ".rtf"}:
            text = source.read_text("utf-8", "replace")
            if suffix in {".html", ".htm"}:
                text = re.sub(r"<[^>]+>", " ", text)
            elif suffix == ".rtf":
                text = re.sub(r"\\[A-Za-z]+-?\d* ?|[{}]", "", text)
            blocks.append(text)
        else:
            raise RuntimeError(
                "This legacy document type needs LibreOffice or Calibre on a desktop build.")

        styles = getSampleStyleSheet()
        body = styles["BodyText"]
        body.wordWrap = "CJK"
        story = []
        for block in blocks:
            for paragraph in str(block).replace("\r", "").split("\n"):
                if paragraph.strip():
                    story.extend((Paragraph(escape(paragraph.strip()), body), Spacer(1, 2 * mm)))
        if not story:
            raise RuntimeError("The selected document contains no readable text.")
        SimpleDocTemplate(str(target), pagesize=A4, leftMargin=15 * mm,
                          rightMargin=15 * mm, topMargin=15 * mm,
                          bottomMargin=15 * mm).build(story)

    def start_image_folder_conversion(self, path):
        if self.updating or self.setup_running:
            self.notify("Components are being updated. Try again when the update finishes.", "err")
            return
        if self.is_converting:
            self.notify("A conversion is already running.", "err")
            return
        self.is_converting = True
        threading.Thread(target=self.convert_image_folder, args=(path,), daemon=True).start()

    def convert_image_folder(self, path):
        """Combine a chapter/image folder into a naturally ordered PDF."""
        folder = Path(path)
        partial = None
        self.ui("converting", True, "images")
        try:
            if not folder.is_dir():
                raise RuntimeError("The selected image folder no longer exists.")
            images = [item for item in folder.rglob("*")
                      if item.is_file() and item.suffix.lower() in IMAGE_DOCUMENT_EXTENSIONS]
            images.sort(key=lambda item: natural_sort_key(item.relative_to(folder).as_posix()))
            if not images:
                raise RuntimeError("No supported images were found in that folder.")
            if len(images) > MANGA_PAGE_LIMIT:
                raise RuntimeError(f"Image folders are limited to {MANGA_PAGE_LIMIT} pages.")
            target_parent = (download_dir("converted", self.save_folder.get().strip())
                             if IS_ANDROID else folder.parent)
            target = unique_path(target_parent, safe_filename(
                folder.name, fallback="chapter", suffix=".pdf"))
            partial = target.with_name(target.stem + ".working.pdf")
            self.log(f"Combining {len(images)} ordered images into {target.name}...", "info")
            with tempfile.TemporaryDirectory() as temp:
                merged = self._write_images_pdf(images, partial, Path(temp))
            if not is_valid_pdf_file(partial):
                raise RuntimeError("The image converter returned an incomplete PDF.")
            os.replace(partial, target)
            skipped = len(images) - merged
            detail = (f" ({skipped} unreadable/tiny image(s) skipped)" if skipped else "")
            self.log(f"Created {target.name} from {merged} images{detail}.", "success")
            self.notify(f"Created {target.name} from {merged} ordered images{detail}", "ok")
        except Exception as exc:
            detail = str(exc)[:250]
            self.log(f"Image-folder conversion failed: {detail}", "error")
            self.notify(detail, "err")
        finally:
            if partial:
                partial.unlink(missing_ok=True)
            self.is_converting = False
            self.ui("converting", False, "images")

    def convert_document_file(self, path):
        source = Path(path)
        partial = None
        self.ui("converting", True, "document")
        try:
            if not source.is_file():
                raise RuntimeError("The selected document no longer exists.")
            target_parent = (download_dir("converted", self.save_folder.get().strip())
                             if IS_ANDROID else source.parent)
            target = unique_path(target_parent, safe_filename(
                source.stem, fallback="document", suffix=".pdf"))
            partial = target.with_name(target.stem + ".working.pdf")
            partial.unlink(missing_ok=True)
            suffix = source.suffix.lower()
            self.log(f"Converting {source.name} to PDF...", "info")

            with tempfile.TemporaryDirectory() as temp:
                workspace = Path(temp)
                if suffix in IMAGE_DOCUMENT_EXTENSIONS:
                    self._write_images_pdf([source], partial, workspace)
                elif suffix == ".cbz":
                    with zipfile.ZipFile(source) as archive:
                        members = [member for member in archive.infolist()
                                   if not member.is_dir()
                                   and Path(member.filename).suffix.lower()
                                   in IMAGE_DOCUMENT_EXTENSIONS]
                        members.sort(key=lambda member: natural_sort_key(member.filename))
                        if len(members) > MANGA_PAGE_LIMIT:
                            raise RuntimeError(f"Comic archives are limited to {MANGA_PAGE_LIMIT} pages.")
                        images = []
                        for index, member in enumerate(members, 1):
                            # Extract to generated flat names instead of trusting
                            # archive paths, and stream each member so large CBZ
                            # files do not keep every page in RAM simultaneously.
                            extracted = workspace / (
                                f"cbz-{index:05d}{Path(member.filename).suffix.lower()}")
                            with (archive.open(member) as source_image,
                                  extracted.open("wb") as target_image):
                                shutil.copyfileobj(source_image, target_image, length=256 * 1024)
                            images.append(extracted)
                    self._write_images_pdf(images, partial, workspace)
                elif suffix == ".pdf":
                    pdf_rewrite(source, partial)
                elif IS_ANDROID and suffix in (EBOOK_CONVERT_EXTENSIONS |
                                               OFFICE_DOCUMENT_EXTENSIONS):
                    self._android_document_to_pdf(source, partial)
                elif suffix in EBOOK_CONVERT_EXTENSIONS:
                    converter = self._converter_program("ebook-convert")
                    if not converter:
                        raise RuntimeError(
                            "Calibre is required for EPUB, MOBI, AZW, AZW3 and FB2 conversion.")
                    result = subprocess.run(
                        [converter, str(source), str(partial)], capture_output=True, text=True,
                        encoding="utf-8", errors="replace",
                        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
                    if result.returncode != 0:
                        raise RuntimeError((result.stderr or result.stdout).strip()[-220:])
                elif suffix in OFFICE_DOCUMENT_EXTENSIONS:
                    converter = self._converter_program("soffice")
                    if not converter:
                        raise RuntimeError(
                            "LibreOffice is required for Office and OpenDocument conversion.")
                    result = subprocess.run(
                        [converter, "--headless", "--convert-to", "pdf", "--outdir",
                         str(workspace), str(source)], capture_output=True, text=True,
                        encoding="utf-8", errors="replace",
                        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
                    generated = workspace / f"{source.stem}.pdf"
                    if result.returncode != 0 or not generated.is_file():
                        raise RuntimeError((result.stderr or result.stdout or
                                            "LibreOffice returned no PDF").strip()[-220:])
                    shutil.copyfile(generated, partial)
                else:
                    raise RuntimeError("That document type cannot be converted to PDF.")

            if not is_valid_pdf_file(partial):
                raise RuntimeError("The converter returned an incomplete or invalid PDF.")
            os.replace(partial, target)
            self.log(f"Converted: {target.name}", "success")
            self.notify(f"Created {target.name}", "ok")
        except Exception as exc:
            detail = str(exc)[:250]
            self.log(f"Document conversion failed: {detail}", "error")
            self.notify(detail, "err")
        finally:
            if partial:
                partial.unlink(missing_ok=True)
            self.is_converting = False
            self.ui("converting", False, "document")

    def on_close(self):
        self.closing = True
        if self.process:
            try:
                self.process.terminate()
            except Exception:
                pass
        try:
            self.executor.shutdown(wait=False, cancel_futures=True)
        except Exception:
            pass


class Api:
    """The js_api surface: only these methods reach the page.

    pywebview walks every public attribute of js_api into the JS bridge, so the
    engine, the window, and anything unserialisable stay underscored.
    """

    def __init__(self, app):
        self._app = app
        self._window = None

    # -- lifecycle ------------------------------------------------------
    def hello(self):
        self._app.begin()
        return self._app.snapshot()

    def poll(self):
        return self._app.drain()

    def set_option(self, name, value):
        self._app.set_option(name, value)

    # -- window chrome --------------------------------------------------
    def minimize(self):
        self._window.minimize()

    def toggle_maximize(self):
        self._window.toggle_fullscreen()

    def close(self):
        self._app.on_close()
        self._window.destroy()

    # -- actions --------------------------------------------------------
    def start_download(self, url, fmt, quality):
        self._app.start_download(url, fmt, quality)

    def start_media(self, kind, url, quality, audio_only, limit):
        self._app.start_media_download(kind, url, quality, audio_only, limit or None)

    def start_pdf(self, url, limit):
        self._app.start_pdf_download(url, limit or None)

    def start_manga(self, url, limit):
        self._app.start_manga_download(url, limit or None)

    def stop_download(self):
        self._app.stop_download()

    def set_preset(self, proxy_type, url):
        self._app.set_proxy_preset(proxy_type, url)
        self._app.ui("proxy", proxy_type, url, True)

    def test_proxy(self):
        self._app.test_proxy()

    def check_setup(self):
        threading.Thread(target=self._app.check_setup, daemon=True).start()

    def sign_in(self, kind):
        self._app.sign_in(kind)

    def finish_sign_in(self):
        threading.Thread(target=self._app.finish_sign_in, daemon=True).start()

    def cookie_status(self):
        return self._app.cookie_status()

    def update_tools(self):
        self._app.update_tools()

    def open_folder(self):
        self._app.open_folder()

    def open_telegram(self):
        webbrowser.open(TELEGRAM_URL)

    def clipboard(self):
        return read_clipboard()

    def copy_text(self, text):
        """Fallback for pages whose browser refuses navigator.clipboard."""
        return write_clipboard(str(text or ""))

    def about(self):
        self._app.notify(
            f"{APP_NAME} {APP_VERSION}\n"
            f"Created by {CREATOR}\n"
            f"{TELEGRAM_URL}\n\n"
            "Spotify, YouTube, YouTube Music, TikTok, Instagram, SoundCloud, X,\n"
            "documents, ebooks, manga/manhwa and general video sites.\n"
            "FFmpeg, spotDL, yt-dlp, Deno, Streamlink and gallery-dl ship inside the app.\n"
            "The core command-line tools update from\n"
            "mirrored sources. V2RayN compatible.",
            "info",
        )

    def browse(self):
        picked = self._window.create_file_dialog(webview.FOLDER_DIALOG)
        if picked:
            folder = picked[0]
            self._app.set_option("folder", folder)
            self._app.ui("folder", folder)

    def convert_media(self, options):
        if not self._app.ffmpeg_cmd:
            self._app.notify("FFmpeg is not ready yet.", "err")
            return
        picked = self._window.create_file_dialog(
            webview.OPEN_DIALOG,
            file_types=(
                "Media files (*.mp4;*.mkv;*.webm;*.mov;*.avi;*.m4v;*.mp3;*.flac;*.wav;*.m4a;*.aac;*.ogg)",
                "All files (*.*)",
            ),
        )
        if picked:
            self._app.start_media_conversion(picked[0], options)

    def convert_document(self):
        picked = self._window.create_file_dialog(
            webview.OPEN_DIALOG,
            file_types=(
                "Documents (*.pdf;*.doc;*.docx;*.xls;*.xlsx;*.ppt;*.pptx;*.odt;*.ods;*.odp;*.rtf;*.txt;*.csv;*.html;*.epub;*.mobi;*.azw;*.azw3;*.fb2;*.cbz;*.jpg;*.jpeg;*.png;*.webp;*.avif;*.gif;*.tif;*.tiff;*.bmp)",
                "All files (*.*)",
            ),
        )
        if picked:
            self._app.start_document_conversion(picked[0])

    def convert_image_folder(self):
        picked = self._window.create_file_dialog(webview.FOLDER_DIALOG)
        if picked:
            self._app.start_image_folder_conversion(picked[0])


def main():
    app = SpotifyDownloader()
    api = Api(app)
    window = webview.create_window(
        f"{APP_NAME} {APP_VERSION}",
        url=web_index(),
        js_api=api,
        width=1210,
        height=820,
        min_size=(1005, 697),
        frameless=True,
        easy_drag=False,
        background_color="#05080d",
        text_select=False,
    )
    api._window = window
    # private_mode is pywebview's default, and it throws the browser profile away
    # on exit — which is why a sign-in never survived a restart. Give the webview
    # a folder of its own and the account stays signed in, like any browser.
    webview.start(private_mode=False,
                  storage_path=str(Path(os.environ.get("LOCALAPPDATA", Path.home()))
                                   / "SpotifyDownloader" / "browser"))
    app.on_close()
    return 0


if __name__ == "__main__":
    if "--validate-engine-update" in sys.argv:
        try:
            _probe_index = sys.argv.index("--validate-engine-update")
            _probe_paths = json.loads(sys.argv[_probe_index + 1])
            if not isinstance(_probe_paths, list) or not all(
                    isinstance(item, str) for item in _probe_paths):
                raise ValueError("engine paths must be a JSON list of strings")
        except (ValueError, IndexError, json.JSONDecodeError):
            raise SystemExit(4)
        raise SystemExit(0 if validate_engine_update(_probe_paths) else 4)
    elif "--self-test" in sys.argv:
        spotdl = bundled_tool(f"spotdl{EXE}")
        ffmpeg = bundled_tool(f"ffmpeg{EXE}")
        ytdlp = bundled_tool(f"yt-dlp{EXE}")
        deno = bundled_tool(f"deno{EXE}")
        if not ffmpeg:
            local_ffmpeg = app_data_dir() / "ffmpeg" / "bin" / f"ffmpeg{EXE}"
            ffmpeg = str(local_ffmpeg) if local_ffmpeg.is_file() else None
        # Only the tools this platform actually has a build of are required.
        if not ffmpeg or not ytdlp or ("spotdl" in TOOLS and not spotdl) or (
                "deno" in TOOLS and not deno):
            raise SystemExit(2)
        # --ffmpeg-location must name something yt-dlp can actually find a
        # program in: the folder only when it holds the real names, the binary
        # itself otherwise. Getting this wrong costs every merge and every MP3.
        _loc_probe = SpotifyDownloader.__new__(SpotifyDownloader)
        with tempfile.TemporaryDirectory() as _tmp:
            _shimmed = Path(_tmp) / "shimmed"
            _shimmed.mkdir()
            (_shimmed / f"ffmpeg{EXE}").write_bytes(b"")
            _loc_probe.ffmpeg_cmd = str(_shimmed / f"ffmpeg{EXE}")
            assert SpotifyDownloader.ffmpeg_location(_loc_probe) == str(_shimmed)

            _packaged = Path(_tmp) / "lib"
            _packaged.mkdir()
            (_packaged / "libffmpeg.so").write_bytes(b"")
            _loc_probe.ffmpeg_cmd = str(_packaged / "libffmpeg.so")
            # No file called ffmpeg here, so the folder would find nothing.
            assert SpotifyDownloader.ffmpeg_location(_loc_probe) == \
                str(_packaged / "libffmpeg.so")

        # A run of log lines must reach the page as one event, and anything
        # else in the middle must break the run rather than swallow it.
        _drain_probe = SpotifyDownloader.__new__(SpotifyDownloader)
        _drain_probe.ui_queue = queue.Queue()
        for _n in range(5):
            _drain_probe.ui_queue.put(("batch_log", [(f"line {_n}", "info")]))
        _drain_probe.ui_queue.put(("progress", (50,)))
        _drain_probe.ui_queue.put(("batch_log", [("after", "info")]))
        _drained = SpotifyDownloader.drain(_drain_probe)
        assert [event[0] for event in _drained] == ["log", "progress", "log"], _drained
        assert len(_drained[0][1]) == 5, _drained[0]
        assert _drained[0][1][0] == ["line 0", "info"], _drained[0]

        assert MEDIA_PATTERNS["youtube"].search("https://youtu.be/abc")
        assert MEDIA_PATTERNS["tiktok"].search("https://www.tiktok.com/@u/video/1")
        assert normalize_media_url("soundcloud", "https://soundcloud.com/a/song")
        assert normalize_media_url("x", "https://x.com/user/status/1")
        assert normalize_media_url("general", "https://example.com/watch/1")
        for blocked in ("https://youtube.com/watch?v=1", "https://m.youtube.com/shorts/1",
                        "https://youtu.be/1", "https://www.youtube-nocookie.com/embed/1"):
            try:
                normalize_media_url("general", blocked)
            except ValueError:
                pass
            else:
                raise AssertionError(f"General downloader accepted YouTube URL: {blocked}")
        assert normalize_media_url("ytmusic", "https://music.youtube.com/watch?v=1")
        # YouTube Music is a narrower door than YouTube, not the same one.
        for wrong in ("https://www.youtube.com/watch?v=1", "https://youtu.be/1"):
            try:
                normalize_media_url("ytmusic", wrong)
            except ValueError:
                pass
            else:
                raise AssertionError(f"YouTube Music accepted a plain YouTube URL: {wrong}")
        assert cookie_site("ytmusic") == "youtube" and cookie_site("x") == "x"
        assert SIGNIN_PAGES["x"][0] == "https://x.com/i/flow/login"
        assert SIGNIN_PAGES["soundcloud"][0] == "https://soundcloud.com/signin"
        assert "auth_token" in SITE_SESSION_COOKIES["x"]
        assert "oauth_token" in SITE_SESSION_COOKIES["soundcloud"]
        assert "soundcloud" in COOKIE_ON_DEMAND
        assert "certifi" in PYTHON_COMPONENTS
        assert f"v9.9.9/{SPOTDL_ASSETS[OS_TAG].format(v='9.9.9')}" in (
            SpotifyDownloader.tool_download_urls("spotdl", "9.9.9")[0])
        _app = SpotifyDownloader()
        _app.ytdlp_cmd, _app.ffmpeg_cmd, _app.deno_cmd = ytdlp, ffmpeg, deno
        _ym_cmd = _app._ytdlp_cmd_for(
            "ytmusic", "https://music.youtube.com/watch?v=1", Path(tempfile.gettempdir()),
            "best", True, None)
        _sc_cmd = _app._ytdlp_cmd_for(
            "soundcloud", "https://soundcloud.com/a/song", Path(tempfile.gettempdir()),
            "best", True, None)
        _sc_cookie_cmd = _app._ytdlp_cmd_for(
            "soundcloud", "https://soundcloud.com/a/song", Path(tempfile.gettempdir()),
            "best", True, None, ("file", str(Path(tempfile.gettempdir()) / "sc-cookies.txt")))
        _general_cmd = _app._ytdlp_cmd_for(
            "general", "https://video.example.com/watch/1", Path(tempfile.gettempdir()),
            "best", False, None, impersonate=True)
        _fresh_cmd = _app._ytdlp_cmd_for(
            "youtube", "https://youtu.be/g9ilhvSurpQ", Path(tempfile.gettempdir()),
            "best", False, None, clients="android_vr", fresh=True)
        assert "--embed-metadata" in _ym_cmd and "--embed-thumbnail" in _ym_cmd
        assert "--js-runtimes" in _ym_cmd and any(str(deno) in part for part in _ym_cmd)
        assert "--js-runtimes" in _sc_cmd and "--js-runtimes" in _general_cmd
        assert "--sleep-requests" in _ym_cmd
        assert "--ignore-config" in _fresh_cmd
        assert "--no-continue" in _fresh_cmd
        assert "--fragment-retries" in _fresh_cmd and "--extractor-retries" in _fresh_cmd
        assert any("player_client=android_vr" in part for part in _fresh_cmd)
        _spent = set()
        _step = _app._next_signin_step(
            "youtube", "https://youtu.be/g9ilhvSurpQ",
            {"cookies": None, "clients": None, "impersonate": False, "fresh": False},
            _spent, challenge=True)
        assert _step["clients"] == "android_vr" and _step["fresh"]
        _step = _app._next_signin_step(
            "youtube", "https://youtu.be/g9ilhvSurpQ", _step, _spent, challenge=True)
        assert _step["clients"] == YT_CLIENT_LADDER[1] and _step["fresh"]
        _app.youtube_po_token.set("mweb.gvs+" + "A" * 32)
        _pot_cmd = _app._ytdlp_cmd_for(
            "youtube", "https://youtube.com/watch?v=1", Path(tempfile.gettempdir()),
            "best", False, None)
        assert any("po_token=mweb.gvs+" in part and "player_client=mweb,default" in part
                   for part in _pot_cmd)
        assert "--extract-audio" in _sc_cmd and "--embed-metadata" not in _sc_cmd
        assert "--cookies" in _sc_cookie_cmd
        # A replayed session goes out as the browser that earned it, but only
        # where that is known — which is Android, never the desktop, where the
        # jar may have come out of Firefox or Edge.
        assert ("--user-agent" in _sc_cookie_cmd) == IS_ANDROID, _sc_cookie_cmd
        if IS_ANDROID:
            assert _sc_cookie_cmd[_sc_cookie_cmd.index("--user-agent") + 1] == BROWSER_UA
        assert "--user-agent" not in _sc_cmd            # no jar, no claim to make
        _yt_cookie_cmd = _app._ytdlp_cmd_for(
            "youtube", "https://youtu.be/g9ilhvSurpQ", Path(tempfile.gettempdir()),
            "best", False, None, ("file", str(Path(tempfile.gettempdir()) / "yt-cookies.txt")))
        assert "--cookies" in _yt_cookie_cmd and "--user-agent" not in _yt_cookie_cmd
        # The identity is one string everywhere, and it never claims the wrong
        # kind of device.
        assert ("Android" in BROWSER_UA) == IS_ANDROID, BROWSER_UA
        # Opening on android_vr must take that rung out of the ladder, or the
        # first escalation would retry what already failed.
        # The rungs that need neither a token nor a solved player come first,
        # and the one that can carry cookies comes last. Checked against
        # yt-dlp's own client table when it is importable — the frozen desktop
        # builds ship yt-dlp as an executable, so there is nothing to import
        # there and the ladder is taken on trust rather than failing the run.
        try:
            from yt_dlp.extractor.youtube._base import INNERTUBE_CLIENTS as _YT_CLIENTS
        except Exception:
            _YT_CLIENTS = None
        if _YT_CLIENTS:
            for _rung in YT_CLIENT_LADDER:
                _needs_nothing = all(
                    not _YT_CLIENTS[c].get("REQUIRE_PO_TOKEN")
                    and not _YT_CLIENTS[c].get("REQUIRE_JS_PLAYER", True)
                    for c in _rung.split(",") if c in _YT_CLIENTS)
                _carries_session = any(
                    _YT_CLIENTS[c].get("SUPPORTS_COOKIES")
                    for c in _rung.split(",") if c in _YT_CLIENTS)
                # Every rung must earn its place: it either asks the site for
                # nothing, or it can spend a session. A rung that does neither
                # is one that cannot succeed where the one above it failed.
                assert _needs_nothing or _carries_session, _rung
                assert _carries_session == (_rung in YT_COOKIE_CLIENTS), _rung
            # The opening rung and the first escalation are what they have
            # always been; the additions sit below them.
            assert YT_CLIENT_LADDER[:2] == ("android_vr", "web_embedded,tv,default")

        # A phone can only ever read a jar this app wrote; a desktop must keep
        # every browser it could read one out of.
        assert any(c[0] == "browser" for c in cookie_candidates()) != IS_ANDROID

        _lead_spent = {"clients-0"}
        _lead_step = _app._next_signin_step(
            "youtube", "https://youtu.be/g9ilhvSurpQ",
            {"cookies": None, "clients": YT_CLIENT_LADDER[0], "impersonate": False,
             "fresh": False},
            _lead_spent, challenge=True)
        assert _lead_step["clients"] == YT_CLIENT_LADDER[1], _lead_step
        assert _general_cmd[_general_cmd.index("--impersonate") + 1] == "chrome"
        assert "--sleep-requests" not in _general_cmd
        assert not any("youtube:" in part or "po_token=" in part or "player_client=" in part
                       for part in _general_cmd)
        assert cookie_domain_for("general", "https://video.Example.com/watch") == "video.example.com"
        assert cookie_domain_for("youtube", "https://music.youtube.com/watch") == "youtube.com"
        for phrase in ("HTTP Error 403: Forbidden", "Cloudflare CAPTCHA",
                       "Please verify that you are human"):
            assert GENERAL_CHALLENGE.search(phrase), phrase
        assert SIGNIN_DEMANDED.search("Sign in to confirm you’re not a bot")
        assert set(MEDIA_LABELS) <= set(SOURCES), "every media source needs a folder"
        assert "pdf" in SOURCES

        # A server picks these names, so they are untrusted input.
        assert safe_filename("../../../../boot.ini") == "boot.pdf"
        assert safe_filename("report.pdf") == "report.pdf"
        assert safe_filename("a/b/CON.pdf") == "_CON.pdf"
        assert safe_filename('bad:name*?.pdf') == "bad_name__.pdf"
        assert safe_filename("") == "document.pdf"
        assert "/" not in safe_filename("x/y") and "\\" not in safe_filename("x\\y")
        for bad in ("ftp://host/a.pdf", "file:///c:/a.pdf", "https://u:p@host/a.pdf", "notaurl"):
            try:
                normalize_page_url(bad)
            except ValueError:
                pass
            else:
                raise AssertionError(f"PDF downloader accepted {bad}")
        assert normalize_page_url("https://example.com/docs")
        _p = PdfLinkParser()
        _p.feed('<a href="/a.pdf">x</a><iframe src="b.pdf"></iframe><a href="#top">y</a>')
        assert _p.links == ["/a.pdf", "b.pdf", "#top"], _p.links
        assert detect_document_extension(b"%PDF-1.7", "https://host/download") == ".pdf"
        assert detect_document_extension(b"PK\x03\x04data", "https://host/book.epub") == ".zip"
        _mobi = bytearray(80)
        _mobi[60:68] = b"BOOKMOBI"
        assert detect_document_extension(bytes(_mobi), "https://host/book.azw3") == ".azw3"
        assert detect_document_extension(b"<FictionBook/>", "https://host/book.fb2") == ".fb2"
        assert detect_document_extension(b"<html>error</html>", "https://host/book.epub") is None
        with tempfile.TemporaryDirectory() as _tmp:
            _complete_pdf = Path(_tmp) / "complete.pdf"
            _complete_pdf.write_bytes(b"%PDF-1.7\nbody\n%%EOF\n")
            _truncated_pdf = Path(_tmp) / "truncated.pdf"
            _truncated_pdf.write_bytes(b"%PDF-1.7\nbody")
            assert is_valid_pdf_file(_complete_pdf)
            assert not is_valid_pdf_file(_truncated_pdf)
            _book = Path(_tmp) / "book.zip"
            with zipfile.ZipFile(_book, "w") as _archive:
                _archive.writestr("mimetype", "application/epub+zip")
                _archive.writestr("META-INF/container.xml", "<container/>")
            assert inspect_zip_document(_book) == ".epub"
            _broken_book = Path(_tmp) / "missing-container.zip"
            with zipfile.ZipFile(_broken_book, "w") as _archive:
                _archive.writestr("mimetype", "application/epub+zip")
            assert inspect_zip_document(_broken_book) is None
            _office = Path(_tmp) / "office.zip"
            with zipfile.ZipFile(_office, "w") as _archive:
                _archive.writestr("[Content_Types].xml", "<Types/>")
                _archive.writestr("word/document.xml", "<document/>")
            assert inspect_zip_document(_office) == ".docx"
        _media = HtmlMediaParser()
        _media.feed('<video src="/movie.mp4"></video><meta property="og:video" '
                    'content="https://cdn.test/master.m3u8"><a href="sound.mp3">x</a>')
        assert _media.links == ["/movie.mp4", "https://cdn.test/master.m3u8", "sound.mp3"]
        _images = HtmlImageParser()
        _images.feed('<img data-src="/001.jpg"><img src="pixel.gif" '
                     'srcset="/002-small.webp 320w, /002.webp 1280w">')
        assert _images.links == ["/001.jpg", "/002.webp"]
        _bulk_images = HtmlImageParser()
        _bulk_images.feed("".join(
            f'<img data-lazy-src="/chapter/page-{index:03d}.jpg">'
            for index in range(1, 101)))
        assert len(_bulk_images.links) == 100
        assert _bulk_images.links[0].endswith("001.jpg")
        assert _bulk_images.links[-1].endswith("100.jpg")
        _streams = {"360p": object(), "720p": object(), "1080p60": object()}
        assert SpotifyDownloader._pick_streamlink_stream(_streams, "720")[0] == "720p"
        assert sorted(["page10.jpg", "page2.jpg", "page1.jpg"], key=natural_sort_key) == [
            "page1.jpg", "page2.jpg", "page10.jpg"]
        assert versions_equal("2026.07.22", "2026.7.22")
        assert not versions_equal("8.5.0", "8.6.0")
        _convert_options = SpotifyDownloader.normalize_media_conversion_options({
            "format": "mp4", "codec": "h265", "quality": "master",
            "resize": "1080p", "fps": "30", "frame_mode": "motion",
            "upscale": "lanczos", "video_enhance": "detail",
            "audio_enhance": "normalize", "audio_bitrate": "256",
        })
        _convert_cmd = SpotifyDownloader.media_conversion_command(
            ffmpeg, Path("input.mkv"), Path("output.mp4"), _convert_options)
        assert "libx265" in _convert_cmd and "minterpolate=fps=30" in " ".join(_convert_cmd)
        assert "scale=1920:1080" in " ".join(_convert_cmd) and "+faststart" in _convert_cmd
        assert "flags=lanczos" in " ".join(_convert_cmd)
        assert "unsharp=" in " ".join(_convert_cmd) and "loudnorm=" in " ".join(_convert_cmd)
        assert _convert_cmd[_convert_cmd.index("-preset") + 1] == "slow"
        _audio_options = SpotifyDownloader.normalize_media_conversion_options({
            "format": "ogg", "audio_bitrate": "320", "audio_enhance": "music"})
        _audio_cmd = SpotifyDownloader.media_conversion_command(
            ffmpeg, Path("input.mp4"), Path("output.ogg"), _audio_options)
        assert "-vn" in _audio_cmd and "libopus" in _audio_cmd and "320k" in _audio_cmd
        assert "acompressor=" in " ".join(_audio_cmd) and "48000" in _audio_cmd
        with tempfile.TemporaryDirectory() as _tmp:
            _chapter_dir, _chapter_pdf, _chapter_cbz = unique_chapter_paths(
                Path(_tmp), "chapter 12")
            assert _chapter_dir.name == "chapter 12"
            assert _chapter_pdf.name == "chapter 12.pdf" and _chapter_cbz.name == "chapter 12.cbz"
            _chapter_dir.mkdir()
            assert unique_chapter_paths(Path(_tmp), "chapter 12")[0].name == "chapter 12 (2)"

        # yt-dlp saying "no video here" must route to another engine; yt-dlp
        # failing for any other reason must not.
        for phrase in ("ERROR: [twitter] 1: No video could be found in this tweet",
                       "ERROR: Unsupported URL: https://example.com/gallery",
                       "ERROR: no video formats found"):
            assert NO_VIDEO_FOR_YTDLP.search(phrase), phrase
        for phrase in ("ERROR: unable to download video data: HTTP Error 403",
                       "ERROR: [DRM] The requested site is known to use DRM protection",
                       "ERROR: Video unavailable"):
            assert not NO_VIDEO_FOR_YTDLP.search(phrase), phrase
        assert YOUTUBE_MEDIA_DENIED.search(
            "ERROR: unable to download video data: HTTP Error 403: Forbidden")
        # Only sources with a second engine may hand off at all.
        assert set(ENGINES) <= set(MEDIA_LABELS), set(ENGINES) - set(MEDIA_LABELS)
        assert not ({"youtube", "ytmusic"} & set(ENGINES)), (
            "YouTube and YouTube Music have their own client ladder, not general fallbacks")
        assert "gallery-dl" in ENGINES["x"] and "direct" in ENGINES["general"]
        assert ENGINES["general"][:2] == ("streamlink", "html5")
        # Dynamic plugin imports are easy for a freezer to miss. Exercise both
        # engines inside the packaged executable without touching the network.
        import gallery_dl
        from gallery_dl.job import DataJob
        from streamlink import Streamlink
        import img2pdf
        import pikepdf
        from PIL import Image
        with tempfile.TemporaryDirectory() as _tmp:
            _page = Path(_tmp) / "page.png"
            Image.new("RGB", (8, 8), "white").save(_page)
            _pdf_bytes = img2pdf.convert([str(_page)])
            assert _pdf_bytes.startswith(PDF_MAGIC)
            with pikepdf.Pdf.open(io.BytesIO(_pdf_bytes)) as _pdf:
                assert len(_pdf.pages) == 1
        _gallery_job = DataJob("https://x.com/example/status/1", file=io.StringIO())
        assert _gallery_job.extractor.category == "twitter"
        _stream_name, _, _ = Streamlink().resolve_url("hls://example.com/live.m3u8")
        assert _stream_name == "hls"
        # Media must keep its own extension; only the PDF page forces .pdf.
        assert safe_filename("clip", suffix=".mp4") == "clip.mp4"
        assert safe_filename("../x/photo.jpg", suffix=".jpg") == "photo.jpg"

        assert not MEDIA_PATTERNS["tiktok"].search("https://open.spotify.com/track/1")
        assert not MEDIA_PATTERNS["youtube"].search("https://nyoutube.com.evil.tld/x")
        # Every tool must be reachable through the same generic lookup.
        for _name in TOOLS:
            assert bundled_tool(TOOLS[_name]["exe"]), _name
            assert len(TOOLS[_name]["urls"]) >= 2, f"{_name} needs a fallback mirror"
        checks = ([spotdl, "--version"], [ffmpeg, "-version"], [ytdlp, "--version"],
                  [deno, "--version"])
        for command in checks:
            result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=30)
            if result.returncode != 0:
                raise SystemExit(3)
        raise SystemExit(0)
    else:
        main()
